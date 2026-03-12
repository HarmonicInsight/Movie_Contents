"""共通ヘルパー関数 - 通信制高校選び完全ガイド YouTube動画用"""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 定数 ──
SLIDE_WIDTH = 9144000
SLIDE_HEIGHT = 5143500
MARGIN = Emu(457200)
CONTENT_W = Emu(8229600)

# カラーパレット
C_DARK = RGBColor(0x33, 0x33, 0x33)
C_GOLD = RGBColor(0x8B, 0x75, 0x36)
C_GOLD_LIGHT = RGBColor(0xC9, 0xA8, 0x4C)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_LIGHT = RGBColor(0xF5, 0xF0, 0xE6)
C_BROWN = RGBColor(0x3C, 0x2F, 0x1E)
C_ACCENT_BLUE = RGBColor(0x2B, 0x5C, 0x8A)
C_ACCENT_GREEN = RGBColor(0x3A, 0x7D, 0x5E)
C_ACCENT_ORANGE = RGBColor(0xC4, 0x6B, 0x2B)
C_ACCENT_PURPLE = RGBColor(0x6B, 0x4C, 0x9A)
C_RED = RGBColor(0x8B, 0x2E, 0x2E)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def new_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height, text, font_name="Calibri",
                font_size=Pt(14), bold=False, color=C_DARK, alignment=PP_ALIGN.LEFT,
                word_wrap=True):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    return txbox


def add_multi_text(slide, left, top, width, height, lines):
    """lines: list of (text, font_size, bold, color)"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, font_size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = "Calibri"
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return txbox


def add_rounded_rect(slide, left, top, width, height, fill_color, border=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not border:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_footer(slide, page_num, series_label=""):
    add_rect(slide, 0, Emu(4800600), Emu(SLIDE_WIDTH), Emu(342900), C_BROWN)
    label = "H A R M O N I C   i n s i g h t"
    if series_label:
        label += f"  |  {series_label}"
    add_textbox(slide, MARGIN, Emu(4800600), Emu(6400800), Emu(342900),
                label, "Calibri", Pt(8), False, C_WHITE, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(8229600), Emu(4800600), Emu(731520), Emu(342900),
                str(page_num), "Calibri", Pt(8), False, C_WHITE, PP_ALIGN.CENTER)


def add_title_bar(slide, title_text, subtitle_text=None):
    add_textbox(slide, MARGIN, Emu(182880), CONTENT_W, Emu(502920),
                title_text, "Calibri", Pt(22), True, C_DARK, PP_ALIGN.LEFT)
    add_rect(slide, MARGIN, Emu(685800), Emu(1097280), Emu(36576), C_GOLD_LIGHT)
    if subtitle_text:
        add_textbox(slide, MARGIN, Emu(868680), CONTENT_W, Emu(320040),
                    subtitle_text, "Calibri", Pt(13), True, C_GOLD, PP_ALIGN.LEFT)


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_title_slide(prs, main_title, sub_title, series_text, date_text):
    """共通タイトルスライド"""
    slide = new_slide(prs)
    add_rect(slide, 0, 0, Emu(SLIDE_WIDTH), Emu(54864), C_GOLD_LIGHT)
    add_rect(slide, 0, Emu(5088636), Emu(SLIDE_WIDTH), Emu(54864), C_GOLD_LIGHT)

    # Hi icon
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(4114800), Emu(914400),
                                   Emu(914400), Emu(822960))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_BG_LIGHT
    shape.line.fill.background()
    add_textbox(slide, Emu(4114800), Emu(960120), Emu(914400), Emu(822960),
                "Hi", "Georgia", Pt(22), False, C_GOLD_LIGHT, PP_ALIGN.CENTER)

    # decorative lines
    add_rect(slide, Emu(914400), Emu(2057400), Emu(2926080), Emu(9144), C_GOLD_LIGHT)
    add_rect(slide, Emu(5303520), Emu(2057400), Emu(2926080), Emu(9144), C_GOLD_LIGHT)

    # Main title
    add_textbox(slide, MARGIN, Emu(2194560), CONTENT_W, Emu(502920),
                main_title, "Georgia", Pt(32), False, C_BROWN, PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN, Emu(2697480), CONTENT_W, Emu(457200),
                sub_title, "Georgia", Pt(24), False, C_BROWN, PP_ALIGN.CENTER)

    add_rect(slide, Emu(2286000), Emu(3200400), Emu(4572000), Emu(9144), C_GOLD_LIGHT)

    add_textbox(slide, MARGIN, Emu(3383280), CONTENT_W, Emu(365760),
                series_text, "Calibri", Pt(14), False, C_GOLD, PP_ALIGN.CENTER)

    # branding
    add_textbox(slide, MARGIN, Emu(3886200), CONTENT_W, Emu(274320),
                "H A R M O N I C", "Calibri", Pt(14), False, C_BROWN, PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN, Emu(4114800), CONTENT_W, Emu(274320),
                "i n s i g h t", "Calibri", Pt(14), False, C_GOLD_LIGHT, PP_ALIGN.CENTER)
    add_textbox(slide, MARGIN, Emu(4434840), CONTENT_W, Emu(274320),
                date_text, "Calibri", Pt(10), False,
                RGBColor(0x8B, 0x73, 0x55), PP_ALIGN.CENTER)
    return slide


def add_end_slide(prs, summary_items, next_video_title, next_video_desc):
    """共通エンドスライド - 次の動画への誘導を重視"""
    slide = new_slide(prs)
    add_rect(slide, 0, 0, Emu(SLIDE_WIDTH), Emu(54864), C_GOLD_LIGHT)
    add_rect(slide, 0, Emu(5088636), Emu(SLIDE_WIDTH), Emu(54864), C_GOLD_LIGHT)

    add_textbox(slide, MARGIN, Emu(274320), CONTENT_W, Emu(411480),
                "今日のまとめ", "Calibri", Pt(22), True, C_BROWN, PP_ALIGN.CENTER)

    y = Emu(731520)
    for i, s in enumerate(summary_items):
        num = str(i + 1)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(1828800), y,
                                        Emu(320040), Emu(320040))
        circle.fill.solid()
        circle.fill.fore_color.rgb = C_GOLD
        circle.line.fill.background()
        add_textbox(slide, Emu(1828800), y + Emu(27432), Emu(320040), Emu(274320),
                    num, "Calibri", Pt(12), True, C_WHITE, PP_ALIGN.CENTER)
        add_textbox(slide, Emu(2286000), y, Emu(5486400), Emu(320040),
                    s, "Calibri", Pt(13), False, C_DARK, PP_ALIGN.LEFT)
        y += Emu(365760)

    # 次の動画への誘導（最重要CTA）
    add_rect(slide, Emu(914400), y + Emu(182880), Emu(7315200), Emu(9144), C_GOLD_LIGHT)

    add_rounded_rect(slide, Emu(1371600), y + Emu(365760), Emu(6400800), Emu(914400), C_ACCENT_BLUE)
    add_textbox(slide, Emu(1463040), y + Emu(411480), Emu(6217920), Emu(320040),
                f"次に見るべき動画 >>>  {next_video_title}",
                "Calibri", Pt(14), True, C_WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, Emu(1463040), y + Emu(731520), Emu(6217920), Emu(457200),
                next_video_desc, "Calibri", Pt(11), False, RGBColor(0xCC, 0xDD, 0xFF), PP_ALIGN.CENTER)

    # branding
    add_textbox(slide, MARGIN, Emu(4571040), CONTENT_W, Emu(228600),
                "H A R M O N I C   i n s i g h t", "Calibri", Pt(10), False, C_BROWN, PP_ALIGN.CENTER)
    return slide
