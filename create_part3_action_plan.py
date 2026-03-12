#!/usr/bin/env python3
"""
後編：学費攻略＋学校選び実践術＋行動プラン（約12分 / 8スライド）
YouTube SEOキーワード: 通信制高校 学費 就学支援金 申請方法 学校選び方
"""
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx_helpers import *

SERIES = "通信制高校選び完全ガイド【後編】"
prs = create_presentation()

# ════════════════════════════════════════════════════
# SLIDE 1: タイトル
# ════════════════════════════════════════════════════
slide = add_title_slide(
    prs,
    "通信制高校の学費攻略＋行動プラン",
    "知らないと損する支援制度と今日からの一歩",
    "〜 通信制高校選び完全ガイド【後編】 〜",
    "Harmonic Insight 2026年3月"
)

# フィードバック②: 問題提起から入る
add_notes(slide,
    "「通信制高校って、結局いくらかかるの？」\n"
    "「うちの家計で通わせられるの？」\n"
    "保護者の方なら、真っ先にそう思いますよね。\n\n"
    "結論から言います。\n"
    "就学支援金を使えば、私立の通信制でも\n"
    "月額1万円以下で通える可能性があります。\n"
    "スマホの月額料金と同じくらいです。\n\n"
    "こんにちは、Harmonic Insightです。\n"
    "この動画は「通信制高校選び完全ガイド」の後編です。\n"
    "学費の不安を解消する支援制度の話と、\n"
    "後悔しない学校選びの実践テクニック、\n"
    "そして今日からすぐ始められる行動プランをお伝えします。\n\n"
    "前編・中編をまだ見ていない方は、概要欄にリンクがありますので\n"
    "ぜひそちらも合わせてご覧ください。"
)


# ════════════════════════════════════════════════════
# SLIDE 2: 学費と支援制度
# フィードバック⑦: 月額換算 + 比喩で「自分ごと化」
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "通信制高校の学費、実際いくら？",
              "支援制度を使えば、スマホ代と同じくらい")
add_footer(slide, 1, SERIES)

# 学費テーブル（月額換算つき）
add_rounded_rect(slide, MARGIN, Emu(960120), CONTENT_W, Emu(1645920), C_WHITE)
add_textbox(slide, Emu(548640), Emu(960120), CONTENT_W, Emu(274320),
            "学費の目安（支援金適用前 → 適用後）", "Calibri", Pt(13), True, C_GOLD, PP_ALIGN.LEFT)

headers = ["学校タイプ", "年間学費", "支援金適用後", "月額換算"]
col_w = Emu(2057400)
for i, h in enumerate(headers):
    x = MARGIN + col_w * i
    add_rect(slide, x, Emu(1234440), col_w, Emu(274320), C_GOLD)
    add_textbox(slide, x, Emu(1234440), col_w, Emu(274320),
                h, "Calibri", Pt(10), True, C_WHITE, PP_ALIGN.CENTER)

rows = [
    ("公立通信制", "約3〜5万円", "実質0〜2万円", "月0〜1,700円"),
    ("私立（ネット型）", "約30〜50万円", "約0〜14万円", "月0〜1.2万円"),
    ("私立（通学型）", "約70〜120万円", "約30〜80万円", "月2.5〜6.7万円"),
]
for ri, row in enumerate(rows):
    bg = C_BG_LIGHT if ri % 2 == 0 else C_WHITE
    y = Emu(1508760) + Emu(274320) * ri
    for ci, val in enumerate(row):
        x = MARGIN + col_w * ci
        clr = C_ACCENT_GREEN if ci == 3 else C_DARK
        bld = ci == 3
        add_rect(slide, x, y, col_w, Emu(274320), bg)
        add_textbox(slide, x, y, col_w, Emu(274320),
                    val, "Calibri", Pt(10), bld, clr, PP_ALIGN.CENTER)

# 比喩で伝える
add_rounded_rect(slide, MARGIN, Emu(2468880), CONTENT_W, Emu(640080), C_ACCENT_GREEN)
add_textbox(slide, Emu(548640), Emu(2468880), Emu(8046720), Emu(320040),
            "つまり...共働き世帯年収590万円以下なら",
            "Calibri", Pt(12), False, C_WHITE, PP_ALIGN.CENTER)
add_textbox(slide, Emu(548640), Emu(2743200), Emu(8046720), Emu(365760),
            "私立ネット型でも月額1万円以下 ≒ スマホ代と同じくらいで通える！",
            "Calibri", Pt(16), True, C_WHITE, PP_ALIGN.CENTER)

# 支援制度リスト
add_rounded_rect(slide, MARGIN, Emu(3291840), CONTENT_W, Emu(1280160), C_WHITE)
add_textbox(slide, Emu(548640), Emu(3291840), CONTENT_W, Emu(274320),
            "使える支援制度", "Calibri", Pt(13), True, C_ACCENT_BLUE, PP_ALIGN.LEFT)

supports = [
    ("高等学校等就学支援金", "年収910万円未満 → 私立は最大39.6万円/年"),
    ("都道府県の上乗せ助成金", "自治体により追加支給あり（要確認）"),
    ("学校独自の奨学金・特待生", "成績や特技で授業料免除の可能性"),
]
y = Emu(3566160)
for name, desc in supports:
    add_textbox(slide, Emu(640080), y, Emu(2651760), Emu(274320),
                name, "Calibri", Pt(10), True, C_DARK, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(3383280), y, Emu(5212080), Emu(274320),
                desc, "Calibri", Pt(9), False, C_GRAY, PP_ALIGN.LEFT)
    y += Emu(274320)

add_notes(slide,
    "学費の話をしましょう。\n"
    "まず結論から。就学支援金を使えば、思った以上に安く通えます。\n\n"
    "公立の通信制なら、支援金適用後は実質ゼロ円に近い。月額ゼロ〜1,700円です。\n"
    "私立のネット型でも、支援金を使えば月額ゼロ〜1万2千円程度。\n"
    "共働きで世帯年収590万円以下の家庭なら、\n"
    "私立の通信制でもスマホの月額料金と同じくらいで通えるんです。\n\n"
    "活用できる支援制度は3つ。\n"
    "1つ目は国の就学支援金。年収910万円未満なら、私立は最大39万6千円が支給されます。\n"
    "2つ目は都道府県の上乗せ助成金。東京都や大阪府など、追加で支給している自治体があります。\n"
    "3つ目は学校独自の奨学金や特待生制度。成績や特技で授業料が免除になる場合もあります。\n\n"
    "次のスライドで、具体的な申請方法をお伝えします。"
)


# ════════════════════════════════════════════════════
# SLIDE 3: 支援制度申請ガイド
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "就学支援金の申請方法【3ステップ】",
              "申請漏れゼロ！手続きは意外と簡単")
add_footer(slide, 2, SERIES)

# 3ステップ
step_cards = [
    ("STEP 1", "書類を受け取る", "入学時に学校から\n申請書類が配布される\n（転入の場合も同様）", C_ACCENT_BLUE),
    ("STEP 2", "所得情報を提出", "マイナンバーカード等で\n保護者の所得証明を提出\n（オンライン申請も可能）", C_ACCENT_GREEN),
    ("STEP 3", "自動で減額", "審査後、授業料から\n支援金が差し引かれる\n（入金手続きは不要）", C_ACCENT_ORANGE),
]

card_w = Emu(2651760)
gap = Emu(182880)
for i, (step, title, desc, color) in enumerate(step_cards):
    x = MARGIN + (card_w + gap) * i
    y = Emu(1005840)
    add_rounded_rect(slide, x, y, card_w, Emu(2103120), C_WHITE)
    add_rounded_rect(slide, x + Emu(91440), y + Emu(91440), Emu(914400), Emu(320040), color)
    add_textbox(slide, x + Emu(91440), y + Emu(91440), Emu(914400), Emu(320040),
                step, "Calibri", Pt(11), True, C_WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, x + Emu(91440), y + Emu(502920), card_w - Emu(182880), Emu(365760),
                title, "Calibri", Pt(16), True, C_DARK, PP_ALIGN.CENTER)
    add_rect(slide, x + Emu(457200), y + Emu(868680), card_w - Emu(914400), Emu(9144), C_BG_LIGHT)
    add_textbox(slide, x + Emu(137160), y + Emu(960120), card_w - Emu(274320), Emu(914400),
                desc, "Calibri", Pt(11), False, C_DARK, PP_ALIGN.CENTER)

# 注意事項
add_rounded_rect(slide, MARGIN, Emu(3291840), CONTENT_W, Emu(1280160), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Emu(548640), Emu(3291840), CONTENT_W, Emu(274320),
            "見落としがちな注意点", "Calibri", Pt(13), True, C_ACCENT_ORANGE, PP_ALIGN.LEFT)

warnings = [
    "申請期限は入学後すぐ → 入学前から準備しておこう",
    "都道府県の上乗せ助成金は別途申請が必要",
    "7月頃に収入状況届出（継続届）の提出を忘れずに",
    "転入・編入の場合は前校の在籍期間分が差し引かれる場合あり",
]
wy = Emu(3566160)
for w in warnings:
    add_textbox(slide, Emu(640080), wy, Emu(8046720), Emu(274320),
                f"  {w}", "Calibri", Pt(10), False, C_DARK, PP_ALIGN.LEFT)
    wy += Emu(274320)

add_notes(slide,
    "申請方法は実はとても簡単で、3ステップで完了します。\n\n"
    "ステップ1、入学時に学校から申請書類が配られます。\n"
    "転入の場合も同じく書類が用意されます。\n\n"
    "ステップ2、マイナンバーカード等で保護者の所得情報を提出します。\n"
    "最近はオンラインで申請できるケースも増えています。\n\n"
    "ステップ3、審査後は授業料から自動的に支援金が差し引かれます。\n"
    "自分で入金手続きをする必要はありません。\n\n"
    "注意点が4つあります。\n"
    "1つ目、申請期限は入学後すぐ。入学前から準備しておきましょう。\n"
    "2つ目、都道府県の上乗せ助成金は国の支援金とは別に申請が必要です。\n"
    "3つ目、7月頃に継続届の提出が必要です。忘れると支給が止まります。\n"
    "4つ目、転入の場合は前の学校での在籍期間分が差し引かれることがあります。"
)


# ════════════════════════════════════════════════════
# SLIDE 4: 学校選び3ステップ
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "後悔しない！学校選び3ステップ",
              "ネットの「ランキング」だけで決めない")
add_footer(slide, 3, SERIES)

steps = [
    ("STEP 1", "情報収集", "公式サイトを読み込む\nパンフレットを取り寄せる\nSNSの口コミは参考程度に", C_ACCENT_BLUE),
    ("STEP 2", "体験・比較", "オープンキャンパスに参加\n個別相談で本音を聞く\n最低3校は比較する", C_ACCENT_GREEN),
    ("STEP 3", "家族で対話", "親子で費用・将来を話し合う\n「自分で決めた」実感を大切に\n最終決定は本人に委ねる", C_ACCENT_ORANGE),
]

card_w = Emu(2651760)
gap = Emu(182880)
for i, (step, title, desc, color) in enumerate(steps):
    x = MARGIN + (card_w + gap) * i
    y = Emu(1005840)
    add_rounded_rect(slide, x, y, card_w, Emu(2743200), C_WHITE)
    add_rounded_rect(slide, x + Emu(91440), y + Emu(91440), Emu(914400), Emu(320040), color)
    add_textbox(slide, x + Emu(91440), y + Emu(91440), Emu(914400), Emu(320040),
                step, "Calibri", Pt(11), True, C_WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, x + Emu(91440), y + Emu(502920), card_w - Emu(182880), Emu(411480),
                title, "Calibri", Pt(18), True, C_DARK, PP_ALIGN.CENTER)
    add_rect(slide, x + Emu(457200), y + Emu(914400), card_w - Emu(914400), Emu(9144), C_BG_LIGHT)
    lines = desc.split("\n")
    dy = y + Emu(1005840)
    for line in lines:
        add_textbox(slide, x + Emu(137160), dy, card_w - Emu(274320), Emu(365760),
                    f"  {line}", "Calibri", Pt(11), False, C_DARK, PP_ALIGN.LEFT)
        dy += Emu(365760)

# フィードバック⑧: 親にシェア促進
add_rounded_rect(slide, MARGIN, Emu(3886200), CONTENT_W, Emu(640080), RGBColor(0xE3, 0xF2, 0xFD))
add_textbox(slide, Emu(548640), Emu(3886200), Emu(8046720), Emu(320040),
            "この部分だけでも保護者の方に見てもらいたい！",
            "Calibri", Pt(12), True, C_ACCENT_BLUE, PP_ALIGN.CENTER)
add_textbox(slide, Emu(548640), Emu(4160520), Emu(8046720), Emu(320040),
            "→ この動画のURLをLINEで送る or スクショを撮って見せてください",
            "Calibri", Pt(11), False, C_ACCENT_BLUE, PP_ALIGN.CENTER)

add_notes(slide,
    "学校選びで後悔しないための3ステップです。\n\n"
    "ステップ1は情報収集。公式サイトを丁寧に読み込んで、パンフレットを取り寄せましょう。\n"
    "SNSの口コミは参考程度に。ネガティブな意見が目立ちやすいのがSNSの特性です。\n\n"
    "ステップ2は体験と比較。必ずオープンキャンパスに参加してください。\n"
    "サイトやパンフレットでは分からない「空気感」が体感できます。\n"
    "最低3校は比較することをおすすめします。1校だけでは良し悪しが分かりません。\n\n"
    "ステップ3は家族での対話です。\n"
    "費用のこと、将来のこと、親子で本音を共有してください。\n"
    "そして最終決定は必ず本人に委ねること。\n"
    "「自分で決めた」という実感が、入学後のモチベーションに直結します。\n\n"
    "この部分は保護者の方にもぜひ見ていただきたいです。\n"
    "今すぐこの動画のURLをLINEでお父さんお母さんに送ってあげてください。\n"
    "あるいは、このスライドのスクショを撮って見せてもOKです。"
)


# ════════════════════════════════════════════════════
# SLIDE 5: 親子コミュニケーション
# フィードバック⑧: 「親にこの動画を見せる」動線
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "保護者の方へ ＆ 生徒の皆さんへ",
              "進路を一緒に考えるためのヒント")
add_footer(slide, 4, SERIES)

# Left: 保護者向け
add_rounded_rect(slide, MARGIN, Emu(1005840), Emu(4023360), Emu(2834640), C_WHITE)
add_rect(slide, MARGIN, Emu(1005840), Emu(4023360), Emu(54864), C_ACCENT_BLUE)
add_textbox(slide, Emu(548640), Emu(1097280), Emu(3840480), Emu(320040),
            "保護者の方へ", "Calibri", Pt(14), True, C_ACCENT_BLUE, PP_ALIGN.LEFT)

parent_tips = [
    "  子どもの話を最後まで聞く（否定しない）",
    "  「なぜ通信制？」→「何を学びたい？」に変換",
    "  一緒に学校を調べる（同じ情報を共有）",
    "  最終決定は本人に委ねる",
]
y = Emu(1463040)
for tip in parent_tips:
    add_textbox(slide, Emu(548640), y, Emu(3840480), Emu(320040),
                tip, "Calibri", Pt(11), False, C_DARK, PP_ALIGN.LEFT)
    y += Emu(365760)

# Right: 生徒向け
add_rounded_rect(slide, Emu(4663440), Emu(1005840), Emu(4023360), Emu(2834640), C_WHITE)
add_rect(slide, Emu(4663440), Emu(1005840), Emu(4023360), Emu(54864), C_ACCENT_GREEN)
add_textbox(slide, Emu(4754880), Emu(1097280), Emu(3840480), Emu(320040),
            "生徒の皆さんへ", "Calibri", Pt(14), True, C_ACCENT_GREEN, PP_ALIGN.LEFT)

student_tips = [
    "  自分の気持ちを具体的に言葉にする",
    "  調べた情報を親にも共有する",
    "  不安なことは正直に伝える",
    "  「こうしたい」を明確にする",
]
y = Emu(1463040)
for tip in student_tips:
    add_textbox(slide, Emu(4754880), y, Emu(3840480), Emu(320040),
                tip, "Calibri", Pt(11), False, C_DARK, PP_ALIGN.LEFT)
    y += Emu(365760)

# bottom
add_rounded_rect(slide, MARGIN, Emu(4023360), CONTENT_W, Emu(457200), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(4069080), Emu(8046720), Emu(365760),
            "家族で納得して決めた進路は、入学後の大きな支えになります",
            "Calibri", Pt(12), True, C_GOLD, PP_ALIGN.CENTER)

add_notes(slide,
    "保護者の方に向けたメッセージです。\n"
    "お子さんが「通信制高校に行きたい」と言ったとき、\n"
    "まずは話を最後まで聞いてあげてください。否定せずに受け止める。\n\n"
    "「なぜ通信制なの？」という質問を、\n"
    "「何を学びたいの？」「どんな環境がいい？」に変えるだけで、\n"
    "会話が全く違うものになります。\n\n"
    "一緒に学校のサイトを見たり、オープンキャンパスに同行したり。\n"
    "同じ情報を共有することで、建設的な話し合いができます。\n\n"
    "生徒の皆さんへ。\n"
    "「なんとなく通信制がいい」ではなく、\n"
    "「こういう理由で、この学校に行きたい」と具体的に伝えましょう。\n"
    "目標が明確なほど、親も安心してくれますよ。\n\n"
    "家族みんなが納得した上での決断は、入学後に必ず力になります。"
)


# ════════════════════════════════════════════════════
# SLIDE 6: 先輩ストーリー（コンパクト）
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "先輩たちの多様な成功ストーリー",
              "通信制高校から広がる可能性")
add_footer(slide, 5, SERIES)

stories = [
    ("不登校→再起", "通信制で心身を立て直し\n新たな興味を発見", C_ACCENT_GREEN),
    ("好き→プロ", "eスポーツ・動画編集で\nフリーランスとして活躍", C_ACCENT_BLUE),
    ("通信制→難関大", "自分のペースで学習し\n難関国立大に現役合格", C_ACCENT_PURPLE),
    ("在学中に起業", "高校時代にビジネスを\n立ち上げて成功", C_ACCENT_ORANGE),
]

card_w = Emu(1966440)
gap = Emu(137160)
for i, (title, desc, color) in enumerate(stories):
    x = MARGIN + (card_w + gap) * i
    y = Emu(1005840)
    add_rounded_rect(slide, x, y, card_w, Emu(2286000), C_WHITE)
    add_rect(slide, x, y, card_w, Emu(54864), color)
    cx = x + card_w // 2 - Emu(228600)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, y + Emu(137160),
                                    Emu(457200), Emu(457200))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide, cx, y + Emu(228600), Emu(457200), Emu(274320),
                str(i + 1), "Georgia", Pt(18), True, C_WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, x + Emu(45720), y + Emu(685800), card_w - Emu(91440), Emu(365760),
                title, "Calibri", Pt(13), True, C_DARK, PP_ALIGN.CENTER)
    add_textbox(slide, x + Emu(45720), y + Emu(1097280), card_w - Emu(91440), Emu(914400),
                desc, "Calibri", Pt(10), False, C_GRAY, PP_ALIGN.CENTER)

# bottom
add_rounded_rect(slide, MARGIN, Emu(3474720), CONTENT_W, Emu(457200), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(3520440), Emu(8046720), Emu(365760),
            "「回り道」も立派な選択肢。既存のレールにとらわれない生き方がある。",
            "Calibri", Pt(12), True, C_GOLD, PP_ALIGN.CENTER)

add_notes(slide,
    "最後に、実際に通信制高校を選んだ先輩たちのストーリーをお伝えします。\n\n"
    "1人目は、不登校から再起した先輩。\n"
    "通信制の安心できる環境で心身を立て直し、新しい興味を見つけました。\n\n"
    "2人目は、好きなことでプロになった先輩。\n"
    "eスポーツや動画編集のスキルを在学中に磨き、フリーランスとして活躍中です。\n\n"
    "3人目は、通信制から難関国立大に現役合格した先輩。\n"
    "通信制の時間の自由さを活かして、戦略的に受験勉強に取り組みました。\n\n"
    "4人目は、在学中に起業した先輩。\n"
    "高校時代からビジネスを立ち上げ、今も活躍しています。\n\n"
    "「回り道」だって立派な選択肢です。\n"
    "あなたの未来は、あなた自身が創るものです。"
)


# ════════════════════════════════════════════════════
# SLIDE 7: 今日やること1つだけ
# フィードバック⑨: CTAを1つに絞る
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "今日やることは、たった1つだけ")
add_footer(slide, 6, SERIES)

# 大きなCTA
add_rounded_rect(slide, Emu(1371600), Emu(1188720), Emu(6400800), Emu(1554480), C_ACCENT_BLUE)
add_textbox(slide, Emu(1463040), Emu(1371600), Emu(6217920), Emu(502920),
            "気になった学校の公式サイトを開いて",
            "Calibri", Pt(22), True, C_WHITE, PP_ALIGN.CENTER)
add_textbox(slide, Emu(1463040), Emu(1874520), Emu(6217920), Emu(502920),
            "「資料請求」ボタンを押す",
            "Georgia", Pt(28), True, C_WHITE, PP_ALIGN.CENTER)
add_textbox(slide, Emu(1463040), Emu(2377440), Emu(6217920), Emu(365760),
            "それだけでいいです。",
            "Calibri", Pt(16), False, C_WHITE, PP_ALIGN.CENTER)

# 補足
add_rounded_rect(slide, MARGIN, Emu(2926080), CONTENT_W, Emu(1554480), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(2971800), Emu(8046720), Emu(320040),
            "資料請求すると...", "Calibri", Pt(14), True, C_GOLD, PP_ALIGN.LEFT)

benefits = [
    "  パンフレットが届くので家族で一緒に見られる",
    "  オープンキャンパスや説明会の日程が分かる",
    "  学費の詳細と支援金の情報が手に入る",
    "  「自分のための学校選び」が始まった実感が持てる",
]
by = Emu(3291840)
for b in benefits:
    add_textbox(slide, Emu(640080), by, Emu(8046720), Emu(274320),
                b, "Calibri", Pt(11), False, C_DARK, PP_ALIGN.LEFT)
    by += Emu(274320)

add_notes(slide,
    "さて、ここまで聞いて「いろいろ分かったけど、何から始めればいいの？」\n"
    "と思った方もいるかもしれません。\n\n"
    "今日やることは、たった1つだけです。\n"
    "気になった学校の公式サイトを開いて、「資料請求」ボタンを押してください。\n"
    "それだけでいいです。\n\n"
    "資料請求すると、パンフレットが届くので家族で一緒に見られます。\n"
    "オープンキャンパスの日程も分かります。\n"
    "学費の詳細や支援金の情報も手に入ります。\n\n"
    "そして何より、「自分の学校選びが始まった」という実感が持てます。\n"
    "完璧な選択をする必要はありません。まずは一歩踏み出すことが大切です。\n"
    "チェックリストやFAQの詳細は概要欄に載せていますので、\n"
    "そちらもぜひチェックしてくださいね。"
)


# ════════════════════════════════════════════════════
# SLIDE 8: エンドカード
# フィードバック⑩: 次の動画への誘導
# ════════════════════════════════════════════════════
slide = add_end_slide(
    prs,
    summary_items=[
        "就学支援金で私立でも月額1万円以下の可能性あり",
        "オープンキャンパスは最低3校参加しよう",
        "今日やること → 資料請求ボタンを押すだけ",
    ],
    next_video_title="海外オンライン教育と日本の比較",
    next_video_desc="Dタイプの方は必見！世界の教育トレンドを徹底比較"
)

add_notes(slide,
    "後編のまとめです。\n"
    "1つ目、就学支援金を使えば、私立の通信制でも月額1万円以下で通える可能性があります。\n"
    "2つ目、オープンキャンパスは最低3校参加。自分の目で確かめましょう。\n"
    "3つ目、今日やることは1つだけ。資料請求ボタンを押してください。\n\n"
    "3本の動画を通してお伝えしてきた「通信制高校選び完全ガイド」、\n"
    "いかがでしたか？\n\n"
    "前編でタイプを診断してまだ見ていない方は、概要欄からぜひチェックしてください。\n"
    "Dタイプだった方には、海外の教育事情をまとめた動画もおすすめです。\n"
    "画面に表示されている動画をクリックしてくださいね。\n\n"
    "コメント欄で「資料請求しました！」と報告してくれたら嬉しいです。\n"
    "皆さんの学校選びを応援しています。ありがとうございました。"
)

# ── 保存 ──
output = "通信制高校ガイド_後編_学費攻略と行動プラン.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
