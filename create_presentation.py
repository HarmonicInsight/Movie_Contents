#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
建設業基幹システム カスタマイズ多発ポイント - 勉強会用パワーポイント生成スクリプト
Harmonic Insight テンプレートスタイルに準拠
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ===== Color Palette (Harmonic Insight Style) =====
GOLD = RGBColor(0x8B, 0x75, 0x36)
GOLD_LIGHT = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_ACCENT = RGBColor(0xB8, 0x86, 0x0B)
DARK_BG = RGBColor(0x1A, 0x1A, 0x1A)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
CREAM2 = RGBColor(0xF5, 0xF0, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
TEXT_BROWN = RGBColor(0x3C, 0x2F, 0x1E)
BROWN_LIGHT = RGBColor(0x8B, 0x73, 0x55)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
BLUE_ACCENT = RGBColor(0x2C, 0x7A, 0xB0)
ORANGE_ACCENT = RGBColor(0xE6, 0x7E, 0x22)
TABLE_HEADER_BG = RGBColor(0x8B, 0x75, 0x36)
TABLE_ROW_LIGHT = RGBColor(0xF9, 0xF6, 0xF0)
TABLE_ROW_DARK = RGBColor(0xF0, 0xE8, 0xD5)

SLIDE_WIDTH = Emu(9144000)  # 10 inches
SLIDE_HEIGHT = Emu(5143500)  # 5.625 inches

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ===== Helper Functions =====
def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name=None):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    if font_name:
        p.font.name = font_name
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, default_size=14, default_color=TEXT_DARK):
    """lines: list of (text, size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        text = line_info[0]
        size = line_info[1] if len(line_info) > 1 else default_size
        color = line_info[2] if len(line_info) > 2 else default_color
        bold = line_info[3] if len(line_info) > 3 else False
        align = line_info[4] if len(line_info) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
    return txBox

def add_footer(slide, page_num):
    add_text_box(slide, Inches(0.3), Inches(5.2), Inches(3), Inches(0.3),
                 "H A R M O N I C   i n s i g h t", 8, WHITE)
    add_text_box(slide, Inches(9.2), Inches(5.2), Inches(0.5), Inches(0.3),
                 str(page_num), 8, WHITE)
    # Footer bar
    add_shape(slide, Emu(0), Emu(4914900), SLIDE_WIDTH, Emu(228600), GOLD)

def add_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text

def add_title_bar(slide, title_text):
    add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, Emu(685800), GOLD)
    add_text_box(slide, Inches(0.5), Inches(0.08), Inches(9), Inches(0.6),
                 title_text, 22, WHITE, True, PP_ALIGN.LEFT)

def add_arrow(slide, left, top, width, height, fill_color=GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_chevron(slide, left, top, width, height, fill_color=GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_oval(slide, left, top, width, height, fill_color=GOLD):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_table_to_slide(slide, left, top, width, height, rows_data, col_widths=None):
    """rows_data: list of lists. First row = header."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    
    for r_idx, row_data in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)
                if r_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = TEXT_DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Cell fill
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_DARK
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_LIGHT
    return table_shape

# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, CREAM)
add_shape(slide, Emu(0), Emu(0), Emu(3200000), SLIDE_HEIGHT, GOLD_ACCENT)
add_shape(slide, Emu(3200000), Emu(0), Emu(100000), SLIDE_HEIGHT, GOLD_LIGHT)

add_text_box(slide, Inches(0.3), Inches(1.5), Inches(3.0), Inches(1.0),
             "Hi", 22, GOLD_LIGHT, font_name="Calibri")

add_multiline_text(slide, Inches(3.8), Inches(0.8), Inches(5.8), Inches(3.5), [
    ("建設業基幹システム", 32, TEXT_BROWN, True),
    ("リプレースにおける", 32, TEXT_BROWN, True),
    ("カスタマイズ多発ポイント", 32, TEXT_BROWN, True),
    ("", 12, TEXT_DARK),
    ("～ 業務担当者のための わかりやすい解説 ～", 16, GOLD),
])

add_multiline_text(slide, Inches(3.8), Inches(3.8), Inches(3), Inches(0.6), [
    ("H A R M O N I C", 14, TEXT_BROWN),
])
add_multiline_text(slide, Inches(5.8), Inches(3.8), Inches(3), Inches(0.6), [
    ("i n s i g h t", 14, GOLD_LIGHT),
])
add_text_box(slide, Inches(3.8), Inches(4.5), Inches(5), Inches(0.3),
             "Harmonic Insight 2026年3月12日", 10, BROWN_LIGHT)
add_notes(slide, """皆さん、こんにちは。Harmonic Insightの勉強会にようこそお越しくださいました。
本日は「建設業基幹システム リプレースにおけるカスタマイズ多発ポイント」について、業務担当者の皆さんにわかりやすく解説します。
建設業のシステムリプレースでは、必ずと言っていいほどカスタマイズの問題が発生します。
今日はその原因と解決策を一緒に考えていきましょう。""")

# ============================================================
# SLIDE 2: 今回の勉強会の目的
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "今回の勉強会の目的")

items = [
    ("カスタマイズ地獄とは?", "なぜ建設業でカスタマイズが多発するのかを理解する"),
    ("8つの業務領域", "カスタマイズが集中する業務領域を把握する"),
    ("標準化の考え方", "カスタマイズに頼らない解決策を知る"),
    ("明日からの行動", "自社の業務で何ができるかを考える"),
]

for i, (title, desc) in enumerate(items):
    y = Inches(1.2) + i * Inches(1.0)
    add_rounded_rect(slide, Inches(0.5), y, Inches(0.7), Inches(0.7), CREAM2)
    add_text_box(slide, Inches(0.55), y + Emu(50000), Inches(0.6), Inches(0.5),
                 str(i+1), 20, GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), y, Inches(3), Inches(0.4),
                 title, 14, TEXT_DARK, True)
    add_text_box(slide, Inches(1.5), y + Emu(230000), Inches(4), Inches(0.4),
                 desc, 11, TEXT_GRAY)

# Right side - key message box
add_rounded_rect(slide, Inches(5.8), Inches(1.2), Inches(3.8), Inches(3.8), CREAM)
add_text_box(slide, Inches(6.0), Inches(1.4), Inches(3.4), Inches(0.5),
             "本勉強会のゴール", 13, GOLD, True, PP_ALIGN.CENTER)
add_multiline_text(slide, Inches(6.0), Inches(2.0), Inches(3.4), Inches(2.5), [
    ("「うちは特殊だから」", 14, TEXT_DARK, True),
    ("↓", 18, GOLD, True, PP_ALIGN.CENTER),
    ("「標準化で解決できる」", 14, GREEN_ACCENT, True),
    ("", 10, TEXT_DARK),
    ("この発想の転換が", 11, TEXT_GRAY),
    ("最大の目標です", 11, TEXT_GRAY),
])
add_footer(slide, 1)
add_notes(slide, """本日の勉強会の目的は大きく4つあります。
まず、カスタマイズ地獄とは何か。なぜ建設業でカスタマイズが多発するのかを理解しましょう。
次に、カスタマイズが集中する8つの業務領域を把握します。原価管理から帳票まで、具体的に見ていきます。
そして、カスタマイズに頼らない標準化の考え方を学びます。
最後に、明日から自社の業務で何ができるかを考えていただきます。
右側に今日のゴールを書いています。「うちは特殊だから」という発想を「標準化で解決できる」に変えること。この発想の転換が最大の目標です。""")

# ============================================================
# SLIDE 3: なぜ「カスタマイズ地獄」に陥るのか?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, 'なぜ建設業は「カスタマイズ地獄」に陥るのか?')

# Illustration: Vicious cycle diagram
center_x = Inches(5.0)
center_y = Inches(3.0)

# Draw cycle arrows concept with boxes
cycle_items = [
    (Inches(1.5), Inches(1.3), "パッケージ\nシステム導入", CREAM),
    (Inches(4.2), Inches(1.3), "「うちは特殊」\nの声が続出", RGBColor(0xFF, 0xEB, 0xEB)),
    (Inches(7.0), Inches(1.3), "大量の\nカスタマイズ要求", RGBColor(0xFF, 0xE0, 0xE0)),
    (Inches(7.0), Inches(3.3), "費用・時間が\n大幅に超過", RGBColor(0xFF, 0xD5, 0xD5)),
    (Inches(4.2), Inches(3.3), "個別開発と\n変わらない結果", RGBColor(0xFF, 0xCA, 0xCA)),
    (Inches(1.5), Inches(3.3), "次のリプレースで\nまた同じ問題", RGBColor(0xFF, 0xC0, 0xC0)),
]

for x, y, text, color in cycle_items:
    add_rounded_rect(slide, x, y, Inches(2.2), Inches(1.2), color)
    lines = text.split('\n')
    for li, line in enumerate(lines):
        add_text_box(slide, x + Emu(50000), y + Emu(150000 + li * 250000), Inches(2.0), Inches(0.4),
                     line, 12, TEXT_DARK, True, PP_ALIGN.CENTER)

# Arrows between boxes
add_arrow(slide, Inches(3.5), Inches(1.7), Inches(0.7), Inches(0.3), GOLD)
add_arrow(slide, Inches(6.2), Inches(1.7), Inches(0.7), Inches(0.3), RED_ACCENT)
# Down arrow (right side)
shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(7.8), Inches(2.5), Inches(0.3), Inches(0.7))
shape.fill.solid()
shape.fill.fore_color.rgb = RED_ACCENT
shape.line.fill.background()
# Left arrows
shape = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(6.2), Inches(3.7), Inches(0.7), Inches(0.3))
shape.fill.solid()
shape.fill.fore_color.rgb = RED_ACCENT
shape.line.fill.background()
shape = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(3.5), Inches(3.7), Inches(0.7), Inches(0.3))
shape.fill.solid()
shape.fill.fore_color.rgb = RED_ACCENT
shape.line.fill.background()

add_text_box(slide, Inches(3.5), Inches(4.7), Inches(3.5), Inches(0.4),
             "この悪循環を断ち切るのが「標準化戦略」", 13, GOLD, True, PP_ALIGN.CENTER)
add_footer(slide, 2)
add_notes(slide, """なぜ建設業はカスタマイズ地獄に陥るのか。その悪循環の構造を見てみましょう。
まず、パッケージシステムを導入します。しかしすぐに「うちは特殊だ」という声が現場から続出します。
すると大量のカスタマイズ要求が発生し、費用と時間が大幅に超過。
結果として、個別開発と変わらないシステムが出来上がり、次のリプレースでまた同じ問題が繰り返される。
この悪循環を断ち切るのが「標準化戦略」です。
大切なのは、「うちは特殊」と感じるその業務が、実は業界全体で共通の課題であるケースが非常に多いということ。
それを今日の勉強会で実感していただきたいと思います。""")

# ============================================================
# SLIDE 4: 建設業特有の6つの複雑性
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "建設業特有の6つの複雑性")

items_6 = [
    ("工事の個別性", "同じ工事は二つとない\n「一品生産」の世界"),
    ("現場と本社の壁", "地理的・時間的な隔たり\n情報共有が困難"),
    ("多様な契約形態", "請負・単価・委託など\n複雑な決済ルール"),
    ("法令・管理要件", "建設業法・下請法など\n厳格な法令対応"),
    ("業務の属人化", "ベテランの経験と勘\n暗黙知の形式知化が困難"),
    ("会計の複雑化", "進行基準・完成基準\n未成工事の管理"),
]

colors_6 = [GOLD, GOLD_ACCENT, BLUE_ACCENT, RED_ACCENT, ORANGE_ACCENT, GREEN_ACCENT]

for i, (title, desc) in enumerate(items_6):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + col * Inches(3.2)
    y = Inches(1.2) + row * Inches(2.0)
    
    add_rounded_rect(slide, x, y, Inches(2.8), Inches(1.7), CREAM)
    # Number circle
    add_oval(slide, x + Emu(50000), y + Emu(50000), Inches(0.4), Inches(0.4), colors_6[i])
    add_text_box(slide, x + Emu(50000), y + Emu(55000), Inches(0.4), Inches(0.4),
                 str(i+1), 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Emu(350000), y + Emu(60000), Inches(2.0), Inches(0.4),
                 title, 14, TEXT_DARK, True)
    lines = desc.split('\n')
    for li, line in enumerate(lines):
        add_text_box(slide, x + Emu(100000), y + Emu(400000 + li * 220000), Inches(2.5), Inches(0.3),
                     line, 10, TEXT_GRAY)

add_text_box(slide, Inches(1.0), Inches(4.8), Inches(8.0), Inches(0.3),
             "これらの特性がシステムの「標準機能」では対応しきれないカスタマイズ要求を生み出す", 11, GOLD, True, PP_ALIGN.CENTER)
add_footer(slide, 3)
add_notes(slide, """建設業にはシステム標準化を難しくする6つの特有の複雑性があります。
1つ目は工事の個別性。同じ工事は二つとない一品生産の世界です。工場のように同じ製品を大量生産するのとは根本的に違います。
2つ目は現場と本社の壁。地理的・時間的な隔たりがあり、情報共有が困難です。
3つ目は多様な契約形態。請負、単価、委託など、複雑な決済ルールが存在します。
4つ目は法令・管理要件。建設業法や下請法など、厳格な法令対応が求められます。
5つ目は業務の属人化。ベテランの経験と勘に頼っており、暗黙知を形式知にするのが困難です。
6つ目は会計の複雑化。進行基準と完成基準の使い分け、未成工事の管理など、独特の会計処理があります。
これらの特性が、標準機能では対応しきれないカスタマイズ要求を生み出しているのです。""")

# ============================================================
# SLIDE 5: Section Divider - カスタマイズ多発 8つの業務領域
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), Emu(2926080), SLIDE_HEIGHT, GOLD)
add_shape(slide, Emu(2926080), Emu(0), Emu(6217920), SLIDE_HEIGHT, DARK_BG)

add_text_box(slide, Inches(3.5), Inches(1.0), Inches(6.0), Inches(1.0),
             "カスタマイズが多発する", 28, WHITE, False)
add_text_box(slide, Inches(3.5), Inches(1.8), Inches(6.0), Inches(1.0),
             "8つの業務領域", 36, GOLD_LIGHT, True)
add_text_box(slide, Inches(3.5), Inches(3.0), Inches(6.0), Inches(1.5),
             "各領域の「よくある課題」と「標準化の方向性」を\n分かりやすく解説します", 14, RGBColor(0xAA, 0xAA, 0xAA))
add_text_box(slide, Inches(3.5), Inches(4.5), Inches(3), Inches(0.3),
             "H A R M O N I C   i n s i g h t", 8, RGBColor(0x88, 0x88, 0x88))
add_notes(slide, """ここからが本日のメインコンテンツです。
カスタマイズが多発する8つの業務領域を、一つずつ見ていきます。
各領域の「よくある課題」と「標準化の方向性」を、できるだけ分かりやすく解説します。
業務担当者の皆さんにとって、「あるある」と感じる内容が多いと思います。
ぜひ自社の状況と照らし合わせながら聞いてください。""")

# ============================================================
# SLIDE 6: 8領域マップ（全体俯瞰）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "カスタマイズ多発 8つの業務領域マップ")

areas = [
    ("1", "原価管理\n工事採算", GOLD_ACCENT, "最重要"),
    ("2", "契約・請求\n入金管理", BLUE_ACCENT, "複雑"),
    ("3", "経費精算\n現場管理", GREEN_ACCENT, "頻出"),
    ("4", "購買・在庫\n資材管理", ORANGE_ACCENT, "現場密着"),
    ("5", "ワークフロー\n承認プロセス", RGBColor(0x8E, 0x44, 0xAD), "全社横断"),
    ("6", "出面管理\n労務管理", RGBColor(0x16, 0xA0, 0x85), "法令関連"),
    ("7", "個別帳票\nレポート", RGBColor(0xD3, 0x54, 0x00), "要望多"),
    ("8", "標準化\n戦略", RED_ACCENT, "解決策"),
]

for i, (num, label, color, tag) in enumerate(areas):
    col = i % 4
    row = i // 4
    x = Inches(0.3) + col * Inches(2.4)
    y = Inches(1.1) + row * Inches(2.1)
    
    add_rounded_rect(slide, x, y, Inches(2.1), Inches(1.8), CREAM)
    add_oval(slide, x + Inches(0.75), y + Emu(80000), Inches(0.55), Inches(0.55), color)
    add_text_box(slide, x + Inches(0.75), y + Emu(85000), Inches(0.55), Inches(0.55),
                 num, 18, WHITE, True, PP_ALIGN.CENTER)
    lines = label.split('\n')
    for li, line in enumerate(lines):
        add_text_box(slide, x + Emu(50000), y + Emu(550000 + li * 220000), Inches(2.0), Inches(0.3),
                     line, 12, TEXT_DARK, True, PP_ALIGN.CENTER)
    # Tag
    add_rounded_rect(slide, x + Inches(0.5), y + Inches(1.4), Inches(1.1), Inches(0.3), color)
    add_text_box(slide, x + Inches(0.5), y + Inches(1.4), Inches(1.1), Inches(0.3),
                 tag, 9, WHITE, True, PP_ALIGN.CENTER)

add_footer(slide, 4)
add_notes(slide, """カスタマイズが多発する8つの業務領域の全体マップです。
第1が原価管理・工事採算。これが最重要領域で、カスタマイズ要望が最も多い分野です。
第2が契約・請求・入金管理。複雑な取引形態への対応が求められます。
第3が経費精算・現場管理。日常的に頻出する課題です。
第4が購買・在庫・資材管理。現場に密着した処理が必要です。
第5がワークフロー・承認プロセス。全社横断的な課題です。
第6が出面管理・労務管理。法令関連の要件が厳しい分野です。
第7が個別帳票・レポート。要望が非常に多い領域です。
そして第8が標準化戦略。これが解決策となります。
では、第1章から順に見ていきましょう。""")

# ============================================================
# SLIDE 7: 第1章 原価管理・工事採算管理（概要）
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "第1章 原価管理・工事採算管理")

add_text_box(slide, Inches(0.5), Inches(0.95), Inches(9), Inches(0.4),
             "建設業の生命線 ― 見えないコストを可視化する標準化戦略", 13, GOLD, True)

# Left: 3 problem areas
problems = [
    ("実行予算管理", "設計変更のたびにExcel作り直し\n予算残高をリアルタイムで見たい"),
    ("共通費の配賦", "間接費の按分ルールが複雑\n工事間で不公平感が生まれる"),
    ("工事進行基準", "進捗率の計算方法が統一されない\n収益認識が企業ごとにバラバラ"),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(1.5) + i * Inches(1.15)
    add_rounded_rect(slide, Inches(0.3), y, Inches(4.5), Inches(1.0), RGBColor(0xFF, 0xF5, 0xE8))
    add_oval(slide, Inches(0.4), y + Emu(100000), Inches(0.35), Inches(0.35), RED_ACCENT)
    add_text_box(slide, Inches(0.42), y + Emu(100000), Inches(0.35), Inches(0.35),
                 "!", 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.9), y + Emu(50000), Inches(1.5), Inches(0.3),
                 title, 13, TEXT_DARK, True)
    lines = desc.split('\n')
    for li, line in enumerate(lines):
        add_text_box(slide, Inches(0.9), y + Emu(300000 + li * 200000), Inches(3.5), Inches(0.3),
                     line, 9, TEXT_GRAY)

# Right: Solution direction
add_rounded_rect(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(3.4), RGBColor(0xE8, 0xF8, 0xE8))
add_text_box(slide, Inches(5.4), Inches(1.6), Inches(4.0), Inches(0.4),
             "標準化の方向性", 14, GREEN_ACCENT, True)

solutions = [
    "業界標準の原価分類マスタをベースに\n企業独自の費目を追加できる仕組み",
    "配賦パターンをテンプレート化し\n選択式で設定できる機能",
    "進捗率の自動計算ロジックを\n複数パターンから選択可能に",
]

for i, sol in enumerate(solutions):
    y = Inches(2.1) + i * Inches(0.9)
    add_text_box(slide, Inches(5.5), y, Inches(0.3), Inches(0.3),
                 ">>", 11, GREEN_ACCENT, True)
    lines = sol.split('\n')
    for li, line in enumerate(lines):
        add_text_box(slide, Inches(5.9), y + Emu(li * 180000), Inches(3.4), Inches(0.3),
                     line, 10, TEXT_DARK)

add_footer(slide, 5)
add_notes(slide, """第1章、原価管理・工事採算管理です。建設業の生命線とも言える領域です。
左側に3つの主要な問題を挙げています。
まず実行予算管理。設計変更のたびにExcelを作り直し、予算残高をリアルタイムで見たいというニーズが強くあります。
次に共通費の配賦。間接費の按分ルールが複雑で、工事間で不公平感が生まれています。
そして工事進行基準。進捗率の計算方法が統一されず、収益認識が企業ごとにバラバラになっています。
右側に標準化の方向性を示しています。
業界標準の原価分類マスタをベースに企業独自の費目を追加できる仕組み、配賦パターンのテンプレート化、進捗率の自動計算ロジックを複数パターンから選択可能にすること。
こうした標準化により、カスタマイズを最小限に抑えられます。""")

# ============================================================
# SLIDE 8: 第1章 実行予算管理 図解
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "実行予算管理の Before / After")

# Before
add_text_box(slide, Inches(0.3), Inches(0.95), Inches(4.5), Inches(0.4),
             "Before : 現状の課題", 14, RED_ACCENT, True)

before_flow = [
    (Inches(0.3), "設計変更\n発生"),
    (Inches(1.8), "Excel\n作り直し"),
    (Inches(3.3), "本社に\n電話確認"),
]
for x, text in before_flow:
    add_rounded_rect(slide, x, Inches(1.5), Inches(1.3), Inches(0.9), RGBColor(0xFF, 0xE8, 0xE8))
    lines = text.split('\n')
    for li, line in enumerate(lines):
        add_text_box(slide, x + Emu(20000), Inches(1.55) + Emu(li * 200000), Inches(1.2), Inches(0.3),
                     line, 11, TEXT_DARK, False, PP_ALIGN.CENTER)

add_arrow(slide, Inches(1.55), Inches(1.8), Inches(0.3), Inches(0.2), RED_ACCENT)
add_arrow(slide, Inches(3.05), Inches(1.8), Inches(0.3), Inches(0.2), RED_ACCENT)

add_rounded_rect(slide, Inches(0.3), Inches(2.6), Inches(4.3), Inches(0.5), RGBColor(0xFF, 0xDD, 0xDD))
add_text_box(slide, Inches(0.5), Inches(2.65), Inches(4.0), Inches(0.4),
             "時間がかかる / ミスが起きやすい / 属人的", 11, RED_ACCENT, True, PP_ALIGN.CENTER)

# After
add_text_box(slide, Inches(5.2), Inches(0.95), Inches(4.5), Inches(0.4),
             "After : 標準化後のイメージ", 14, GREEN_ACCENT, True)

after_flow = [
    (Inches(5.2), "設計変更\nシステム入力"),
    (Inches(6.7), "予算自動\n再計算"),
    (Inches(8.2), "タブレットで\n残高確認"),
]
for x, text in after_flow:
    add_rounded_rect(slide, x, Inches(1.5), Inches(1.3), Inches(0.9), RGBColor(0xE8, 0xF8, 0xE8))
    lines = text.split('\n')
    for li, line in enumerate(lines):
        add_text_box(slide, x + Emu(20000), Inches(1.55) + Emu(li * 200000), Inches(1.2), Inches(0.3),
                     line, 11, TEXT_DARK, False, PP_ALIGN.CENTER)

add_arrow(slide, Inches(6.45), Inches(1.8), Inches(0.3), Inches(0.2), GREEN_ACCENT)
add_arrow(slide, Inches(7.95), Inches(1.8), Inches(0.3), Inches(0.2), GREEN_ACCENT)

add_rounded_rect(slide, Inches(5.2), Inches(2.6), Inches(4.3), Inches(0.5), RGBColor(0xDD, 0xFF, 0xDD))
add_text_box(slide, Inches(5.4), Inches(2.65), Inches(4.0), Inches(0.4),
             "リアルタイム / 自動化 / 誰でも操作可能", 11, GREEN_ACCENT, True, PP_ALIGN.CENTER)

# User voice
add_rounded_rect(slide, Inches(0.3), Inches(3.3), Inches(9.2), Inches(1.5), CREAM)
add_text_box(slide, Inches(0.5), Inches(3.4), Inches(2.0), Inches(0.3),
             "現場の声", 12, GOLD, True)
add_multiline_text(slide, Inches(0.5), Inches(3.8), Inches(4.2), Inches(0.9), [
    ("「設計変更のたびにExcelを作り直してて", 10, TEXT_GRAY),
    ("  時間がかかりすぎる」", 10, TEXT_GRAY),
    ("「タブレットで予算残高が見られれば", 10, TEXT_GRAY),
    ("  無駄な発注も減るはずだ」", 10, TEXT_GRAY),
])
add_multiline_text(slide, Inches(5.0), Inches(3.8), Inches(4.2), Inches(0.9), [
    ("「工事ごとに細かく原価を管理したいのに", 10, TEXT_GRAY),
    ("  システムだと大まかすぎて使い物にならない」", 10, TEXT_GRAY),
    ("「承認ルートが工事規模で変わるから", 10, TEXT_GRAY),
    ("  標準システムだと対応できない」", 10, TEXT_GRAY),
])
add_footer(slide, 6)
add_notes(slide, """実行予算管理のBefore/Afterを図解しています。
左のBefore。設計変更が発生すると、まずExcelを作り直し、そして本社に電話で確認。時間がかかり、ミスが起きやすく、属人的な作業です。
右のAfter。設計変更をシステムに入力すると、予算が自動で再計算され、タブレットで残高をその場で確認できます。リアルタイムで自動化され、誰でも操作できる世界です。
下に現場の声を載せています。「設計変更のたびにExcelを作り直すのは時間がかかりすぎる」「タブレットで予算残高が見られれば無駄な発注も減るはずだ」。
こうした声は、多くの建設会社で共通して聞かれるものです。つまり、カスタマイズしなくても標準機能で対応できる可能性が高い領域なのです。""")

# ============================================================
# SLIDE 9: 第2章 契約・請求・入金管理
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "第2章 契約・請求・入金管理")

add_text_box(slide, Inches(0.5), Inches(0.95), Inches(9), Inches(0.4),
             "複雑な取引を簡素化する標準化戦略", 13, GOLD, True)

# Contract types table
table_data = [
    ["契約形態", "特徴", "支払方法", "カスタマイズの課題"],
    ["請負契約", "工事全体を一括受注", "出来高払い・完成払い", "出来高の計算方法が企業ごとに異なる"],
    ["単価契約", "数量×単価で精算", "月次精算", "単価改定・追加数量の処理が複雑"],
    ["業務委託", "人工ベースの契約", "月末締め翌月払い", "労務実績との突合が手作業"],
]
add_table_to_slide(slide, Inches(0.3), Inches(1.5), Inches(9.2), Inches(2.0), table_data,
                   [Inches(1.5), Inches(2.3), Inches(2.3), Inches(3.1)])

# Key points
add_rounded_rect(slide, Inches(0.3), Inches(3.7), Inches(4.3), Inches(1.2), RGBColor(0xE8, 0xF0, 0xFF))
add_text_box(slide, Inches(0.5), Inches(3.8), Inches(4.0), Inches(0.3),
             "下請支払管理の重要ポイント", 12, BLUE_ACCENT, True)
add_multiline_text(slide, Inches(0.5), Inches(4.2), Inches(4.0), Inches(0.6), [
    ("建設業法・下請法への確実な対応が必須", 10, TEXT_DARK),
    ("建退共の証紙管理・CCUSとの連携も重要", 10, TEXT_DARK),
])

add_rounded_rect(slide, Inches(5.0), Inches(3.7), Inches(4.5), Inches(1.2), RGBColor(0xE8, 0xF8, 0xE8))
add_text_box(slide, Inches(5.2), Inches(3.8), Inches(4.0), Inches(0.3),
             "標準化の方向性", 12, GREEN_ACCENT, True)
add_multiline_text(slide, Inches(5.2), Inches(4.2), Inches(4.0), Inches(0.6), [
    ("契約パターンマスタで類型化し選択式に", 10, TEXT_DARK),
    ("出来高計算テンプレートの標準提供", 10, TEXT_DARK),
    ("法令チェック機能の標準実装", 10, TEXT_DARK),
])
add_footer(slide, 7)
add_notes(slide, """第2章、契約・請求・入金管理です。
建設業には請負契約、単価契約、業務委託と多様な契約形態があり、それぞれ支払方法も異なります。
請負契約は出来高払いや完成払いですが、出来高の計算方法が企業ごとに異なるのが課題です。
単価契約は月次精算ですが、単価改定や追加数量の処理が複雑になります。
業務委託は月末締め翌月払いですが、労務実績との突合が手作業で行われています。
左下の重要ポイントとして、建設業法・下請法への対応、建退共の証紙管理やCCUSとの連携も必要です。
右下の標準化の方向性としては、契約パターンマスタで類型化して選択式にすること、出来高計算テンプレートの標準提供、そして法令チェック機能の標準実装が有効です。""")

# ============================================================
# SLIDE 10: 第3章 経費精算・現場管理
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "第3章 経費精算・現場管理")

add_text_box(slide, Inches(0.5), Inches(0.95), Inches(9), Inches(0.4),
             "現場のリアルを捉える経費処理の最適化", 13, GOLD, True)

# Two-column layout
# Left: Current problems (building analogy)
add_rounded_rect(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(3.3), RGBColor(0xFF, 0xF5, 0xE8))
add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.0), Inches(0.3),
             "現場で起きていること", 13, ORANGE_ACCENT, True)

expense_items = [
    ("現場経費 vs 一般経費", "どこまでが現場の費用？\n分類基準が人によってバラバラ"),
    ("立替・仮払い", "現場監督が自腹で立て替え\n精算が月末に集中してパンク"),
    ("安全対策費", "安全大会・保護具・標識…\n工事原価に入れる？入れない？"),
]

for i, (title, desc) in enumerate(expense_items):
    y = Inches(2.0) + i * Inches(1.0)
    add_text_box(slide, Inches(0.6), y, Inches(2.0), Inches(0.3),
                 title, 11, TEXT_DARK, True)
    lines = desc.split('\n')
    for li, line in enumerate(lines):
        add_text_box(slide, Inches(0.6), y + Emu(200000 + li * 170000), Inches(4.0), Inches(0.3),
                     line, 9, TEXT_GRAY)

# Right: Solution
add_rounded_rect(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(3.3), RGBColor(0xE8, 0xF8, 0xE8))
add_text_box(slide, Inches(5.4), Inches(1.6), Inches(4.0), Inches(0.3),
             "標準化アプローチ", 13, GREEN_ACCENT, True)

sol_items = [
    "経費分類マスタの標準化\n→ 業界共通の分類ルールを適用",
    "モバイルアプリで現場から即時申請\n→ 写真・GPS情報を自動付与",
    "工事別・費目別の自動配賦\n→ 設定ベースで柔軟に対応",
]

for i, sol in enumerate(sol_items):
    y = Inches(2.1) + i * Inches(1.0)
    add_oval(slide, Inches(5.4), y + Emu(10000), Inches(0.25), Inches(0.25), GREEN_ACCENT)
    add_text_box(slide, Inches(5.43), y + Emu(10000), Inches(0.25), Inches(0.25),
                 str(i+1), 9, WHITE, True, PP_ALIGN.CENTER)
    lines = sol.split('\n')
    for li, line in enumerate(lines):
        add_text_box(slide, Inches(5.8), y + Emu(li * 180000), Inches(3.5), Inches(0.3),
                     line, 10, TEXT_DARK)
add_footer(slide, 8)
add_notes(slide, """第3章、経費精算・現場管理です。現場のリアルを捉える経費処理の最適化が課題です。
左側に現場で起きていることを3つ挙げています。
まず「現場経費と一般経費の区分」。どこまでが現場の費用なのか、分類基準が人によってバラバラです。
次に「立替・仮払い」。現場監督が自腹で立て替え、精算が月末に集中してパンクする。よくある話ですね。
そして「安全対策費」。安全大会や保護具の費用を工事原価に入れるのか入れないのか、判断が分かれます。
右側の標準化アプローチでは3つの解決策を提示しています。
1つ目、経費分類マスタの標準化で業界共通のルールを適用する。
2つ目、モバイルアプリで現場から即時申請。写真やGPS情報を自動付与することで精度を上げる。
3つ目、工事別・費目別の自動配賦を設定ベースで柔軟に対応する。""")

# ============================================================
# SLIDE 11: 第4章 購買・在庫管理
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "第4章 購買・在庫管理")

add_text_box(slide, Inches(0.5), Inches(0.95), Inches(9), Inches(0.4),
             "現場起点で効率化する資材・重機・仮設材管理", 13, GOLD, True)

# Flow diagram: procurement flow
flow_items = [
    (Inches(0.2), "現場で\n資材が必要", ORANGE_ACCENT),
    (Inches(2.1), "現場直送\n都度購入", RGBColor(0xFF, 0xE0, 0xC0)),
    (Inches(4.0), "発注・\n検収処理", RGBColor(0xFF, 0xE0, 0xC0)),
    (Inches(5.9), "在庫管理\n不要?", RGBColor(0xFF, 0xE0, 0xC0)),
    (Inches(7.8), "原価\n計上", RGBColor(0xFF, 0xE0, 0xC0)),
]

for i, (x, text, color) in enumerate(flow_items):
    if i == 0:
        add_rounded_rect(slide, x, Inches(1.4), Inches(1.5), Inches(0.9), color)
        lines = text.split('\n')
        for li, line in enumerate(lines):
            add_text_box(slide, x + Emu(20000), Inches(1.45) + Emu(li * 200000), Inches(1.4), Inches(0.3),
                         line, 11, WHITE, True, PP_ALIGN.CENTER)
    else:
        add_rounded_rect(slide, x, Inches(1.4), Inches(1.5), Inches(0.9), color)
        lines = text.split('\n')
        for li, line in enumerate(lines):
            add_text_box(slide, x + Emu(20000), Inches(1.45) + Emu(li * 200000), Inches(1.4), Inches(0.3),
                         line, 11, TEXT_DARK, False, PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_arrow(slide, x + Inches(1.5), Inches(1.7), Inches(0.5), Inches(0.2), GOLD)

# Two key areas
add_rounded_rect(slide, Inches(0.3), Inches(2.6), Inches(4.5), Inches(2.2), CREAM)
add_text_box(slide, Inches(0.5), Inches(2.7), Inches(4.0), Inches(0.3),
             "現場直送・都度購入の課題", 12, ORANGE_ACCENT, True)
add_multiline_text(slide, Inches(0.5), Inches(3.1), Inches(4.0), Inches(1.5), [
    ("「倉庫を通さず現場に直送するから、", 10, TEXT_GRAY),
    ("  在庫管理の概念が当てはまらない」", 10, TEXT_GRAY),
    ("", 8, TEXT_DARK),
    ("「現場で急に必要になった資材を", 10, TEXT_GRAY),
    ("  その場で発注できないと困る」", 10, TEXT_GRAY),
    ("", 8, TEXT_DARK),
    ("→ 発注=検収=原価計上の一気通貫処理が必要", 10, ORANGE_ACCENT, True),
])

add_rounded_rect(slide, Inches(5.2), Inches(2.6), Inches(4.3), Inches(2.2), CREAM)
add_text_box(slide, Inches(5.4), Inches(2.7), Inches(4.0), Inches(0.3),
             "レンタル・仮設材の課題", 12, BLUE_ACCENT, True)
add_multiline_text(slide, Inches(5.4), Inches(3.1), Inches(4.0), Inches(1.5), [
    ("「重機レンタルの日割り計算が複雑」", 10, TEXT_GRAY),
    ("「仮設材を他の現場に転用したい」", 10, TEXT_GRAY),
    ("", 8, TEXT_DARK),
    ("→ レンタル期間の自動計算テンプレート", 10, BLUE_ACCENT, True),
    ("→ 仮設材の現場間移動トラッキング機能", 10, BLUE_ACCENT, True),
    ("→ 月額・日額の自動按分処理", 10, BLUE_ACCENT, True),
])
add_footer(slide, 9)
add_notes(slide, """第4章、購買・在庫管理です。建設業の購買は、工場の製造業とは全く違う特徴があります。
上部のフロー図をご覧ください。現場で資材が必要になると、倉庫を通さず現場に直送するケースが多い。都度購入が基本で、従来の在庫管理の概念が当てはまりません。
左下の課題にある通り、「倉庫を通さないから在庫管理ができない」「現場で急に必要な資材をその場で発注したい」という声があります。発注イコール検収イコール原価計上の一気通貫処理が必要なのです。
右下はレンタル・仮設材の課題です。重機レンタルの日割り計算が複雑だったり、仮設材を他の現場に転用したいというニーズがあります。
レンタル期間の自動計算テンプレート、仮設材の現場間移動トラッキング、月額・日額の自動按分処理。これらを標準機能として提供することで、カスタマイズを減らせます。""")

# ============================================================
# SLIDE 12: 第5章 ワークフロー・承認プロセス
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "第5章 ワークフロー・承認プロセス")

add_text_box(slide, Inches(0.5), Inches(0.95), Inches(9), Inches(0.4),
             "柔軟性と統制を両立させる電子化", 13, GOLD, True)

# Approval flow diagram
add_text_box(slide, Inches(0.5), Inches(1.5), Inches(3.0), Inches(0.3),
             "建設業の承認フロー例", 12, RGBColor(0x8E, 0x44, 0xAD), True)

# Multi-level approval flow
flow_levels = [
    ("現場担当", RGBColor(0xE8, 0xD5, 0xF5)),
    ("現場所長", RGBColor(0xD5, 0xC0, 0xE8)),
    ("部門長", RGBColor(0xC0, 0xA8, 0xD8)),
    ("経理部", RGBColor(0xAA, 0x90, 0xC8)),
    ("社長", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (label, color) in enumerate(flow_levels):
    x = Inches(0.3) + i * Inches(1.5)
    add_chevron(slide, x, Inches(2.0), Inches(1.4), Inches(0.6), color)
    text_color = WHITE if i == len(flow_levels) - 1 else TEXT_DARK
    add_text_box(slide, x + Emu(200000), Inches(2.05), Inches(1.0), Inches(0.5),
                 label, 10, text_color, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(8.0), Inches(2.05), Inches(1.5), Inches(0.5),
             "金額に応じて\nルート変更!", 10, RED_ACCENT, True)

# Problem and solution table
table_data2 = [
    ["よくある課題", "標準化の方向性"],
    ["金額・工種で承認ルートが変わる", "条件分岐テンプレート（金額・工種マトリックス）"],
    ["現場にいる承認者がなかなか承認できない", "モバイル承認 + 代理承認機能の標準実装"],
    ["緊急時の事後承認フローがない", "緊急承認ワークフローの標準パターン提供"],
    ["紙の稟議書とシステムが二重管理", "紙帳票の電子化テンプレート提供"],
]
add_table_to_slide(slide, Inches(0.3), Inches(2.9), Inches(9.2), Inches(2.0), table_data2,
                   [Inches(4.6), Inches(4.6)])
add_footer(slide, 10)
add_notes(slide, """第5章、ワークフロー・承認プロセスです。柔軟性と統制を両立させる電子化がテーマです。
建設業の承認フローは特殊で、現場担当から現場所長、部門長、経理部、社長と多段階の承認が必要です。しかも、金額によって承認ルートが変わります。
表に4つのよくある課題と標準化の方向性を示しています。
金額や工種で承認ルートが変わる課題には、条件分岐テンプレートで対応。
現場の承認者がなかなか承認できない問題には、モバイル承認と代理承認機能を標準搭載。
緊急時の事後承認フローがない課題には、緊急承認ワークフローの標準パターンを提供。
紙の稟議書とシステムの二重管理には、紙帳票の電子化テンプレートを提供。
これらは多くの建設会社で共通する課題ですので、標準機能として実装すべき内容です。""")

# ============================================================
# SLIDE 13: 第6章 出面管理・労務管理
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "第6章 出面管理・労務管理")

add_text_box(slide, Inches(0.5), Inches(0.95), Inches(9), Inches(0.4),
             "現場の勤怠を正確に把握し、法令遵守を支援する", 13, GOLD, True)

# Payment types comparison
add_text_box(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(0.3),
             "多様な支払形態", 12, RGBColor(0x16, 0xA0, 0x85), True)

payment_table = [
    ["支払形態", "内容", "課題"],
    ["常用（人工）", "日当×日数で支払い", "日報との突合が煩雑"],
    ["月極", "月額固定で支払い", "変動分の追加精算"],
    ["出来高払い", "作業量に応じて支払い", "数量の確認方法"],
]
add_table_to_slide(slide, Inches(0.3), Inches(1.9), Inches(4.5), Inches(1.5), payment_table,
                   [Inches(1.3), Inches(1.8), Inches(1.4)])

# Right side: Legal compliance
add_rounded_rect(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.9), RGBColor(0xE8, 0xF5, 0xF0))
add_text_box(slide, Inches(5.4), Inches(1.6), Inches(4.0), Inches(0.3),
             "法令遵守のポイント", 12, RGBColor(0x16, 0xA0, 0x85), True)
add_multiline_text(slide, Inches(5.4), Inches(2.0), Inches(4.0), Inches(1.2), [
    ("外国人労働者の在留資格管理", 10, TEXT_DARK, True),
    ("  → 期限アラート機能の標準搭載", 10, TEXT_GRAY),
    ("", 6, TEXT_DARK),
    ("建退共の掛金管理・CCUS連携", 10, TEXT_DARK, True),
    ("  → 就業実績から証紙枚数を自動計算", 10, TEXT_GRAY),
    ("", 6, TEXT_DARK),
    ("労働時間の上限規制対応", 10, TEXT_DARK, True),
    ("  → 残業時間の自動集計・アラート", 10, TEXT_GRAY),
])

# Bottom: Solution
add_rounded_rect(slide, Inches(0.3), Inches(3.7), Inches(9.2), Inches(1.1), CREAM)
add_text_box(slide, Inches(0.5), Inches(3.8), Inches(3.0), Inches(0.3),
             "標準化のポイント", 12, GOLD, True)
add_multiline_text(slide, Inches(0.5), Inches(4.1), Inches(8.5), Inches(0.6), [
    ("GPS・顔認証・ICカードなどによる出退勤データの自動取得  →  日報・出面との自動連携  →  支払計算の自動化", 10, TEXT_DARK),
    ("現場の「手書き日報」から脱却し、モバイル入力で正確なデータをリアルタイムに収集する仕組みを標準機能として提供", 10, TEXT_GRAY),
])
add_footer(slide, 11)
add_notes(slide, """第6章、出面管理・労務管理です。現場の勤怠を正確に把握し、法令遵守を支援する領域です。
支払形態が多様なのが建設業の特徴です。常用は日当掛ける日数、月極は月額固定、出来高払いは作業量に応じて支払う。それぞれ異なる課題があります。
右側の法令遵守のポイントは特に重要です。
外国人労働者の在留資格管理では、期限アラート機能が必須です。これを見逃すと法令違反になります。
建退共の掛金管理やCCUS連携では、就業実績から証紙枚数を自動計算する仕組みが必要です。
そして労働時間の上限規制。2024年4月から建設業にも適用された残業規制に対応するため、残業時間の自動集計とアラートが不可欠です。
下の標準化のポイントにある通り、GPSや顔認証、ICカードによる出退勤データの自動取得から、日報との連携、支払計算の自動化まで一気通貫で標準化すべきです。""")

# ============================================================
# SLIDE 14: 第7章 個別帳票・レポート
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "第7章 個別帳票・レポート")

add_text_box(slide, Inches(0.5), Inches(0.95), Inches(9), Inches(0.4),
             "企業独自の管理要件を満たす情報可視化", 13, GOLD, True)

# Problem illustration
add_rounded_rect(slide, Inches(0.3), Inches(1.5), Inches(4.5), Inches(1.8), RGBColor(0xFF, 0xF0, 0xE0))
add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4.0), Inches(0.3),
             "帳票カスタマイズが多発する理由", 12, RGBColor(0xD3, 0x54, 0x00), True)

reasons = [
    "施主ごとに提出フォーマットが違う",
    "社内の管理帳票が独自フォーマット",
    "法定帳票の書式が頻繁に改定される",
    "経営層が「この形で見たい」にこだわる",
]
for i, reason in enumerate(reasons):
    y = Inches(2.0) + i * Inches(0.3)
    add_text_box(slide, Inches(0.7), y, Inches(4.0), Inches(0.3),
                 f"  {reason}", 10, TEXT_DARK)

# Solution
add_rounded_rect(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.8), RGBColor(0xE8, 0xF8, 0xE8))
add_text_box(slide, Inches(5.4), Inches(1.6), Inches(4.0), Inches(0.3),
             "標準化アプローチ", 12, GREEN_ACCENT, True)

solutions_report = [
    ("BIツール連携", "ノーコードでレポート作成"),
    ("帳票テンプレート", "業界標準の帳票を標準搭載"),
    ("Excel出力", "自由な加工を可能に"),
    ("ダッシュボード", "リアルタイムの経営指標"),
]
for i, (title, desc) in enumerate(solutions_report):
    y = Inches(2.0) + i * Inches(0.3)
    add_text_box(slide, Inches(5.5), y, Inches(1.8), Inches(0.3),
                 f"  {title}", 10, GREEN_ACCENT, True)
    add_text_box(slide, Inches(7.2), y, Inches(2.2), Inches(0.3),
                 desc, 10, TEXT_GRAY)

# Bottom: Before/After comparison
add_text_box(slide, Inches(0.3), Inches(3.5), Inches(4.5), Inches(0.3),
             "Before : 帳票ごとに個別開発", 12, RED_ACCENT, True)

# Before boxes
for i, label in enumerate(["工事月報", "原価報告書", "安全書類", "施主報告"]):
    x = Inches(0.3) + i * Inches(1.2)
    add_rounded_rect(slide, x, Inches(3.9), Inches(1.0), Inches(0.5), RGBColor(0xFF, 0xDD, 0xDD))
    add_text_box(slide, x, Inches(3.95), Inches(1.0), Inches(0.4),
                 label, 9, RED_ACCENT, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(5.2), Inches(3.5), Inches(4.5), Inches(0.3),
             "After : テンプレート + BIで自由に", 12, GREEN_ACCENT, True)

add_rounded_rect(slide, Inches(5.2), Inches(3.9), Inches(1.5), Inches(0.5), RGBColor(0xDD, 0xFF, 0xDD))
add_text_box(slide, Inches(5.2), Inches(3.95), Inches(1.5), Inches(0.4),
             "共通データ基盤", 9, GREEN_ACCENT, True, PP_ALIGN.CENTER)
add_arrow(slide, Inches(6.7), Inches(4.05), Inches(0.4), Inches(0.2), GREEN_ACCENT)
add_rounded_rect(slide, Inches(7.2), Inches(3.9), Inches(2.3), Inches(0.5), RGBColor(0xDD, 0xFF, 0xDD))
add_text_box(slide, Inches(7.2), Inches(3.95), Inches(2.3), Inches(0.4),
             "BIツールで自在に出力", 9, GREEN_ACCENT, True, PP_ALIGN.CENTER)

# User voice
add_rounded_rect(slide, Inches(0.3), Inches(4.5), Inches(9.2), Inches(0.4), CREAM)
add_text_box(slide, Inches(0.5), Inches(4.52), Inches(9.0), Inches(0.3),
             "「帳票はもう作るのではなく、選ぶ時代へ」 ― データさえあれば、見せ方は自由に変えられる", 10, GOLD, True, PP_ALIGN.CENTER)
add_footer(slide, 12)
add_notes(slide, """第7章、個別帳票・レポートです。これはカスタマイズ要望が最も多い領域の一つです。
帳票カスタマイズが多発する理由は4つあります。
施主ごとに提出フォーマットが違う、社内の管理帳票が独自フォーマット、法定帳票の書式が頻繁に改定される、そして経営層が「この形で見たい」にこだわる。
標準化アプローチとしては、BIツール連携でノーコードのレポート作成、業界標準の帳票テンプレートの標準搭載、自由な加工ができるExcel出力、そしてリアルタイムの経営指標ダッシュボードです。
下のBefore/Afterを見てください。従来は工事月報、原価報告書、安全書類、施主報告と帳票ごとに個別開発していました。
標準化後は、共通データ基盤からBIツールで自在に出力する形に変わります。
「帳票はもう作るのではなく、選ぶ時代へ」。データさえあれば、見せ方は自由に変えられるのです。""")

# ============================================================
# SLIDE 15: Section Divider - 解決策
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), Emu(2926080), SLIDE_HEIGHT, GOLD)
add_shape(slide, Emu(2926080), Emu(0), Emu(6217920), SLIDE_HEIGHT, DARK_BG)

add_text_box(slide, Inches(3.5), Inches(1.0), Inches(6.0), Inches(1.0),
             "第8章", 28, RGBColor(0xAA, 0xAA, 0xAA))
add_text_box(slide, Inches(3.5), Inches(1.8), Inches(6.0), Inches(1.0),
             "標準化・テンプレート化への", 28, WHITE)
add_text_box(slide, Inches(3.5), Inches(2.5), Inches(6.0), Inches(1.0),
             "戦略的アプローチ", 36, GOLD_LIGHT, True)
add_text_box(slide, Inches(3.5), Inches(3.5), Inches(6.0), Inches(1.0),
             "カスタマイズ地獄から抜け出す5つの鍵", 16, RGBColor(0xAA, 0xAA, 0xAA))
add_text_box(slide, Inches(3.5), Inches(4.5), Inches(3), Inches(0.3),
             "H A R M O N I C   i n s i g h t", 8, RGBColor(0x88, 0x88, 0x88))
add_notes(slide, """ここからは第8章、標準化・テンプレート化への戦略的アプローチに入ります。
7つの業務領域でカスタマイズが多発する課題を見てきましたが、ここからはその解決策です。
カスタマイズ地獄から抜け出す5つの鍵をお伝えします。
この考え方を身につければ、次のシステムリプレースでは全く違うアプローチが可能になります。""")

# ============================================================
# SLIDE 16: 5つの戦略
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "カスタマイズ地獄を脱出する5つの戦略")

strategies = [
    ("01", "業界標準\nテンプレート", "国交省基準をベースに\n業界共通のマスタ・帳票を整備", GOLD_ACCENT),
    ("02", "ノーコード/\nローコード活用", "プログラミング不要で\n設定変更できる仕組みを最大活用", BLUE_ACCENT),
    ("03", "API連携の\n標準化", "外部システムとの\nデータ連携を標準インターフェースで", GREEN_ACCENT),
    ("04", "Fit & Gap\n分析の革新", "「できない」を探すのではなく\n「どう合わせるか」を提案", RGBColor(0x8E, 0x44, 0xAD)),
    ("05", "モバイル\n活用", "現場からの入力を\nスマホ・タブレットで最適化", ORANGE_ACCENT),
]

for i, (num, title, desc, color) in enumerate(strategies):
    x = Inches(0.2) + i * Inches(1.95)
    
    # Number circle
    add_oval(slide, x + Inches(0.6), Inches(1.1), Inches(0.6), Inches(0.6), color)
    add_text_box(slide, x + Inches(0.6), Inches(1.15), Inches(0.6), Inches(0.55),
                 num, 16, WHITE, True, PP_ALIGN.CENTER)
    
    # Card
    add_rounded_rect(slide, x, Inches(1.8), Inches(1.8), Inches(3.0), CREAM)
    
    # Title
    lines = title.split('\n')
    for li, line in enumerate(lines):
        add_text_box(slide, x + Emu(50000), Inches(1.9) + Emu(li * 200000), Inches(1.7), Inches(0.3),
                     line, 12, color, True, PP_ALIGN.CENTER)
    
    # Description
    desc_lines = desc.split('\n')
    for li, line in enumerate(desc_lines):
        add_text_box(slide, x + Emu(50000), Inches(2.6) + Emu(li * 180000), Inches(1.7), Inches(0.3),
                     line, 9, TEXT_GRAY, False, PP_ALIGN.CENTER)

add_footer(slide, 13)
add_notes(slide, """カスタマイズ地獄を脱出する5つの戦略です。
第1に業界標準テンプレート。国交省基準をベースに、業界共通のマスタや帳票を整備します。
第2にノーコード・ローコード活用。プログラミング不要で設定変更できる仕組みを最大限活用します。これにより、従来カスタマイズとして開発していた変更を、設定変更で対応できるようになります。
第3にAPI連携の標準化。外部システムとのデータ連携を標準インターフェースで実現します。
第4にFit&Gap分析の革新。「できないこと」を探すのではなく、「どう合わせるか」を提案する発想の転換です。これが最も重要かもしれません。
第5にモバイル活用。現場からの入力をスマホやタブレットで最適化します。
この5つの戦略を組み合わせることで、カスタマイズを大幅に削減できます。""")

# ============================================================
# SLIDE 17: Fit & Gap分析の新しい考え方
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "Fit & Gap 分析の新しい考え方")

# Traditional vs New approach
add_text_box(slide, Inches(0.3), Inches(1.0), Inches(4.3), Inches(0.4),
             "従来のアプローチ", 14, RED_ACCENT, True)
add_rounded_rect(slide, Inches(0.3), Inches(1.5), Inches(4.3), Inches(1.5), RGBColor(0xFF, 0xF0, 0xF0))

trad_steps = [
    ("1.", "標準機能を見せる"),
    ("2.", "「できない」ことを洗い出す"),
    ("3.", "「できないならカスタマイズ」"),
    ("4.", "費用が膨らむ..."),
]
for i, (num, text) in enumerate(trad_steps):
    y = Inches(1.6) + i * Inches(0.3)
    add_text_box(slide, Inches(0.5), y, Inches(4.0), Inches(0.3),
                 f"  {num} {text}", 11, TEXT_DARK if i < 3 else RED_ACCENT, i == 3)

add_text_box(slide, Inches(5.2), Inches(1.0), Inches(4.3), Inches(0.4),
             "新しいアプローチ", 14, GREEN_ACCENT, True)
add_rounded_rect(slide, Inches(5.2), Inches(1.5), Inches(4.3), Inches(1.5), RGBColor(0xF0, 0xFF, 0xF0))

new_steps = [
    ("1.", "業務プロセスを深く理解する"),
    ("2.", "「本当に必要か」を問い直す"),
    ("3.", "標準機能で「どう実現するか」提案"),
    ("4.", "業務の方を変える選択肢も提示"),
]
for i, (num, text) in enumerate(new_steps):
    y = Inches(1.6) + i * Inches(0.3)
    add_text_box(slide, Inches(5.4), y, Inches(4.0), Inches(0.3),
                 f"  {num} {text}", 11, TEXT_DARK if i < 3 else GREEN_ACCENT, i == 3)

# Key message
add_rounded_rect(slide, Inches(0.3), Inches(3.3), Inches(9.2), Inches(1.5), CREAM)
add_text_box(slide, Inches(0.5), Inches(3.4), Inches(8.8), Inches(0.4),
             "Fit & Gap 分析の3段階アプローチ", 13, GOLD, True, PP_ALIGN.CENTER)

gap_levels = [
    (Inches(0.5), "Level 1\n設定で対応", "標準の設定・パラメータ\nで解決できないか？", GREEN_ACCENT),
    (Inches(3.5), "Level 2\n運用で対応", "業務プロセスの見直しで\n標準に合わせられないか？", BLUE_ACCENT),
    (Inches(6.5), "Level 3\nアドオンで対応", "本当にカスタマイズが\n必要な最小限の範囲は？", ORANGE_ACCENT),
]

for x, title, desc, color in gap_levels:
    add_rounded_rect(slide, x, Inches(3.8), Inches(2.7), Inches(0.9), WHITE)
    lines = title.split('\n')
    add_text_box(slide, x + Emu(30000), Inches(3.85), Inches(1.0), Inches(0.3),
                 lines[0], 11, color, True)
    add_text_box(slide, x + Emu(30000), Inches(4.1), Inches(1.0), Inches(0.3),
                 lines[1], 11, color, True)
    desc_lines = desc.split('\n')
    for li, line in enumerate(desc_lines):
        add_text_box(slide, x + Inches(1.1), Inches(3.85) + Emu(li * 180000), Inches(1.5), Inches(0.3),
                     line, 9, TEXT_GRAY)

add_footer(slide, 14)
add_notes(slide, """Fit&Gap分析の新しい考え方について詳しく見ていきます。
左側が従来のアプローチ。標準機能を見せて、「できないこと」を洗い出し、「できないならカスタマイズ」となり、費用が膨らむ。この流れに心当たりがある方も多いのではないでしょうか。
右側が新しいアプローチ。まず業務プロセスを深く理解し、「本当に必要か」を問い直します。そして標準機能で「どう実現するか」を提案し、場合によっては業務の方を変える選択肢も提示します。
下のFit&Gap分析3段階アプローチが実践的なフレームワークです。
レベル1は設定で対応。標準のパラメータ設定で解決できないかをまず検討します。
レベル2は運用で対応。業務プロセスの見直しで標準に合わせられないかを考えます。
レベル3がアドオンで対応。本当にカスタマイズが必要な最小限の範囲だけを特定します。
この順番で検討することで、カスタマイズ量を劇的に減らせるのです。""")

# ============================================================
# SLIDE 18: カスタマイズ影響度マトリックス
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "カスタマイズ影響度マトリックス")

# Matrix table
matrix_data = [
    ["業務領域", "カスタマイズ\n発生頻度", "影響範囲", "標準化\n難易度", "優先度"],
    ["原価管理・工事採算", "非常に高い", "経営判断に直結", "高", "最優先"],
    ["契約・請求・入金", "高い", "取引先との関係", "中", "優先"],
    ["経費精算・現場管理", "中程度", "日常業務効率", "低", "早期着手"],
    ["購買・在庫管理", "中程度", "原価精度に影響", "中", "計画的"],
    ["ワークフロー", "高い", "全社的な影響", "低", "早期着手"],
    ["出面・労務管理", "高い", "法令遵守", "中", "優先"],
    ["個別帳票", "非常に高い", "情報可視化", "低", "BI活用"],
]
add_table_to_slide(slide, Inches(0.3), Inches(1.0), Inches(9.2), Inches(3.5), matrix_data,
                   [Inches(2.0), Inches(1.5), Inches(2.0), Inches(1.2), Inches(1.5)])

add_rounded_rect(slide, Inches(0.3), Inches(4.6), Inches(9.2), Inches(0.3), CREAM)
add_text_box(slide, Inches(0.5), Inches(4.62), Inches(9.0), Inches(0.25),
             "標準化難易度が「低」の領域から着手し、段階的にシステム全体を標準化していくのが効果的", 10, GOLD, True, PP_ALIGN.CENTER)
add_footer(slide, 15)
add_notes(slide, """カスタマイズ影響度マトリックスです。7つの業務領域を4つの軸で評価しています。
カスタマイズ発生頻度、影響範囲、標準化の難易度、そして優先度です。
原価管理・工事採算は発生頻度が非常に高く、経営判断に直結するため最優先です。ただし標準化難易度も高いので、段階的に取り組む必要があります。
契約・請求・入金と出面・労務管理は優先度が高い。法令遵守や取引先との関係に影響するためです。
経費精算・現場管理とワークフローは標準化難易度が低く、早期に着手すべき領域です。すぐに効果が出やすいところです。
個別帳票は発生頻度が非常に高いですが、BIツールの活用で解決できます。
下のメッセージの通り、標準化難易度が低い領域から着手し、段階的にシステム全体を標準化していくのが最も効果的なアプローチです。""")

# ============================================================
# SLIDE 19: まとめ ― カスタマイズを越えて
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, WHITE)
add_title_bar(slide, "まとめ：カスタマイズを越え、建設業の未来を創造する")

# 3 key takeaways
takeaways = [
    ("理解する", "建設業の「特殊性」は\n本当に特殊なのか？\n多くは業界共通の課題", GOLD_ACCENT),
    ("標準化する", "テンプレート・設定・\nBIツールの活用で\nカスタマイズを最小限に", GREEN_ACCENT),
    ("提案する", "「カスタマイズありき」\nではない新しい\nシステム導入の形を", BLUE_ACCENT),
]

for i, (title, desc, color) in enumerate(takeaways):
    x = Inches(0.3) + i * Inches(3.2)
    add_rounded_rect(slide, x, Inches(1.1), Inches(2.9), Inches(2.5), CREAM)
    add_oval(slide, x + Inches(1.05), Inches(1.2), Inches(0.7), Inches(0.7), color)
    add_text_box(slide, x + Inches(1.05), Inches(1.25), Inches(0.7), Inches(0.65),
                 str(i+1), 22, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Emu(50000), Inches(2.0), Inches(2.8), Inches(0.4),
                 title, 16, color, True, PP_ALIGN.CENTER)
    desc_lines = desc.split('\n')
    for li, line in enumerate(desc_lines):
        add_text_box(slide, x + Emu(50000), Inches(2.5) + Emu(li * 200000), Inches(2.8), Inches(0.3),
                     line, 11, TEXT_GRAY, False, PP_ALIGN.CENTER)

# Bottom message
add_rounded_rect(slide, Inches(0.3), Inches(3.8), Inches(9.2), Inches(1.0), GOLD)
add_multiline_text(slide, Inches(0.5), Inches(3.85), Inches(8.8), Inches(0.9), [
    ("「うちは特殊だから」を「標準化で解決できる」に変える", 16, WHITE, True, PP_ALIGN.CENTER),
    ("これが、ITコンサルタントの真の価値です", 12, GOLD_LIGHT, False, PP_ALIGN.CENTER),
])
add_footer(slide, 16)
add_notes(slide, """最後のまとめです。3つのキーメッセージを確認しましょう。
1つ目は「理解する」。建設業の特殊性は本当に特殊なのか。今日見てきた通り、多くは業界共通の課題です。自社だけの問題だと思っていたことが、実は他の建設会社でも同じように起きているのです。
2つ目は「標準化する」。テンプレート、設定変更、BIツールの活用で、カスタマイズを最小限に抑えられます。設定で対応、運用で対応、最後の手段としてアドオンで対応。この3段階の考え方を忘れないでください。
3つ目は「提案する」。カスタマイズありきではない、新しいシステム導入の形を提案していきましょう。業務を変える選択肢も含めた、より良い解決策を提示することが大切です。
「うちは特殊だから」を「標準化で解決できる」に変える。これがITコンサルタントの真の価値です。""")

# ============================================================
# SLIDE 20: End Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Emu(0), Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT, CREAM)
add_shape(slide, Emu(0), Emu(0), Emu(3200000), SLIDE_HEIGHT, GOLD_ACCENT)
add_shape(slide, Emu(3200000), Emu(0), Emu(100000), SLIDE_HEIGHT, GOLD_LIGHT)

add_text_box(slide, Inches(4.0), Inches(1.5), Inches(5.5), Inches(1.0),
             "ご清聴ありがとうございました", 28, TEXT_BROWN, True)

add_multiline_text(slide, Inches(4.0), Inches(2.8), Inches(3), Inches(0.6), [
    ("H A R M O N I C", 14, TEXT_BROWN),
])
add_multiline_text(slide, Inches(6.0), Inches(2.8), Inches(3), Inches(0.6), [
    ("i n s i g h t", 14, GOLD_LIGHT),
])

add_text_box(slide, Inches(4.0), Inches(3.5), Inches(5.5), Inches(0.5),
             "建設業DXの推進を、共に。", 14, BROWN_LIGHT)

add_text_box(slide, Inches(4.0), Inches(4.3), Inches(5), Inches(0.3),
             "Harmonic Insight 2026年3月12日", 10, BROWN_LIGHT)
add_notes(slide, """以上で本日の勉強会は終了です。ご清聴ありがとうございました。
今日お伝えしたかったのは、建設業のシステムリプレースにおいて、カスタマイズは避けられないものではなく、戦略的に減らせるものだということです。
明日からの業務で、「これは本当にカスタマイズが必要だろうか」「標準機能で代替できないだろうか」と一度立ち止まって考えていただければ幸いです。
ご質問やご相談がありましたら、お気軽にお申し付けください。
Harmonic Insightは、建設業DXの推進を、皆さんと共に歩んでまいります。
本日はありがとうございました。""")

# ===== Save =====
output_path = "建設業基幹システム_カスタマイズ多発ポイント_勉強会資料.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
