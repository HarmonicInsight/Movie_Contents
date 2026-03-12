#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
コンサルタントが知っておくべき未来年表 - 勉強会用パワーポイント生成
Harmonic Insight テンプレートスタイルに準拠
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ===== Color Palette =====
GOLD = RGBColor(0x8B, 0x75, 0x36)
GOLD_LIGHT = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_ACCENT = RGBColor(0xB8, 0x86, 0x0B)
DARK_BG = RGBColor(0x1A, 0x1A, 0x1A)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
CREAM2 = RGBColor(0xF5, 0xF0, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
TEXT_BROWN = RGBColor(0x3C, 0x2F, 0x1E)
BROWN_LIGHT = RGBColor(0x8B, 0x73, 0x55)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
BLUE_ACCENT = RGBColor(0x2C, 0x7A, 0xB0)
ORANGE_ACCENT = RGBColor(0xE6, 0x7E, 0x22)
PURPLE_ACCENT = RGBColor(0x8E, 0x44, 0xAD)
TEAL_ACCENT = RGBColor(0x16, 0xA0, 0x85)
TBL_HDR = RGBColor(0x8B, 0x75, 0x36)
TBL_R1 = RGBColor(0xF9, 0xF6, 0xF0)
TBL_R2 = RGBColor(0xF0, 0xE8, 0xD5)

SW = Emu(9144000)
SH = Emu(5143500)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH

# ===== Helpers =====
def rect(s, l, t, w, h, fc=None, lc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if fc:
        sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else:
        sh.fill.background()
    if lc:
        sh.line.color.rgb = lc; sh.line.width = Pt(1)
    return sh

def rrect(s, l, t, w, h, fc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if fc:
        sh.fill.solid(); sh.fill.fore_color.rgb = fc
    else:
        sh.fill.background()
    return sh

def oval(s, l, t, w, h, fc=GOLD):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.fill.background()
    return sh

def arrow_r(s, l, t, w, h, fc=GOLD):
    sh = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.fill.background()
    return sh

def arrow_d(s, l, t, w, h, fc=GOLD):
    sh = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.fill.background()
    return sh

def chevron(s, l, t, w, h, fc=GOLD):
    sh = s.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.fill.background()
    return sh

def tb(s, l, t, w, h, txt, sz=14, c=TEXT_DARK, b=False, al=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.alignment = al
    return bx

def ml(s, l, t, w, h, lines):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, li in enumerate(lines):
        txt = li[0]; sz = li[1] if len(li)>1 else 14; c = li[2] if len(li)>2 else TEXT_DARK
        bd = li[3] if len(li)>3 else False; al = li[4] if len(li)>4 else PP_ALIGN.LEFT
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = bd; p.alignment = al
    return bx

def footer(s, n):
    rect(s, Emu(0), Emu(4914900), SW, Emu(228600), GOLD)
    tb(s, Inches(0.3), Inches(5.2), Inches(3), Inches(0.3), "H A R M O N I C   i n s i g h t", 8, WHITE)
    tb(s, Inches(9.2), Inches(5.2), Inches(0.5), Inches(0.3), str(n), 8, WHITE)

def add_notes(s, text):
    notes = s.notes_slide
    tf = notes.notes_text_frame
    tf.text = text

def title_bar(s, txt):
    rect(s, Emu(0), Emu(0), SW, Emu(685800), GOLD)
    tb(s, Inches(0.5), Inches(0.08), Inches(9), Inches(0.6), txt, 22, WHITE, True)

def subtitle(s, txt):
    tb(s, Inches(0.5), Inches(0.95), Inches(9), Inches(0.4), txt, 13, GOLD, True)

def add_table(s, l, t, w, h, data, cw=None):
    nr = len(data); nc = len(data[0])
    ts = s.shapes.add_table(nr, nc, l, t, w, h)
    tbl = ts.table
    if cw:
        for i, cwidth in enumerate(cw):
            tbl.columns[i].width = cwidth
    for r, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(r, ci); cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                if r == 0:
                    p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
                else:
                    p.font.color.rgb = TEXT_DARK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_HDR
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_R2
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = TBL_R1
    return ts

def section_divider(num, title1, title2, desc):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, Emu(0), Emu(0), Emu(2926080), SH, GOLD)
    rect(s, Emu(2926080), Emu(0), Emu(6217920), SH, DARK_BG)
    tb(s, Inches(0.3), Inches(1.5), Inches(2.5), Inches(1.5), num, 60, WHITE, True, PP_ALIGN.CENTER)
    tb(s, Inches(3.5), Inches(1.2), Inches(6), Inches(0.8), title1, 28, WHITE)
    tb(s, Inches(3.5), Inches(2.0), Inches(6), Inches(0.8), title2, 34, GOLD_LIGHT, True)
    tb(s, Inches(3.5), Inches(3.2), Inches(6), Inches(1), desc, 14, RGBColor(0xAA,0xAA,0xAA))
    tb(s, Inches(3.5), Inches(4.5), Inches(3), Inches(0.3), "H A R M O N I C   i n s i g h t", 8, RGBColor(0x88,0x88,0x88))
    return s

# ============================================================
# SLIDE 1: タイトル
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, CREAM)
rect(s, Emu(0), Emu(0), Emu(3200000), SH, GOLD_ACCENT)
rect(s, Emu(3200000), Emu(0), Emu(100000), SH, GOLD_LIGHT)
tb(s, Inches(0.3), Inches(1.5), Inches(3.0), Inches(1.0), "Hi", 22, GOLD_LIGHT)
ml(s, Inches(3.8), Inches(0.6), Inches(5.8), Inches(3.0), [
    ("コンサルタントが", 32, TEXT_BROWN, True),
    ("知っておくべき", 32, TEXT_BROWN, True),
    ("未来年表 2025-2050", 34, GOLD, True),
    ("", 12, TEXT_DARK),
    ("～ 3年後・5年後のロードマップ策定のために ～", 16, BROWN_LIGHT),
])
ml(s, Inches(3.8), Inches(3.8), Inches(3), Inches(0.6), [("H A R M O N I C", 14, TEXT_BROWN)])
ml(s, Inches(6.0), Inches(3.8), Inches(3), Inches(0.6), [("i n s i g h t", 14, GOLD_LIGHT)])
tb(s, Inches(3.8), Inches(4.5), Inches(5), Inches(0.3), "Harmonic Insight 2026年3月12日", 10, BROWN_LIGHT)
add_notes(s, """皆さん、こんにちは。Harmonic Insightの勉強会にようこそお越しくださいました。
本日は「コンサルタントが知っておくべき未来年表 2025年から2050年」というテーマでお話しします。
クライアントへのロードマップ策定や中長期計画の提案に、この未来予測データをどう活かすかを一緒に学んでいきましょう。
約40分の勉強会です。どうぞ最後までお付き合いください。""")

# ============================================================
# SLIDE 2: 勉強会の目的
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "なぜコンサルタントに「未来年表」が必要なのか")

# Left: 3 reasons
reasons = [
    ("ロードマップ策定", "3年後・5年後の提案には\n未来の変化を織り込む必要がある", GOLD_ACCENT),
    ("クライアントへの説得力", "数値に基づく未来予測で\n提案の説得力が格段に上がる", BLUE_ACCENT),
    ("リスク回避", "想定外の変化に備え\n先手を打つ戦略が可能になる", RED_ACCENT),
]
for i, (t, d, c) in enumerate(reasons):
    y = Inches(1.2) + i * Inches(1.2)
    rrect(s, Inches(0.3), y, Inches(5.0), Inches(1.0), CREAM)
    oval(s, Inches(0.4), y + Emu(100000), Inches(0.45), Inches(0.45), c)
    tb(s, Inches(0.42), y + Emu(105000), Inches(0.45), Inches(0.4), str(i+1), 16, WHITE, True, PP_ALIGN.CENTER)
    tb(s, Inches(1.0), y + Emu(50000), Inches(2.0), Inches(0.3), t, 13, TEXT_DARK, True)
    for li, line in enumerate(d.split('\n')):
        tb(s, Inches(1.0), y + Emu(280000 + li*180000), Inches(4.0), Inches(0.3), line, 10, TEXT_GRAY)

# Right: Scope
rrect(s, Inches(5.6), Inches(1.2), Inches(4.0), Inches(3.5), CREAM)
tb(s, Inches(5.8), Inches(1.3), Inches(3.6), Inches(0.4), "本資料のカバー範囲", 13, GOLD, True, PP_ALIGN.CENTER)
scope = [
    "人口・社会構造の変化",
    "業界別AI代替率と年収変動",
    "GDP・国際競争力の推移",
    "地域格差の拡大予測",
    "世代別の影響と特性",
    "生活・働き方の変容",
    "環境・リスク要因",
    "コンサル実務への活用法",
]
for i, item in enumerate(scope):
    y = Inches(1.8) + i * Inches(0.32)
    tb(s, Inches(6.0), y, Inches(3.5), Inches(0.3), f"  {item}", 10, TEXT_DARK)
footer(s, 1)
add_notes(s, """まず、なぜコンサルタントに未来年表が必要なのかをお話しします。
3つの理由があります。
第一に、ロードマップ策定です。3年後・5年後の提案には、社会の変化を織り込む必要があります。
第二に、クライアントへの説得力です。数値に基づく未来予測があれば、提案の信頼性が格段に上がります。
第三に、リスク回避です。想定外の変化に備え、先手を打つ戦略が可能になります。
右側に本資料のカバー範囲を示しています。人口・社会構造から、業界別AI代替率、GDP推移、地域格差、世代別影響、生活の変容、環境リスク、そしてコンサル実務への活用法まで、幅広くカバーします。""")

# ============================================================
# SLIDE 3: 日本の全体像 2025→2050 ダッシュボード
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "日本の未来ダッシュボード 2025 → 2050")

# Key metrics cards
metrics = [
    ("人口", "1.23億人", "0.95億人", "-23%", RED_ACCENT),
    ("高齢化率", "28.6%", "39.0%", "+10pt", ORANGE_ACCENT),
    ("GDP世界順位", "3位", "8位", "↓5", RED_ACCENT),
    ("外国人比率", "3.5%", "12.8%", "+9pt", BLUE_ACCENT),
    ("AI技術力順位", "6位", "3位", "↑3", GREEN_ACCENT),
    ("労働生産性順位", "28位", "15位", "↑13", GREEN_ACCENT),
]

for i, (label, v25, v50, delta, c) in enumerate(metrics):
    col = i % 3; row = i // 3
    x = Inches(0.3) + col * Inches(3.2)
    y = Inches(1.0) + row * Inches(1.9)
    rrect(s, x, y, Inches(2.9), Inches(1.6), CREAM)
    tb(s, x + Emu(50000), y + Emu(30000), Inches(2.8), Inches(0.3), label, 11, TEXT_GRAY, False, PP_ALIGN.CENTER)
    # 2025
    tb(s, x + Emu(50000), y + Emu(260000), Inches(1.2), Inches(0.25), "2025年", 8, TEXT_GRAY, False, PP_ALIGN.CENTER)
    tb(s, x + Emu(50000), y + Emu(440000), Inches(1.2), Inches(0.4), v25, 16, TEXT_DARK, True, PP_ALIGN.CENTER)
    # Arrow
    arrow_r(s, x + Inches(1.2), y + Emu(500000), Inches(0.4), Inches(0.2), GOLD)
    # 2050
    tb(s, x + Emu(1300000), y + Emu(260000), Inches(1.2), Inches(0.25), "2050年", 8, TEXT_GRAY, False, PP_ALIGN.CENTER)
    tb(s, x + Emu(1300000), y + Emu(440000), Inches(1.2), Inches(0.4), v50, 16, c, True, PP_ALIGN.CENTER)
    # Delta badge
    rrect(s, x + Inches(1.0), y + Inches(1.15), Inches(0.9), Inches(0.3), c)
    tb(s, x + Inches(1.0), y + Inches(1.15), Inches(0.9), Inches(0.3), delta, 10, WHITE, True, PP_ALIGN.CENTER)

footer(s, 2)
add_notes(s, """こちらは日本の未来ダッシュボードです。2025年から2050年の25年間で、日本がどう変わるかを6つの主要指標で俯瞰します。
まず人口。1.23億人から0.95億人へ、23%も減少します。これは4人に1人がいなくなる計算です。
高齢化率は28.6%から39%へ。国民の約4割が65歳以上になります。
GDP世界順位は3位から8位へ下がります。ただし、これは他国の成長によるもので、日本が衰退するわけではありません。
外国人比率は3.5%から12.8%へ。人手不足を外国人労働者で補う形が進みます。
そして明るい材料もあります。AI技術力は6位から3位へ上昇、労働生産性も28位から15位へ大幅に改善します。
このダッシュボードは、クライアントへの提案の冒頭に使うと非常にインパクトがあります。""")

# ============================================================
# SLIDE 4: 社会変化マトリクス（2025-2050）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "社会変化マトリクス 2025-2050")
subtitle(s, "人口・経済・技術・生活の4軸で未来を俯瞰する")

tdata = [
    ["", "2025年", "2030年", "2040年", "2050年"],
    ["人口", "1.23億人\n高齢化率28.6%", "1.16億人\n高齢化率31.2%", "1.11億人\n高齢化率35.3%", "0.95億人\n高齢化率39.0%"],
    ["経済", "GDP 540兆円\n世界3位", "GDP 520兆円\n世界4位", "GDP 480兆円\n世界6位", "GDP 430兆円\n世界8位"],
    ["技術", "スマホ+AI\n業務支援AI開始", "ARグラス普及\nAGI実現", "体内センサー\n創造協働AI", "脳-PC接続\n意識融合AI"],
    ["働き方", "リモート標準化\n終身雇用30%", "専門職主流\nAI協働開始", "プロジェクト型\n実力主義社会", "創造=労働\nAI全業務支援"],
    ["平均年収", "420万円\n(実質400万円)", "480万円\n(実質410万円)", "530万円\n(実質430万円)", "580万円\n(実質450万円)"],
]
add_table(s, Inches(0.2), Inches(1.4), Inches(9.5), Inches(3.3), tdata,
          [Inches(1.2), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.0)])
rrect(s, Inches(0.3), Inches(4.75), Inches(9.2), Inches(0.2), CREAM)
tb(s, Inches(0.5), Inches(4.73), Inches(9.0), Inches(0.2),
   "平均年収は上昇するが、中央値はほぼ横ばい → 格差拡大が最大のテーマ", 9, RED_ACCENT, True, PP_ALIGN.CENTER)
footer(s, 3)
add_notes(s, """この社会変化マトリクスは、人口・経済・技術・働き方・平均年収の5軸で、2025年から2050年までを5年刻みで俯瞰しています。
人口は一貫して減少し、2050年には1億人を割ります。
経済面ではGDPの絶対額が540兆円から430兆円に減少。世界順位も3位から8位に後退します。
技術面では劇的な変化があります。2025年のスマホとAI業務支援から始まり、2030年にはARグラスが普及しAGIが実現。2040年には体内センサーや創造協働AIが登場し、2050年には脳とPCが直接接続される時代が来ると予測されています。
働き方も大きく変わります。終身雇用は2025年の30%からさらに減少し、プロジェクト型・実力主義の社会に移行します。
注目すべきは平均年収です。名目では420万円から580万円に上がりますが、実質ベースでは400万円から450万円。つまり、ほとんど変わらない。これが「格差拡大」という最大テーマにつながります。""")

# ============================================================
# SLIDE 5: Section - 業界別AI代替
# ============================================================
sd1 = section_divider("01", "業界別", "AI代替率と年収変動", "あなたのクライアントの業界は\nどのくらい変わるのか？")
add_notes(sd1, """ここからセクション1、業界別のAI代替率と年収変動に入ります。
皆さんのクライアントの業界が、AIによってどのくらい変わるのか。
これを正確に把握することが、的確な提案につながります。
それでは具体的な数字を見ていきましょう。""")

# ============================================================
# SLIDE 6: 業界別AI代替率 一覧
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "業界別AI代替率の推移 2025→2050")
subtitle(s, "業務時間ベースでAIが実行可能になる割合")

tdata2 = [
    ["業界", "2025年", "2030年", "2040年", "2050年", "中央値年収変動"],
    ["製造業", "30%", "45%", "75%", "90%", "-5% (格差拡大)"],
    ["金融・保険", "25%", "35%", "65%", "85%", "-21% (極端な二極化)"],
    ["小売・流通", "20%", "40%", "70%", "80%", "-13% (低賃金化)"],
    ["運輸・物流", "15%", "30%", "60%", "85%", "-8% (自動運転)"],
    ["建設・不動産", "20%", "25%", "50%", "65%", "+14% (人手不足)"],
    ["医療・介護", "10%", "15%", "25%", "40%", "+19% (人間性価値)"],
    ["教育", "5%", "20%", "45%", "60%", "+18% (創造性)"],
    ["農業", "25%", "40%", "70%", "80%", "+36% (食料安保)"],
    ["行政・公務", "5%", "15%", "35%", "55%", "+24% (専門性)"],
]
add_table(s, Inches(0.2), Inches(1.3), Inches(9.5), Inches(3.6), tdata2,
          [Inches(1.3), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(4.6)])
rrect(s, Inches(0.3), Inches(4.75), Inches(9.2), Inches(0.2), RGBColor(0xFF, 0xF0, 0xF0))
tb(s, Inches(0.5), Inches(4.73), Inches(9.0), Inches(0.2),
   "注意: 平均年収は上がっても中央値は下がる業界が多い → 「平均年収+100%の罠」に注意", 9, RED_ACCENT, True, PP_ALIGN.CENTER)
footer(s, 4)
add_notes(s, """業界別のAI代替率の推移を一覧で示しています。
製造業は2025年の30%から2050年には90%まで代替が進みます。金融・保険は85%、小売・流通は80%と、多くの業界で過半数以上の業務がAIに置き換わります。
一方、医療・介護は40%、教育は60%と、人間性や創造性が求められる分野は代替率が低めです。
ここで最も重要なのは右端の「中央値年収変動」の列です。
金融・保険はマイナス21%と極端な二極化。小売・流通はマイナス13%で低賃金化。
逆に建設・不動産はプラス14%、医療・介護はプラス19%、農業はプラス36%と、人手不足の業界は年収が上がります。
下のメモにある通り、「平均年収プラス100%の罠」に注意してください。平均は上がっても中央値は下がる業界が多いのです。""")

# ============================================================
# SLIDE 7: 平均 vs 中央値の罠（格差の可視化）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "「平均年収+100%」の罠 ― 格差拡大の現実")
subtitle(s, "平均は上がるのに、多くの人の実感は悪化する")

# Manufacturing example
rrect(s, Inches(0.3), Inches(1.4), Inches(4.5), Inches(3.4), RGBColor(0xFF, 0xF5, 0xE8))
tb(s, Inches(0.5), Inches(1.5), Inches(4.0), Inches(0.3), "製造業の年収分布（2050年予測）", 13, ORANGE_ACCENT, True)
tb(s, Inches(0.5), Inches(1.9), Inches(4.0), Inches(0.3), "平均 900万円 の内訳", 12, TEXT_DARK, True)

bars_mfg = [
    ("上位10%  AI設計者", "2,000-3,000万円", Inches(3.5), RGBColor(0xFF, 0xC0, 0x40)),
    ("上位30%  技術管理職", "800-1,200万円", Inches(2.2), RGBColor(0xFF, 0xD0, 0x60)),
    ("中央値   一般作業者", "380万円", Inches(1.0), RGBColor(0xFF, 0xA0, 0xA0)),
    ("下位30%  非正規", "250万円", Inches(0.7), RGBColor(0xFF, 0x80, 0x80)),
]
for i, (label, val, bw, bc) in enumerate(bars_mfg):
    y = Inches(2.3) + i * Inches(0.55)
    tb(s, Inches(0.5), y, Inches(2.5), Inches(0.25), label, 9, TEXT_DARK)
    rrect(s, Inches(0.5), y + Emu(170000), bw, Emu(120000), bc)
    tb(s, Inches(0.5) + bw + Emu(30000), y + Emu(150000), Inches(1.5), Inches(0.2), val, 9, TEXT_DARK, True)

# Finance example
rrect(s, Inches(5.2), Inches(1.4), Inches(4.3), Inches(3.4), RGBColor(0xE8, 0xF0, 0xFF))
tb(s, Inches(5.4), Inches(1.5), Inches(4.0), Inches(0.3), "金融業の年収分布（2050年予測）", 13, BLUE_ACCENT, True)
tb(s, Inches(5.4), Inches(1.9), Inches(4.0), Inches(0.3), "平均 1,100万円 の内訳", 12, TEXT_DARK, True)

bars_fin = [
    ("上位5%   AIアナリスト", "3,000万円以上", Inches(3.3), RGBColor(0x80, 0xB0, 0xFF)),
    ("上位30%  営業・企画", "600-1,000万円", Inches(1.8), RGBColor(0xA0, 0xC8, 0xFF)),
    ("中央値   一般事務", "380万円", Inches(0.9), RGBColor(0xFF, 0xA0, 0xA0)),
    ("下位30%  AI代替職", "200万円", Inches(0.5), RGBColor(0xFF, 0x80, 0x80)),
]
for i, (label, val, bw, bc) in enumerate(bars_fin):
    y = Inches(2.3) + i * Inches(0.55)
    tb(s, Inches(5.4), y, Inches(2.5), Inches(0.25), label, 9, TEXT_DARK)
    rrect(s, Inches(5.4), y + Emu(170000), bw, Emu(120000), bc)
    tb(s, Inches(5.4) + bw + Emu(30000), y + Emu(150000), Inches(1.5), Inches(0.2), val, 9, TEXT_DARK, True)

rrect(s, Inches(0.3), Inches(4.75), Inches(9.2), Inches(0.2), RGBColor(0xFF, 0xE8, 0xE8))
tb(s, Inches(0.5), Inches(4.73), Inches(9.0), Inches(0.2),
   "同じ業界内でも最大12倍の格差 ― クライアントの「どの層」にいるかで提案が変わる", 9, RED_ACCENT, True, PP_ALIGN.CENTER)
footer(s, 5)
add_notes(s, """先ほどの「平均年収プラス100%の罠」を具体的に見てみましょう。
左側は製造業の2050年予測です。平均年収は900万円と一見高く見えます。しかし、その内訳を見ると、上位10%のAI設計者が2,000万から3,000万円を稼ぐ一方、中央値の一般作業者は380万円。下位30%の非正規はわずか250万円です。
右側の金融業はさらに極端です。平均1,100万円ですが、上位5%のAIアナリストが3,000万円以上。中央値の一般事務は380万円。下位30%のAI代替職は200万円にまで落ちます。
同じ業界内で最大12倍の格差が生まれるということです。
コンサルタントとして重要なのは、クライアントの従業員が「どの層」にいるかを見極めることです。それによって提案内容が全く変わってきます。""")

# ============================================================
# SLIDE 8: 建設業 AI代替 詳細（コンサル特に重要）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "建設業 AI代替の詳細 ― 2050年の姿")
subtitle(s, "コンサルタントとして特に知っておくべき業界")

# 4 areas
areas = [
    ("設計業務", [("CAD作成", "80%"), ("構造計算", "90%"), ("法規チェック", "95%")], BLUE_ACCENT),
    ("施工業務", [("測量", "70%"), ("重機操作", "60%"), ("品質検査", "80%")], ORANGE_ACCENT),
    ("管理業務", [("工程管理", "90%"), ("資材発注", "85%"), ("安全監視", "70%")], TEAL_ACCENT),
    ("残存業務", [("創造的設計", "-"), ("現場判断", "-"), ("職人技", "-")], GREEN_ACCENT),
]
for i, (title, items, c) in enumerate(areas):
    x = Inches(0.2) + i * Inches(2.45)
    rrect(s, x, Inches(1.4), Inches(2.3), Inches(2.6), CREAM)
    rrect(s, x, Inches(1.4), Inches(2.3), Inches(0.45), c)
    tb(s, x, Inches(1.42), Inches(2.3), Inches(0.4), title, 13, WHITE, True, PP_ALIGN.CENTER)
    for j, (item, pct) in enumerate(items):
        y = Inches(2.0) + j * Inches(0.6)
        tb(s, x + Emu(80000), y, Inches(1.3), Inches(0.25), item, 11, TEXT_DARK)
        if pct != "-":
            # Bar
            bw_pct = int(pct.replace('%','')) / 100
            rrect(s, x + Emu(80000), y + Emu(200000), Emu(int(1800000 * bw_pct)), Emu(100000), c)
            tb(s, x + Emu(80000) + Emu(int(1800000 * bw_pct)) + Emu(20000), y + Emu(170000), Inches(0.5), Inches(0.2), pct, 9, c, True)
        else:
            tb(s, x + Emu(80000), y + Emu(200000), Inches(2.0), Inches(0.2), "人間が不可欠", 9, GREEN_ACCENT, True)

# Bottom message
rrect(s, Inches(0.3), Inches(4.2), Inches(9.2), Inches(0.7), CREAM)
ml(s, Inches(0.5), Inches(4.25), Inches(8.8), Inches(0.6), [
    ("年収予測: 平均720万円 / 中央値480万円（+14%）― 人手不足で待遇は改善傾向", 10, TEXT_DARK, True),
    ("スマート建築設計者: 年収1,200万円  /  現場職人も技能向上で緩やかに改善", 10, TEXT_GRAY),
])
footer(s, 6)
add_notes(s, """建設業のAI代替を詳しく見ていきます。コンサルタントの皆さんにとって建設業は主要なクライアントセクターですので、しっかり押さえましょう。
設計業務では、CAD作成が80%、構造計算が90%、法規チェックが95%とほぼ完全にAI化されます。
施工業務は、測量が70%、重機操作が60%、品質検査が80%です。
管理業務も工程管理90%、資材発注85%と高い代替率です。
一方、残存業務として「創造的設計」「現場判断」「職人技」は人間が不可欠な領域として残ります。
年収予測は平均720万円、中央値480万円でプラス14%。人手不足により待遇は改善傾向にあります。
特にスマート建築設計者は年収1,200万円と高額です。建設DXの提案では、この「残存業務」と「高付加価値人材」に焦点を当てると効果的です。""")

# ============================================================
# SLIDE 9: Section - 地域・世代格差
# ============================================================
sd2 = section_divider("02", "地域格差と", "世代別の未来", "コンサルの提案先が「どこ」で「誰」なのかで\n未来の見え方が全く異なる")
add_notes(sd2, """セクション2では、地域格差と世代別の未来を見ていきます。
コンサルタントの提案先が「どこ」にあって「誰」が意思決定者なのかによって、未来の見え方は全く異なります。
同じDX提案でも、東京の大企業と地方の中小企業では、まったく違うアプローチが必要です。""")

# ============================================================
# SLIDE 10: 地域格差
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "地域格差の拡大予測 2025→2050")
subtitle(s, "東京圏と地方圏の格差は2.7倍に拡大")

tdata3 = [
    ["地域", "2025年\n人口", "2050年\n人口", "変化率", "2050年\n平均年収", "AI\n活用度"],
    ["東京圏", "3,700万人\n(30.1%)", "2,100万人\n(22.1%)", "-43%", "1,200万円", "95%"],
    ["大阪圏", "1,800万人\n(14.6%)", "1,200万人\n(12.6%)", "-33%", "900万円", "85%"],
    ["名古屋圏", "1,100万人\n(8.9%)", "800万人\n(8.4%)", "-27%", "950万円", "80%"],
    ["地方中核都市", "2,400万人\n(19.5%)", "1,900万人\n(20.0%)", "-21%", "700万円", "70%"],
    ["地方圏", "4,000万人\n(32.5%)", "2,950万人\n(31.1%)", "-26%", "450万円", "40%"],
]
add_table(s, Inches(0.2), Inches(1.3), Inches(9.5), Inches(2.6), tdata3,
          [Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.0), Inches(1.5), Inches(1.0)])

# Key insight
rrect(s, Inches(0.3), Inches(4.1), Inches(4.5), Inches(0.8), RGBColor(0xE8, 0xF0, 0xFF))
tb(s, Inches(0.5), Inches(4.15), Inches(4.2), Inches(0.3), "コンサルへの示唆", 11, BLUE_ACCENT, True)
tb(s, Inches(0.5), Inches(4.45), Inches(4.2), Inches(0.3),
   "クライアントの所在地域でDX提案内容が異なる", 10, TEXT_DARK)

rrect(s, Inches(5.2), Inches(4.1), Inches(4.3), Inches(0.8), RGBColor(0xE8, 0xF8, 0xE8))
tb(s, Inches(5.4), Inches(4.15), Inches(4.0), Inches(0.3), "提案のポイント", 11, GREEN_ACCENT, True)
tb(s, Inches(5.4), Inches(4.45), Inches(4.0), Inches(0.3),
   "地方こそAI活用度の伸びしろが大きい（40%→潜在的80%）", 10, TEXT_DARK)
footer(s, 7)
add_notes(s, """地域格差の拡大予測です。東京圏と地方圏の格差は2.7倍に広がると予測されています。
東京圏は人口が43%減少しますが、平均年収は1,200万円でAI活用度は95%。
一方、地方圏は26%の人口減少で、平均年収は450万円、AI活用度はわずか40%です。
左下のコンサルへの示唆にある通り、クライアントの所在地域によってDX提案の内容を変える必要があります。
そして右下の提案のポイントが重要です。地方こそAI活用度の伸びしろが大きい。現状40%のところを潜在的に80%まで引き上げられる可能性があります。
つまり、地方のクライアントへのAI・DX提案は、むしろビジネスチャンスが大きいと言えます。""")

# ============================================================
# SLIDE 11: 世代別影響
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "世代別 人口・影響力の推移")
subtitle(s, "クライアントの意思決定者は「どの世代」か？")

gen_data = [
    ["世代", "2025年\n人口", "2030年\n影響力", "2050年\n影響力", "主な特徴"],
    ["α世代\n(2020年～)", "500万人", "教育変革", "社会主導", "AIネイティブ\n創造性重視"],
    ["Z世代\n(1997-2019年)", "2,000万人", "社会中核", "指導・伝承", "デジタル変革\n価値観多様化"],
    ["ミレニアル\n(1981-1996年)", "2,200万人", "経営層", "助言・サポート", "ワークライフバランス\n実用主義"],
    ["X世代\n(1965-1980年)", "2,100万人", "引退準備", "知恵継承", "現実主義\n安定志向"],
    ["ベビーブーマー\n(1946-1964年)", "1,800万人", "後期高齢", "100歳現役", "豊富な経験\n伝統価値"],
]
add_table(s, Inches(0.2), Inches(1.3), Inches(9.5), Inches(3.2), gen_data,
          [Inches(1.6), Inches(1.2), Inches(1.5), Inches(1.5), Inches(3.0)])

rrect(s, Inches(0.3), Inches(4.6), Inches(9.2), Inches(0.35), CREAM)
tb(s, Inches(0.5), Inches(4.62), Inches(9.0), Inches(0.3),
   "2030年の意思決定者はミレニアル世代 → ワークライフバランスとデジタルを重視した提案が効果的",
   9, GOLD, True, PP_ALIGN.CENTER)
footer(s, 8)
add_notes(s, """世代別の人口と影響力の推移を見ていきます。
α世代は2020年以降生まれのAIネイティブ世代。2050年には社会を主導する世代になります。
Z世代は現在社会の中核であり、デジタル変革と価値観の多様化を牽引しています。
ミレニアル世代は2030年には経営層の主力。ワークライフバランスと実用主義が特徴です。
X世代は引退準備に入りますが、豊富な知恵を次世代に継承する役割を担います。
ベビーブーマーは後期高齢期ですが、100歳現役時代を迎え、経験と伝統の価値を守ります。
下のメッセージが実務のポイントです。2030年の意思決定者はミレニアル世代です。ワークライフバランスとデジタルを重視した提案が最も効果的です。""")

# ============================================================
# SLIDE 12: Section - 生活と社会の変容
# ============================================================
sd3 = section_divider("03", "生活と社会の", "大変容", "食・住・医療・教育…\n生活のあらゆる領域が変わる")
add_notes(sd3, """セクション3では、生活と社会の大変容について見ていきます。
食・住・医療・教育など、生活のあらゆる領域が劇的に変わります。
これらの変化は、コンサルタントが企業への提案だけでなく、個人のキャリア相談や政策提言にも活用できる重要な情報です。""")

# ============================================================
# SLIDE 13: カテゴリ別変化サマリー
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "生活カテゴリ別 変化サマリー")

cats = [
    ("食生活", "食費25%→階級社会\n米価格倍増→食料安保", "食費格差3倍に拡大\n水取引市場が確立", RGBColor(0xE8, 0x6C, 0x00)),
    ("住環境", "空き家800万戸\n都市部家賃高騰", "所有概念が消失\n住居格差10倍に", BLUE_ACCENT),
    ("働き方", "リモート標準化\n終身雇用30%", "創造=労働\n完全実力主義社会", PURPLE_ACCENT),
    ("健康・医療", "医療費48兆円\n医師不足", "平均寿命90歳\nAI診断100%", GREEN_ACCENT),
    ("教育", "大学機能低下\nリスキリング需要", "教育階級社会完成\n脳直接学習", ORANGE_ACCENT),
    ("お金・資産", "NISA拡充\n格差5倍", "格差社会固定化\n個人価値直接通貨化", RED_ACCENT),
]

for i, (label, now, future, c) in enumerate(cats):
    col = i % 3; row = i // 3
    x = Inches(0.2) + col * Inches(3.25)
    y = Inches(1.1) + row * Inches(2.1)
    rrect(s, x, y, Inches(3.0), Inches(1.8), CREAM)
    rrect(s, x, y, Inches(3.0), Inches(0.4), c)
    tb(s, x, y + Emu(20000), Inches(3.0), Inches(0.35), label, 12, WHITE, True, PP_ALIGN.CENTER)
    # Now
    tb(s, x + Emu(50000), y + Emu(400000), Inches(0.6), Inches(0.2), "現在", 8, c, True)
    for li, line in enumerate(now.split('\n')):
        tb(s, x + Emu(50000), y + Emu(560000 + li*160000), Inches(2.8), Inches(0.2), line, 9, TEXT_GRAY)
    # Future
    tb(s, x + Emu(50000), y + Emu(950000), Inches(0.6), Inches(0.2), "2050", 8, c, True)
    for li, line in enumerate(future.split('\n')):
        tb(s, x + Emu(50000), y + Emu(1110000 + li*160000), Inches(2.8), Inches(0.2), line, 9, TEXT_DARK, True)

footer(s, 9)
add_notes(s, """6つの生活カテゴリ別に、現在と2050年の変化をまとめています。
食生活では、食費が階級社会化し格差が3倍に拡大。米の価格は倍増し、水の取引市場まで確立されます。
住環境は空き家800万戸から所有概念の消失へ。住居格差は10倍に広がります。
働き方は終身雇用30%から完全実力主義社会へ。創造性こそが労働価値になります。
健康・医療は平均寿命90歳、AI診断100%の時代。医療費は48兆円から大幅に増加します。
教育は大学機能の低下から教育階級社会の完成へ。最終的には脳への直接学習も研究が進みます。
お金・資産では格差5倍から固定化へ。個人の価値が直接通貨化される世界が来ます。
コンサルタントとしては、クライアント企業の従業員がこうした変化にどう対応すべきか、包括的な提案ができるようになります。""")

# ============================================================
# SLIDE 14: GDP・国際競争力
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "GDP・国際競争力の推移")
subtitle(s, "日本経済の相対的地位は低下するが、技術力は向上する")

# GDP table
gdp_data = [
    ["指標", "2025年", "2030年", "2040年", "2050年"],
    ["名目GDP", "540兆円", "520兆円", "480兆円", "430兆円"],
    ["世界順位", "3位", "4位", "6位", "8位"],
    ["1人当たりGDP", "430万円", "450万円", "430万円", "450万円"],
    ["AI・技術競争力", "世界6位", "世界4位", "-", "世界3位"],
    ["労働生産性", "世界28位", "世界22位", "-", "世界15位"],
]
add_table(s, Inches(0.2), Inches(1.3), Inches(5.5), Inches(2.5), gdp_data,
          [Inches(1.5), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0)])

# Right: Key factors
rrect(s, Inches(5.9), Inches(1.3), Inches(3.7), Inches(2.5), CREAM)
tb(s, Inches(6.1), Inches(1.4), Inches(3.4), Inches(0.3), "GDP減少の主要因", 12, RED_ACCENT, True)
factors = [
    ("人口減少", "23%減 → 労働力・消費が縮小"),
    ("高齢化", "AI恩恵は限定的"),
    ("内需縮小", "消費・投資の構造的減少"),
    ("他国の成長", "インド・東南アジアの台頭"),
]
for i, (t, d) in enumerate(factors):
    y = Inches(1.8) + i * Inches(0.5)
    tb(s, Inches(6.1), y, Inches(1.3), Inches(0.25), t, 10, TEXT_DARK, True)
    tb(s, Inches(6.1), y + Emu(170000), Inches(3.3), Inches(0.25), d, 9, TEXT_GRAY)

# Bottom: Positive message
rrect(s, Inches(0.3), Inches(4.1), Inches(9.2), Inches(0.8), RGBColor(0xE8, 0xF8, 0xE8))
tb(s, Inches(0.5), Inches(4.15), Inches(9.0), Inches(0.3),
   "明るい材料: AI・技術競争力は米中に次ぐ世界3位へ上昇", 12, GREEN_ACCENT, True)
tb(s, Inches(0.5), Inches(4.5), Inches(9.0), Inches(0.3),
   "GDP規模は縮小するが、1人当たり生産性・技術力では世界トップクラスを維持 → DX推進が鍵", 10, TEXT_DARK)
footer(s, 10)
add_notes(s, """GDP・国際競争力の推移を見ていきます。
名目GDPは540兆円から430兆円に減少し、世界順位は3位から8位に後退します。
右側にGDP減少の主要因を4つ挙げています。人口減少23%による労働力と消費の縮小、高齢化によるAI恩恵の限定、内需の構造的な縮小、そしてインドや東南アジアなど他国の成長です。
ただし、1人当たりGDPは430万円から450万円とほぼ横ばいで、国民一人ひとりの生産性は維持されます。
そして明るい材料があります。AI・技術競争力は米中に次ぐ世界3位に上昇。労働生産性も28位から15位へ大幅に改善されます。
つまり、GDP規模は縮小しても、技術力と生産性では世界トップクラスを維持できる。この「DX推進が鍵」というメッセージが、コンサルタントの提案の核になります。""")

# ============================================================
# SLIDE 15: 環境・リスク要因
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "環境制約・リスク要因")
subtitle(s, "気候変動・水資源・ブラックスワンへの備え")

# Climate impact
rrect(s, Inches(0.3), Inches(1.4), Inches(4.5), Inches(1.8), RGBColor(0xFF, 0xF5, 0xE8))
tb(s, Inches(0.5), Inches(1.5), Inches(4.0), Inches(0.3), "気候変動の経済影響", 12, ORANGE_ACCENT, True)
climate = [
    ["項目", "2030年", "2050年"],
    ["自然災害", "GDP-1.2%", "GDP-2.8%"],
    ["農業生産", "-15%", "-30%"],
    ["エネルギーコスト", "+20%", "+40%"],
    ["適応コスト", "年3兆円", "年8兆円"],
]
add_table(s, Inches(0.4), Inches(1.9), Inches(4.2), Inches(1.2), climate,
          [Inches(1.5), Inches(1.3), Inches(1.4)])

# Black swan
rrect(s, Inches(5.2), Inches(1.4), Inches(4.3), Inches(1.8), RGBColor(0xFF, 0xF0, 0xF0))
tb(s, Inches(5.4), Inches(1.5), Inches(4.0), Inches(0.3), "ブラックスワン・リスク", 12, RED_ACCENT, True)
risks = [
    ["リスク", "確率", "影響度"],
    ["パンデミック再発", "30%", "高"],
    ["大規模自然災害", "60%", "高"],
    ["サイバー攻撃", "70%", "中"],
    ["地政学的緊張", "40%", "高"],
    ["AI制御不能", "10%", "極高"],
]
add_table(s, Inches(5.3), Inches(1.9), Inches(4.1), Inches(1.2), risks,
          [Inches(1.5), Inches(1.0), Inches(1.0)])

# Bottom: implication
rrect(s, Inches(0.3), Inches(3.5), Inches(9.2), Inches(1.4), CREAM)
tb(s, Inches(0.5), Inches(3.6), Inches(9.0), Inches(0.3),
   "コンサルタントとしての活用ポイント", 13, GOLD, True)
ml(s, Inches(0.5), Inches(4.0), Inches(8.8), Inches(0.8), [
    ("ロードマップ策定時に「環境制約」をシナリオに組み込む ― 水資源・エネルギーコストの上昇は全業界に影響", 10, TEXT_DARK),
    ("サイバー攻撃リスク70% → セキュリティ投資は「コスト」ではなく「必須投資」として提案する", 10, TEXT_DARK),
    ("ブラックスワン事象を想定した「レジリエンス設計」をロードマップに標準で盛り込む", 10, TEXT_DARK),
])
footer(s, 11)
add_notes(s, """環境制約とリスク要因について見ていきます。ロードマップ策定には、これらのリスクも織り込む必要があります。
左側の気候変動の経済影響では、2050年に自然災害がGDPのマイナス2.8%、農業生産がマイナス30%、エネルギーコストがプラス40%と深刻な影響が予測されています。適応コストだけで年間8兆円が必要です。
右側のブラックスワン・リスクでは、大規模自然災害の発生確率が60%、サイバー攻撃は70%と高い確率です。
下のコンサルタントとしての活用ポイントをご覧ください。
まず、ロードマップ策定時に環境制約をシナリオに組み込むこと。水資源やエネルギーコストの上昇は全業界に影響します。
次に、サイバー攻撃リスク70%を踏まえ、セキュリティ投資はコストではなく必須投資として提案すること。
そして、ブラックスワン事象を想定したレジリエンス設計をロードマップに標準で盛り込むことです。""")

# ============================================================
# SLIDE 16: Section - コンサル実務への活用
# ============================================================
sd4 = section_divider("04", "コンサル実務への", "活用法", "この未来年表を「3年後・5年後の\nロードマップ」にどう落とし込むか")
add_notes(sd4, """いよいよ最後のセクション、コンサル実務への活用法に入ります。
ここまで見てきた未来年表のデータを、具体的に3年後・5年後のロードマップにどう落とし込むか。
実践的なフレームワークと活用例をお伝えします。""")

# ============================================================
# SLIDE 17: 個人のキャリアシナリオ
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "未来年表で見る 個人のキャリアシナリオ")
subtitle(s, "「今30歳の人」と「今20歳の人」で未来がこれだけ違う")

# 30 years old
rrect(s, Inches(0.3), Inches(1.4), Inches(4.5), Inches(2.8), RGBColor(0xE8, 0xF0, 0xFF))
tb(s, Inches(0.5), Inches(1.5), Inches(4.0), Inches(0.3), "現在30歳（1995年生まれ）", 13, BLUE_ACCENT, True)

career30 = [
    ("30歳 2025年", "中間管理職\nリスキリング必須", "年収500万円"),
    ("35歳 2030年", "AI協働\n専門性特化", "年収650万円"),
    ("45歳 2040年", "創造性重視\nプロジェクト型", "年収800万円"),
    ("55歳 2050年", "知識継承\nコンサル型", "年収750万円"),
]
for i, (age, desc, income) in enumerate(career30):
    y = Inches(1.9) + i * Inches(0.55)
    chevron(s, Inches(0.4), y, Inches(1.3), Inches(0.4), BLUE_ACCENT if i < 3 else RGBColor(0xA0, 0xC0, 0xE0))
    tb(s, Inches(0.55), y + Emu(30000), Inches(1.1), Inches(0.35), age, 9, WHITE, True, PP_ALIGN.CENTER)
    tb(s, Inches(1.8), y + Emu(20000), Inches(2.0), Inches(0.35), desc.replace('\n', ' / '), 9, TEXT_DARK)
    tb(s, Inches(3.8), y + Emu(20000), Inches(0.9), Inches(0.35), income, 9, BLUE_ACCENT, True)

# 20 years old
rrect(s, Inches(5.2), Inches(1.4), Inches(4.3), Inches(2.8), RGBColor(0xF0, 0xE8, 0xFF))
tb(s, Inches(5.4), Inches(1.5), Inches(4.0), Inches(0.3), "現在20歳（α世代）", 13, PURPLE_ACCENT, True)

career20 = [
    ("20歳 2025年", "AIネイティブ学習\n大学変革期", "デジタル完全統合"),
    ("25歳 2030年", "AI協創職\n高度専門化", "創造性重視"),
    ("35歳 2040年", "社会中枢\nリーダー層", "多様性・持続可能"),
    ("45歳 2050年", "社会変革主導\n新文化創造", "ポスト・シンギュラリティ"),
]
for i, (age, desc, trait) in enumerate(career20):
    y = Inches(1.9) + i * Inches(0.55)
    chevron(s, Inches(5.3), y, Inches(1.3), Inches(0.4), PURPLE_ACCENT if i < 3 else RGBColor(0xC0, 0xA0, 0xE0))
    tb(s, Inches(5.45), y + Emu(30000), Inches(1.1), Inches(0.35), age, 9, WHITE, True, PP_ALIGN.CENTER)
    tb(s, Inches(6.7), y + Emu(20000), Inches(1.8), Inches(0.35), desc.replace('\n', ' / '), 9, TEXT_DARK)
    tb(s, Inches(8.5), y + Emu(20000), Inches(1.0), Inches(0.35), trait, 8, PURPLE_ACCENT, True)

rrect(s, Inches(0.3), Inches(4.4), Inches(9.2), Inches(0.5), CREAM)
tb(s, Inches(0.5), Inches(4.42), Inches(9.0), Inches(0.45),
   "クライアントの従業員の年齢層に応じた「キャリアシナリオ」を提示できると、提案の説得力が段違いに上がる",
   10, GOLD, True, PP_ALIGN.CENTER)
footer(s, 12)
add_notes(s, """未来年表を個人のキャリアシナリオに落とし込んだ例です。
左側は現在30歳、1995年生まれの方のシナリオ。30歳で中間管理職・年収500万円からスタートし、35歳でAI協働・専門性特化で650万円、45歳で創造性重視のプロジェクト型で800万円、55歳で知識継承・コンサル型で750万円という道筋です。
右側は現在20歳、α世代の方。AIネイティブとして学び、25歳でAI協創職、35歳で社会中枢のリーダー層、45歳で社会変革を主導するポスト・シンギュラリティ世代です。
下のメッセージが大切です。クライアントの従業員の年齢層に応じたキャリアシナリオを提示できると、提案の説得力が段違いに上がります。
例えば、人材育成プログラムの提案では、30代と20代で全く異なるカリキュラムが必要だということを、このデータで示せるわけです。""")

# ============================================================
# SLIDE 18: ロードマップ策定への活用フレームワーク
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "ロードマップ策定への活用フレームワーク")
subtitle(s, "未来年表を「3年後・5年後」の具体的な計画に落とし込む")

# 3-step framework
steps = [
    ("Step 1\n現状把握", "クライアントの業界・地域・\n世代構成を確認", "未来年表の「どの列」を\n読めばいいか特定する", GOLD_ACCENT),
    ("Step 2\nギャップ分析", "現状と2030年/2050年の\n差分を明確化", "AI代替率・年収変動・\n人口変化から影響度を算出", BLUE_ACCENT),
    ("Step 3\nロードマップ化", "3年後・5年後の\nマイルストーン設定", "環境リスクも織り込んだ\nシナリオプランニング", GREEN_ACCENT),
]

for i, (title, desc, detail, c) in enumerate(steps):
    x = Inches(0.2) + i * Inches(3.3)
    rrect(s, x, Inches(1.4), Inches(3.0), Inches(3.2), CREAM)
    rrect(s, x, Inches(1.4), Inches(3.0), Inches(0.8), c)
    lines = title.split('\n')
    for li, line in enumerate(lines):
        tb(s, x, Inches(1.45) + Emu(li * 200000), Inches(3.0), Inches(0.3), line, 13 if li == 0 else 11, WHITE, True, PP_ALIGN.CENTER)
    for li, line in enumerate(desc.split('\n')):
        tb(s, x + Emu(80000), Inches(2.4) + Emu(li * 180000), Inches(2.8), Inches(0.25), line, 10, TEXT_DARK)
    for li, line in enumerate(detail.split('\n')):
        tb(s, x + Emu(80000), Inches(3.1) + Emu(li * 180000), Inches(2.8), Inches(0.25), line, 9, TEXT_GRAY)
    if i < 2:
        arrow_r(s, x + Inches(2.95), Inches(2.8), Inches(0.35), Inches(0.25), GOLD)

rrect(s, Inches(0.3), Inches(4.7), Inches(9.2), Inches(0.25), GOLD)
tb(s, Inches(0.5), Inches(4.7), Inches(9.0), Inches(0.25),
   "未来年表 × Fit & Gap分析 = 説得力のあるロードマップ", 11, WHITE, True, PP_ALIGN.CENTER)
footer(s, 13)
add_notes(s, """未来年表をロードマップに落とし込む3ステップのフレームワークです。
ステップ1は現状把握。クライアントの業界・地域・世代構成を確認し、未来年表のどの部分を読めばいいか特定します。
ステップ2はギャップ分析。現状と2030年・2050年の差分を明確にし、AI代替率・年収変動・人口変化から影響度を算出します。
ステップ3はロードマップ化。3年後・5年後のマイルストーンを設定し、環境リスクも織り込んだシナリオプランニングを行います。
下のメッセージの通り、「未来年表 かける Fit&Gap分析 イコール 説得力のあるロードマップ」です。
このフレームワークを使えば、どのクライアントに対しても構造的に未来を見据えた提案ができるようになります。""")

# ============================================================
# SLIDE 19: 活用例 ― 3つの戦略視点
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "未来年表の3つの活用視点")

views = [
    ("企業戦略", [
        "人材戦略: 世代別特性に応じた採用・育成",
        "事業戦略: AI代替率を考慮した事業転換",
        "立地戦略: 地域格差を踏まえた拠点配置",
    ], GOLD_ACCENT),
    ("個人キャリア", [
        "スキル: AI代替されにくい創造性・対人力",
        "居住: 地域格差を踏まえた住まい選択",
        "資産: インフレ・技術変化を考慮した投資",
    ], BLUE_ACCENT),
    ("政策提言", [
        "労働: AI代替に伴う雇用対策・再教育",
        "地域: 地域格差是正・持続可能な地域創造",
        "社会保障: 高齢化・人口減への制度設計",
    ], GREEN_ACCENT),
]

for i, (title, items, c) in enumerate(views):
    x = Inches(0.2) + i * Inches(3.3)
    rrect(s, x, Inches(1.1), Inches(3.0), Inches(3.2), CREAM)
    oval(s, x + Inches(1.05), Inches(1.2), Inches(0.8), Inches(0.8), c)
    tb(s, x + Inches(1.05), Inches(1.25), Inches(0.8), Inches(0.7), str(i+1), 24, WHITE, True, PP_ALIGN.CENTER)
    tb(s, x + Emu(50000), Inches(2.1), Inches(2.9), Inches(0.4), title, 15, c, True, PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        y = Inches(2.5) + j * Inches(0.5)
        tb(s, x + Emu(80000), y, Inches(2.8), Inches(0.45), item, 10, TEXT_DARK)

rrect(s, Inches(0.3), Inches(4.5), Inches(9.2), Inches(0.4), CREAM)
tb(s, Inches(0.5), Inches(4.52), Inches(9.0), Inches(0.35),
   "コンサルタントは「企業戦略」「個人キャリア」「政策」の3視点から未来を語れる存在であるべき",
   10, GOLD, True, PP_ALIGN.CENTER)
footer(s, 14)
add_notes(s, """未来年表の3つの活用視点をまとめました。
1つ目は企業戦略。世代別特性に応じた人材の採用・育成、AI代替率を考慮した事業転換、地域格差を踏まえた拠点配置。これらを未来年表の数字で裏付けた提案ができます。
2つ目は個人キャリア。AI代替されにくい創造性・対人力のスキル開発、地域格差を踏まえた住まい選択、インフレや技術変化を考慮した資産形成。個人向けのキャリアコンサルティングにも活用できます。
3つ目は政策提言。AI代替に伴う雇用対策・再教育、地域格差の是正、高齢化・人口減に対応する社会保障制度の設計。自治体や業界団体への提言にも使えます。
コンサルタントは、この3つの視点から未来を語れる存在であるべきです。そうすることで、単なるシステム提案ではなく、クライアントの未来を共に創るパートナーになれるのです。""")

# ============================================================
# SLIDE 20: まとめ
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, WHITE)
title_bar(s, "まとめ：未来を知り、未来を創る")

takeaways = [
    ("人口減少は\n不可避", "2050年に0.95億人\n(-23%)\nしかし技術力で\n世界3位へ", RED_ACCENT),
    ("格差拡大が\n最大テーマ", "平均年収は上がるが\n中央値は横ばい\n地域格差は2.7倍に", ORANGE_ACCENT),
    ("AI代替は\n脅威ではなく機会", "標準化・DX推進で\n生産性は飛躍的向上\n提案のチャンス", GREEN_ACCENT),
    ("未来年表は\nコンサルの武器", "数値に基づく提案で\n説得力が格段に向上\nロードマップの根拠に", GOLD_ACCENT),
]

for i, (title, desc, c) in enumerate(takeaways):
    x = Inches(0.2) + i * Inches(2.45)
    rrect(s, x, Inches(1.1), Inches(2.3), Inches(3.2), CREAM)
    oval(s, x + Inches(0.8), Inches(1.2), Inches(0.6), Inches(0.6), c)
    tb(s, x + Inches(0.8), Inches(1.25), Inches(0.6), Inches(0.55), str(i+1), 18, WHITE, True, PP_ALIGN.CENTER)
    lines = title.split('\n')
    for li, line in enumerate(lines):
        tb(s, x + Emu(30000), Inches(1.9) + Emu(li*200000), Inches(2.2), Inches(0.3), line, 12, c, True, PP_ALIGN.CENTER)
    desc_lines = desc.split('\n')
    for li, line in enumerate(desc_lines):
        tb(s, x + Emu(30000), Inches(2.5) + Emu(li*180000), Inches(2.2), Inches(0.25), line, 9, TEXT_GRAY, False, PP_ALIGN.CENTER)

# Bottom bar
rrect(s, Inches(0.3), Inches(4.4), Inches(9.2), Inches(0.5), GOLD)
tb(s, Inches(0.5), Inches(4.42), Inches(8.8), Inches(0.45),
   "未来を知ることは、未来を創ること ― コンサルタントの真価はここにある",
   14, WHITE, True, PP_ALIGN.CENTER)
footer(s, 15)
add_notes(s, """本日のまとめです。4つのキーメッセージを確認しましょう。
1つ目、人口減少は不可避。2050年に0.95億人、23%減。しかし技術力では世界3位に上昇します。悲観ではなく、戦略的に対応すべきテーマです。
2つ目、格差拡大が最大のテーマ。平均年収は上がるが中央値は横ばい。地域格差は2.7倍に。クライアントの「どの層」「どの地域」かを見極めることが重要です。
3つ目、AI代替は脅威ではなく機会。標準化・DX推進で生産性は飛躍的に向上します。これはコンサルタントにとって最大の提案チャンスです。
4つ目、未来年表はコンサルの武器。数値に基づく提案で説得力が格段に向上し、ロードマップの根拠になります。
「未来を知ることは、未来を創ること」。これがコンサルタントの真価です。""")

# ============================================================
# SLIDE 21: エンドスライド
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, Emu(0), Emu(0), SW, SH, CREAM)
rect(s, Emu(0), Emu(0), Emu(3200000), SH, GOLD_ACCENT)
rect(s, Emu(3200000), Emu(0), Emu(100000), SH, GOLD_LIGHT)
tb(s, Inches(4.0), Inches(1.5), Inches(5.5), Inches(1.0), "ご清聴ありがとうございました", 28, TEXT_BROWN, True)
ml(s, Inches(4.0), Inches(2.8), Inches(3), Inches(0.6), [("H A R M O N I C", 14, TEXT_BROWN)])
ml(s, Inches(6.0), Inches(2.8), Inches(3), Inches(0.6), [("i n s i g h t", 14, GOLD_LIGHT)])
tb(s, Inches(4.0), Inches(3.5), Inches(5.5), Inches(0.5), "未来を知り、共に歩む。", 14, BROWN_LIGHT)
tb(s, Inches(4.0), Inches(4.3), Inches(5), Inches(0.3), "Harmonic Insight 2026年3月12日", 10, BROWN_LIGHT)
add_notes(s, """以上で本日の勉強会は終了です。ご清聴ありがとうございました。
本日お伝えした未来年表のデータは、明日からの提案活動にすぐに活用いただけます。
クライアントとの会話の中で「2050年には〇〇になる」と数字を添えて語るだけで、提案の深みが全く変わります。
ぜひ今日の資料を手元に置いて、ロードマップ策定やクライアントへのプレゼンに活かしてください。
ご質問やディスカッションがあれば、ぜひお気軽にどうぞ。
Harmonic Insightは、未来を知り、共に歩むパートナーとして、皆さんのコンサルティング活動を支援してまいります。""")

# ===== Save =====
out = "コンサルタントが知っておくべき未来年表_勉強会資料.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
