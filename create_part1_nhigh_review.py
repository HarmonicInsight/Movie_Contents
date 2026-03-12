#!/usr/bin/env python3
"""
前編：N高・ゼロ高の本音レビュー＋あなたのタイプ診断（約10分 / 7スライド）
YouTube SEOキーワード: N高 評判 本音 通信制高校 タイプ診断
"""
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx_helpers import *

SERIES = "通信制高校選び完全ガイド【前編】"
prs = create_presentation()

# ════════════════════════════════════════════════════
# SLIDE 1: タイトル（冒頭フック重視）
# ════════════════════════════════════════════════════
slide = add_title_slide(
    prs,
    "N高・ゼロ高の本音レビュー",
    "＋ あなたに合う通信制はどのタイプ？",
    "〜 通信制高校選び完全ガイド【前編】 〜",
    "Harmonic Insight 2026年3月"
)

# フィードバック②: 冒頭を「問題提起」から始める
add_notes(slide,
    "N高やゼロ高に興味を持ったけど、学費を見て驚いた。\n"
    "あるいは、自分に本当に合うのか不安になった。\n"
    "そんな方に今日は伝えたいことがあります。\n"
    "実は通信制高校は全国に300校以上あって、\n"
    "あなたにもっと合う学校がある可能性が高いんです。\n\n"
    "こんにちは、Harmonic Insightです。\n"
    "この動画は「通信制高校選び完全ガイド」の前編です。\n"
    "N高・ゼロ高を本音で分析した上で、\n"
    "あなたにどんなタイプの通信制が合うのか、4つの診断で特定します。\n"
    "中編ではタイプ別のおすすめ校、後編では学費攻略と行動プランをお伝えしますので、\n"
    "ぜひ3本セットでご覧ください。"
)


# ════════════════════════════════════════════════════
# SLIDE 2: なぜ今「通信制高校」なのか？
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "なぜ今「通信制高校」が注目されるのか？")
add_footer(slide, 1, SERIES)

# Left column
add_rounded_rect(slide, MARGIN, Emu(1005840), Emu(3931920), Emu(3520440), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(1051560), Emu(3748440), Emu(365760),
            "数字が示す変化", "Calibri", Pt(16), True, C_GOLD, PP_ALIGN.LEFT)

bg_items = [
    ("300校以上", "全国の通信制高校の数"),
    ("26.5万人", "通信制で学ぶ高校生（2024年）"),
    ("6.6%", "高校生全体に占める割合"),
    ("10年で1.5倍", "通信制の生徒数増加率"),
]
y = Emu(1463040)
for num, desc in bg_items:
    add_textbox(slide, Emu(640080), y, Emu(1280160), Emu(365760),
                num, "Calibri", Pt(16), True, C_ACCENT_BLUE, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(1920240), y + Emu(36576), Emu(2468880), Emu(320040),
                desc, "Calibri", Pt(11), False, C_DARK, PP_ALIGN.LEFT)
    y += Emu(411480)

# Right column
add_rounded_rect(slide, Emu(4572000), Emu(1005840), Emu(4114800), Emu(3520440), C_WHITE)
add_textbox(slide, Emu(4663440), Emu(1051560), Emu(3931920), Emu(365760),
            "でも、こんな不安ありませんか？", "Calibri", Pt(16), True, C_ACCENT_ORANGE, PP_ALIGN.LEFT)

anxieties = [
    "「通信制」って聞くとネガティブなイメージ…",
    "N高やゼロ高は有名だけど、自分に合う？",
    "学費は全日制と比べてどうなの？",
    "卒業後の進路は大丈夫？",
    "友達はできるの？",
]
y = Emu(1463040)
for a in anxieties:
    add_textbox(slide, Emu(4754880), y, Emu(3748440), Emu(365760),
                f"  {a}", "Calibri", Pt(12), False, C_DARK, PP_ALIGN.LEFT)
    y += Emu(365760)

# Bottom
add_rounded_rect(slide, Emu(1828800), Emu(3749040), Emu(5486400), Emu(457200), C_GOLD)
add_textbox(slide, Emu(1828800), Emu(3749040), Emu(5486400), Emu(457200),
            "この動画を見れば、すべての不安に答えが見つかります",
            "Calibri", Pt(13), True, C_WHITE, PP_ALIGN.CENTER)

add_notes(slide,
    "まず数字を見てみましょう。\n"
    "通信制高校は全国に300校以上あります。N高とゼロ高だけじゃないんです。\n"
    "現在26万5千人の高校生が通信制で学んでいて、この10年で1.5倍に増えています。\n"
    "でも、皆さんこんな不安を感じていませんか？\n"
    "「通信制って聞くとネガティブなイメージがある」\n"
    "「N高は有名だけど、本当に自分に合う？」\n"
    "「学費は？進路は？友達は？」\n"
    "この動画を最後まで見ていただければ、\n"
    "すべての不安に対する答えが見つかります。\n"
    "まずはN高とゼロ高を本音で分析していきましょう。"
)


# ════════════════════════════════════════════════════
# SLIDE 3: N高・S高グループの本音分析
# フィードバック④: 独自データ解釈を加える
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "N高等学校・S高等学校グループ【本音分析】", "数字の裏側を読み解く")
add_footer(slide, 2, SERIES)

# Merit box
add_rounded_rect(slide, MARGIN, Emu(1005840), Emu(3931920), Emu(1280160),
                 RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Emu(548640), Emu(1051560), Emu(3748440), Emu(320040),
            "強み", "Calibri", Pt(14), True, C_ACCENT_GREEN, PP_ALIGN.LEFT)
merits = [
    "先進的オンライン学習 + 豊富な課外活動",
    "プログラミング・起業の実践カリキュラム",
    "全国どこからでも入学可能",
]
y = Emu(1371600)
for m in merits:
    add_textbox(slide, Emu(640080), y, Emu(3657600), Emu(274320),
                f"  {m}", "Calibri", Pt(11), False, C_DARK)
    y += Emu(274320)

# Demerit box
add_rounded_rect(slide, Emu(4663440), Emu(1005840), Emu(3931920), Emu(1280160),
                 RGBColor(0xFF, 0xEB, 0xEE))
add_textbox(slide, Emu(4754880), Emu(1051560), Emu(3748440), Emu(320040),
            "注意点", "Calibri", Pt(14), True, RGBColor(0xC6, 0x28, 0x28), PP_ALIGN.LEFT)
demerits = [
    "自己管理能力が必須（ネットコース）",
    "対面サポートが限定的になりがち",
    "生徒数が多く個別対応に限界あり",
]
y = Emu(1371600)
for d in demerits:
    add_textbox(slide, Emu(4846320), y, Emu(3657600), Emu(274320),
                f"  {d}", "Calibri", Pt(11), False, C_DARK)
    y += Emu(274320)

# フィードバック④: 独自データ分析セクション
add_rounded_rect(slide, MARGIN, Emu(2468880), CONTENT_W, Emu(2103120), C_WHITE)
add_textbox(slide, Emu(548640), Emu(2468880), CONTENT_W, Emu(320040),
            "数字の裏側を読み解く", "Calibri", Pt(14), True, C_GOLD, PP_ALIGN.LEFT)

# 独自分析カード
analysis_items = [
    ("在籍2.7万人", "1学年あたり約9,000人\n担任1人あたりの生徒数は\n数百人規模になる計算", C_ACCENT_BLUE),
    ("ネットコース\n年25〜38万円", "月額換算で約2〜3万円\n全日制私立の半額以下\nただしサポート校併用で+α", C_ACCENT_GREEN),
    ("難関大合格実績\nあり", "東大・早慶の実績はあるが\n全体の合格率は非公開\n自主学習力が鍵", C_ACCENT_ORANGE),
    ("口コミ分析で\n見えたパターン", "「自由すぎて不安」が\n最多の不満ポイント\n向き不向きが明確に分かれる", C_ACCENT_PURPLE),
]

card_w = Emu(1920240)
gap = Emu(137160)
for i, (title, desc, color) in enumerate(analysis_items):
    x = Emu(548640) + (card_w + gap) * i
    y_card = Emu(2834640)
    add_rounded_rect(slide, x, y_card, card_w, Emu(1645920), C_BG_LIGHT)
    add_rect(slide, x, y_card, card_w, Emu(45720), color)
    add_textbox(slide, x + Emu(45720), y_card + Emu(91440), card_w - Emu(91440), Emu(457200),
                title, "Calibri", Pt(11), True, color, PP_ALIGN.CENTER)
    add_textbox(slide, x + Emu(45720), y_card + Emu(594360), card_w - Emu(91440), Emu(960120),
                desc, "Calibri", Pt(9), False, C_DARK, PP_ALIGN.CENTER)

add_notes(slide,
    "N高の強みは皆さんご存知の通り、先進的なオンライン学習と豊富な課外活動です。\n"
    "でも今日は、数字の裏側を読み解いてみましょう。\n\n"
    "在籍生徒数2万7千人。これ、1学年あたり約9,000人ということです。\n"
    "担任の先生1人あたり数百人の生徒を見ている計算になります。\n"
    "つまり、個別に手厚いサポートを期待するのは、構造的に難しい面があるんです。\n\n"
    "学費はネットコースで月額換算2万から3万円。全日制の私立と比べれば半額以下です。\n"
    "ただし、サポート校を併用すると追加費用がかかる場合もあります。\n\n"
    "難関大学の合格実績はありますが、全体の合格率は公開されていません。\n"
    "合格しているのは自主学習力が高い一部の生徒であることを理解しておきましょう。\n\n"
    "口コミを多数分析して見えたパターンがあります。\n"
    "最も多い不満は「自由すぎて不安」というもの。\n"
    "つまりN高は向き不向きが明確に分かれる学校なんです。\n"
    "では、自分にはどんな学校が向いているのか？この後の診断で確認しましょう。"
)


# ════════════════════════════════════════════════════
# SLIDE 4: ゼロ高の本音分析
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "ゼロ高等学院【本音分析】", "堀江氏の教育論 × 体験型学習の実態")
add_footer(slide, 3, SERIES)

# 特徴（コンパクト）
add_rounded_rect(slide, MARGIN, Emu(1005840), Emu(5120640), Emu(1280160), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(1051560), Emu(4937760), Emu(320040),
            "ゼロ高のコンセプト", "Calibri", Pt(14), True, C_GOLD, PP_ALIGN.LEFT)

features = [
    "「偏差値0」= 座学にとらわれない教育理念",
    "起業体験・地方創生プロジェクト等の実践中心",
    "堀江貴文氏が主宰、少人数制で密度の濃い体験",
]
y = Emu(1371600)
for f in features:
    add_textbox(slide, Emu(640080), y, Emu(4846320), Emu(274320),
                f"  {f}", "Calibri", Pt(11), False, C_DARK)
    y += Emu(274320)

# 向き不向き
add_rounded_rect(slide, MARGIN, Emu(2468880), Emu(3931920), Emu(2103120),
                 RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Emu(548640), Emu(2514600), Emu(3748440), Emu(274320),
            "こんな人に向いている", "Calibri", Pt(13), True, C_ACCENT_GREEN, PP_ALIGN.LEFT)
suited = [
    "  やりたいことが明確にある",
    "  行動力・チャレンジ精神がある",
    "  起業やフリーランスに興味がある",
    "  「とにかくやってみる」タイプ",
]
y = Emu(2834640)
for s in suited:
    add_textbox(slide, Emu(640080), y, Emu(3657600), Emu(274320),
                s, "Calibri", Pt(10), False, C_DARK)
    y += Emu(274320)

add_rounded_rect(slide, Emu(4572000), Emu(2468880), Emu(4114800), Emu(2103120),
                 RGBColor(0xFF, 0xEB, 0xEE))
add_textbox(slide, Emu(4663440), Emu(2514600), Emu(3931920), Emu(274320),
            "こんな人には合わないかも", "Calibri", Pt(13), True,
            RGBColor(0xC6, 0x28, 0x28), PP_ALIGN.LEFT)
not_suited = [
    "  手厚い学習サポートを求める方",
    "  大学受験を最優先にしたい方",
    "  安定した通学スタイルを望む方",
    "  まだやりたいことが見つかっていない方",
]
y = Emu(2834640)
for ns in not_suited:
    add_textbox(slide, Emu(4754880), y, Emu(3840480), Emu(274320),
                ns, "Calibri", Pt(10), False, C_DARK)
    y += Emu(274320)

add_notes(slide,
    "次にゼロ高等学院です。\n"
    "ゼロ高は堀江貴文さんが主宰する学校で、\n"
    "「偏差値0」つまり座学にとらわれない教育を掲げています。\n\n"
    "正直に言うと、ゼロ高はかなり人を選ぶ学校です。\n"
    "やりたいことが明確にあって、自分から動ける人には最高の環境。\n"
    "でも「まだやりたいことが見つかっていない」人が入ると、\n"
    "何をすればいいか分からなくなる可能性があります。\n\n"
    "手厚い学習サポートを期待する方にも向いていません。\n"
    "ゼロ高は「教えてもらう」場所ではなく「自分で掴み取る」場所だからです。\n\n"
    "N高もゼロ高も素晴らしい学校ですが、\n"
    "大事なのは「自分に合うかどうか」なんです。\n"
    "ここからは、あなたにどんなタイプの学校が合うか、診断していきましょう。"
)


# ════════════════════════════════════════════════════
# SLIDE 5: タイプ診断（インタラクティブ）
# フィードバック③: 具体的シナリオ＋インタラクティブ要素
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "あなたはどのタイプ？ 通信制高校適性診断",
              "一度動画を止めて、A〜Dどれに近いか考えてみてください")
add_footer(slide, 4, SERIES)

types = [
    ("A", "学習自由度\n重視タイプ",
     "朝起きるのが苦手\n部活や習い事を続けたい\n自分のペースで勉強したい",
     C_ACCENT_BLUE),
    ("B", "専門性\n追求タイプ",
     "好きなことに夢中になれる\n将来の夢がある程度決まっている\nプロを目指したい分野がある",
     C_ACCENT_GREEN),
    ("C", "コミュニティ\n重視タイプ",
     "友達と一緒に学びたい\n学校行事が好き\n先生に相談しやすい環境がいい",
     C_ACCENT_ORANGE),
    ("D", "グローバル\n挑戦タイプ",
     "海外に興味がある\n英語力を伸ばしたい\n世界で活躍したい",
     C_ACCENT_PURPLE),
]

card_w = Emu(1920240)
gap = Emu(182880)
for i, (letter, title, scenarios, color) in enumerate(types):
    x = Emu(457200) + (card_w + gap) * i
    y = Emu(1097280)
    add_rounded_rect(slide, x, y, card_w, Emu(3200400), C_WHITE)
    add_rect(slide, x, y, card_w, Emu(54864), color)
    # letter circle
    cx = x + card_w // 2 - Emu(274320)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, y + Emu(137160),
                                    Emu(548640), Emu(548640))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide, cx, y + Emu(228600), Emu(548640), Emu(365760),
                letter, "Georgia", Pt(22), True, C_WHITE, PP_ALIGN.CENTER)
    # title
    add_textbox(slide, x + Emu(45720), y + Emu(777240), card_w - Emu(91440), Emu(457200),
                title, "Calibri", Pt(13), True, C_DARK, PP_ALIGN.CENTER)
    # scenarios (具体的)
    add_textbox(slide, x + Emu(45720), y + Emu(1371600), card_w - Emu(91440), Emu(1554480),
                scenarios, "Calibri", Pt(10), False, C_GRAY, PP_ALIGN.LEFT)

# CTA
add_rounded_rect(slide, Emu(1371600), Emu(4389120), Emu(6400800), Emu(365760), C_GOLD)
add_textbox(slide, Emu(1371600), Emu(4389120), Emu(6400800), Emu(365760),
            "コメント欄にあなたのタイプ（A〜D）を書いてください！",
            "Calibri", Pt(13), True, C_WHITE, PP_ALIGN.CENTER)

add_notes(slide,
    "ここで一度動画を止めてください。\n"
    "あなたに合う通信制高校のタイプを診断します。\n\n"
    "Aタイプは学習自由度重視。\n"
    "朝起きるのが苦手な人、部活や習い事を続けたい人、\n"
    "自分のペースで勉強したい人はこのタイプです。\n\n"
    "Bタイプは専門性追求。\n"
    "好きなことに夢中になれる人、将来の夢がある程度決まっている人、\n"
    "声優やeスポーツ、IT分野などプロを目指したい人です。\n\n"
    "Cタイプはコミュニティ重視。\n"
    "友達と一緒に学びたい人、学校行事が好きな人、\n"
    "先生に気軽に相談できる環境がいい人です。\n\n"
    "Dタイプはグローバル挑戦。\n"
    "海外に興味がある人、英語力を伸ばしたい人、\n"
    "世界で活躍したいと思っている人です。\n\n"
    "どのタイプに近いか、コメント欄にA・B・C・Dで書いてください！\n"
    "複数当てはまる人は、一番強いものを選んでくださいね。\n"
    "次のスライドでさらに詳しく診断します。"
)


# ════════════════════════════════════════════════════
# SLIDE 6: 4つの診断テスト詳細
# ════════════════════════════════════════════════════
slide = new_slide(prs)
add_title_bar(slide, "4つの診断でタイプを確定しよう",
              "迷った人はこの診断で絞り込めます")
add_footer(slide, 5, SERIES)

diagnostics = [
    ("診断1", "学習スタイル", "オンライン完結派？→A\n週数日通学派？→C\n個別指導派？→B/C", C_ACCENT_BLUE),
    ("診断2", "生活リズム", "朝が苦手→A\n部活と両立→A/B\n規則正しく通学→C", C_ACCENT_GREEN),
    ("診断3", "将来の目標", "大学進学→A/C\n専門職・プロ→B\n起業→D\n海外→D", C_ACCENT_ORANGE),
    ("診断4", "サポート", "自分でできる→A/D\n手厚いケア→C\n専門指導→B", C_ACCENT_PURPLE),
]

card_w = Emu(1920240)
gap = Emu(182880)
for i, (num, title, mapping, color) in enumerate(diagnostics):
    x = Emu(457200) + (card_w + gap) * i
    y = Emu(1097280)
    add_rounded_rect(slide, x, y, card_w, Emu(3200400), C_WHITE)
    add_rect(slide, x, y, card_w, Emu(54864), color)
    add_rounded_rect(slide, x + Emu(137160), y + Emu(137160), Emu(822960), Emu(320040), color)
    add_textbox(slide, x + Emu(137160), y + Emu(137160), Emu(822960), Emu(320040),
                num, "Calibri", Pt(11), True, C_WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, x + Emu(45720), y + Emu(548640), card_w - Emu(91440), Emu(365760),
                title, "Calibri", Pt(15), True, C_DARK, PP_ALIGN.CENTER)
    add_rect(slide, x + Emu(274320), y + Emu(960120), card_w - Emu(548640), Emu(9144), C_BG_LIGHT)
    add_textbox(slide, x + Emu(91440), y + Emu(1051560), card_w - Emu(182880), Emu(1828800),
                mapping, "Calibri", Pt(11), False, C_DARK, PP_ALIGN.LEFT)

add_textbox(slide, MARGIN, Emu(4389120), CONTENT_W, Emu(274320),
            "→ 一番多かったアルファベットがあなたのタイプです！中編でおすすめ校を紹介します",
            "Calibri", Pt(11), True, C_GOLD, PP_ALIGN.CENTER)

add_notes(slide,
    "もう少し詳しく診断しましょう。4つの質問に答えてください。\n\n"
    "診断1、学習スタイル。\n"
    "完全にオンラインで学びたい人はAタイプ。\n"
    "週に数日は通学したい人はCタイプ。\n"
    "専門分野の個別指導を受けたい人はBかCです。\n\n"
    "診断2、生活リズム。\n"
    "朝が苦手な人はA。部活と両立したい人はAかB。\n"
    "規則正しく通学したい人はCです。\n\n"
    "診断3、将来の目標。\n"
    "大学進学ならAかC、専門職やプロを目指すならB、\n"
    "起業や海外ならDです。\n\n"
    "診断4、サポートの好み。\n"
    "自分でどんどん進められるならAかD、\n"
    "手厚いケアが欲しいならC、専門的な指導が欲しいならBです。\n\n"
    "一番多かったアルファベットがあなたのメインタイプです。\n"
    "タイプが分かったら、中編の動画でおすすめの学校を紹介しますので、\n"
    "ぜひ続けてご覧ください。"
)


# ════════════════════════════════════════════════════
# SLIDE 7: エンドカード（次の動画誘導重視）
# フィードバック⑩: 次の動画への誘導を最大化
# ════════════════════════════════════════════════════
slide = add_end_slide(
    prs,
    summary_items=[
        "N高は「自由すぎて不安」が最大の落とし穴",
        "ゼロ高は「やりたいことが明確な人」向け",
        "通信制は300校以上、あなたに合う学校が必ずある",
    ],
    next_video_title="【中編】タイプ別おすすめ通信制高校",
    next_video_desc="あなたのタイプに最適な学校を具体的に紹介します"
)

add_notes(slide,
    "前編のまとめです。\n"
    "1つ目、N高の最大の落とし穴は「自由すぎて不安になる」こと。\n"
    "自己管理が得意な人には最高の環境ですが、そうでない人には合わない可能性があります。\n\n"
    "2つ目、ゼロ高は「やりたいことが明確にある人」向け。\n"
    "行動力がある人には最高ですが、まだ模索中の人には他の選択肢がおすすめです。\n\n"
    "3つ目、通信制高校は300校以上あります。\n"
    "N高やゼロ高だけでなく、あなたに合う学校が必ずあります。\n\n"
    "さて、あなたのタイプは分かりましたか？\n"
    "中編では、A・B・C・Dそれぞれのタイプに最適な学校を\n"
    "具体的に紹介します。\n"
    "画面に表示されている中編の動画をぜひクリックしてください。\n"
    "あなたのタイプの学校が見つかるはずです。"
)

# ── 保存 ──
output = "通信制高校ガイド_前編_N高ゼロ高レビューとタイプ診断.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Total slides: {len(prs.slides)}")
