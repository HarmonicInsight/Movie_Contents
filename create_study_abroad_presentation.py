#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
海外教育比較 × グローバルオンライン教育 YouTube勉強会用プレゼンテーション
Harmonic Insight テンプレートスタイル準拠

スピーカーノート = YouTube動画のナレーション台本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ===== Color Palette (Harmonic Insight Style) =====
GOLD = RGBColor(0x8B, 0x75, 0x36)
GOLD_LIGHT = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_ACCENT = RGBColor(0xB8, 0x86, 0x0B)
DARK_BG = RGBColor(0x1A, 0x1A, 0x1A)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_GRAY = RGBColor(0x66, 0x66, 0x66)
BLUE_ACCENT = RGBColor(0x2C, 0x7A, 0xB0)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
ORANGE_ACCENT = RGBColor(0xE6, 0x7E, 0x22)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
PURPLE_ACCENT = RGBColor(0x8E, 0x44, 0xAD)

SLIDE_WIDTH = Emu(9144000)
SLIDE_HEIGHT = Emu(5143500)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ===== Helper Functions =====
def add_shape(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False,
                 color=TEXT_DARK, alignment=PP_ALIGN.LEFT, font_name="Yu Gothic"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, items, font_size=12,
                    color=TEXT_DARK, font_name="Yu Gothic", spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
    return txBox

def add_gold_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape

def set_notes(slide, text):
    """スライドにスピーカーノート（ナレーション台本）を設定"""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def make_title_slide(title, subtitle="", notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=DARK_BG)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Pt(4), fill_color=GOLD)
    add_shape(slide, 0, SLIDE_HEIGHT - Pt(4), SLIDE_WIDTH, Pt(4), fill_color=GOLD)
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.5),
                 title, font_size=32, bold=True, color=GOLD_LIGHT,
                 alignment=PP_ALIGN.CENTER, font_name="Yu Gothic")
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.2), Inches(8), Inches(1),
                     subtitle, font_size=16, color=CREAM,
                     alignment=PP_ALIGN.CENTER, font_name="Yu Gothic")
    add_text_box(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.4),
                 "Harmonic Insight", font_size=11, color=GOLD,
                 alignment=PP_ALIGN.RIGHT, font_name="Yu Gothic")
    if notes:
        set_notes(slide, notes)
    return slide

def make_content_slide(title, bullet_items, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=CREAM)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Pt(3), fill_color=GOLD)
    add_shape(slide, 0, Pt(3), SLIDE_WIDTH, Inches(0.7), fill_color=WHITE)
    add_text_box(slide, Inches(0.5), Inches(0.05), Inches(9), Inches(0.6),
                 title, font_size=20, bold=True, color=GOLD,
                 alignment=PP_ALIGN.LEFT, font_name="Yu Gothic")
    add_gold_line(slide, Inches(0.5), Inches(0.72), Inches(9))
    add_bullet_text(slide, Inches(0.6), Inches(0.9), Inches(8.8), Inches(3.8),
                    bullet_items, font_size=13, color=TEXT_DARK)
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.3),
                 "Harmonic Insight", font_size=9, color=GOLD,
                 alignment=PP_ALIGN.RIGHT, font_name="Yu Gothic")
    if notes:
        set_notes(slide, notes)
    return slide

def make_two_column_slide(title, left_title, left_items, right_title, right_items,
                          notes="", left_color=GOLD, right_color=BLUE_ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=CREAM)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Pt(3), fill_color=GOLD)
    add_shape(slide, 0, Pt(3), SLIDE_WIDTH, Inches(0.7), fill_color=WHITE)
    add_text_box(slide, Inches(0.5), Inches(0.05), Inches(9), Inches(0.6),
                 title, font_size=20, bold=True, color=GOLD,
                 alignment=PP_ALIGN.LEFT, font_name="Yu Gothic")
    add_gold_line(slide, Inches(0.5), Inches(0.72), Inches(9))
    # Left
    add_shape(slide, Inches(0.4), Inches(0.9), Inches(4.3), Inches(0.4), fill_color=left_color)
    add_text_box(slide, Inches(0.5), Inches(0.92), Inches(4.1), Inches(0.35),
                 left_title, font_size=13, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, font_name="Yu Gothic")
    add_bullet_text(slide, Inches(0.5), Inches(1.4), Inches(4.1), Inches(3.4),
                    left_items, font_size=11, color=TEXT_DARK)
    # Right
    add_shape(slide, Inches(5.2), Inches(0.9), Inches(4.3), Inches(0.4), fill_color=right_color)
    add_text_box(slide, Inches(5.3), Inches(0.92), Inches(4.1), Inches(0.35),
                 right_title, font_size=13, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, font_name="Yu Gothic")
    add_bullet_text(slide, Inches(5.3), Inches(1.4), Inches(4.1), Inches(3.4),
                    right_items, font_size=11, color=TEXT_DARK)
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.3),
                 "Harmonic Insight", font_size=9, color=GOLD,
                 alignment=PP_ALIGN.RIGHT, font_name="Yu Gothic")
    if notes:
        set_notes(slide, notes)
    return slide

def make_three_column_slide(title, col_data, notes=""):
    """col_data: list of (col_title, items, color)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=CREAM)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Pt(3), fill_color=GOLD)
    add_shape(slide, 0, Pt(3), SLIDE_WIDTH, Inches(0.7), fill_color=WHITE)
    add_text_box(slide, Inches(0.5), Inches(0.05), Inches(9), Inches(0.6),
                 title, font_size=20, bold=True, color=GOLD,
                 alignment=PP_ALIGN.LEFT, font_name="Yu Gothic")
    add_gold_line(slide, Inches(0.5), Inches(0.72), Inches(9))
    col_width = Inches(2.9)
    for i, (col_title, items, col_color) in enumerate(col_data):
        x = Inches(0.35 + i * 3.1)
        add_shape(slide, x, Inches(0.9), col_width, Inches(0.35), fill_color=col_color)
        add_text_box(slide, x + Inches(0.05), Inches(0.9), col_width - Inches(0.1), Inches(0.35),
                     col_title, font_size=12, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER, font_name="Yu Gothic")
        add_bullet_text(slide, x + Inches(0.1), Inches(1.35), col_width - Inches(0.2), Inches(3.5),
                        items, font_size=10, color=TEXT_DARK, spacing=Pt(4))
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.3),
                 "Harmonic Insight", font_size=9, color=GOLD,
                 alignment=PP_ALIGN.RIGHT, font_name="Yu Gothic")
    if notes:
        set_notes(slide, notes)
    return slide

def make_highlight_slide(title, key_number, key_label, description_items, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=CREAM)
    add_shape(slide, 0, 0, SLIDE_WIDTH, Pt(3), fill_color=GOLD)
    add_shape(slide, 0, Pt(3), SLIDE_WIDTH, Inches(0.7), fill_color=WHITE)
    add_text_box(slide, Inches(0.5), Inches(0.05), Inches(9), Inches(0.6),
                 title, font_size=20, bold=True, color=GOLD,
                 alignment=PP_ALIGN.LEFT, font_name="Yu Gothic")
    add_gold_line(slide, Inches(0.5), Inches(0.72), Inches(9))
    add_shape(slide, Inches(0.5), Inches(1.0), Inches(3.5), Inches(2.0),
              fill_color=RGBColor(0x2A, 0x24, 0x18))
    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(3.3), Inches(1.2),
                 key_number, font_size=40, bold=True, color=GOLD_LIGHT,
                 alignment=PP_ALIGN.CENTER, font_name="Yu Gothic")
    add_text_box(slide, Inches(0.6), Inches(2.2), Inches(3.3), Inches(0.6),
                 key_label, font_size=14, color=CREAM,
                 alignment=PP_ALIGN.CENTER, font_name="Yu Gothic")
    add_bullet_text(slide, Inches(4.3), Inches(1.0), Inches(5.2), Inches(3.5),
                    description_items, font_size=12, color=TEXT_DARK)
    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.3),
                 "Harmonic Insight", font_size=9, color=GOLD,
                 alignment=PP_ALIGN.RIGHT, font_name="Yu Gothic")
    if notes:
        set_notes(slide, notes)
    return slide


# ============================================================
# SLIDE 1: フック＋タイトル
# ============================================================
make_title_slide(
    "通信制・オンラインで学ぶあなたへ\n── 世界と比べたら、意外な事実が見えた ──",
    "教育の選択肢シリーズ 第2回\n日米欧オンライン教育 徹底比較",
    notes=(
        "皆さんこんにちは、Harmonic Insightです。\n"
        "アメリカの高校生は、成績だけじゃなくて課外活動で大学合否が決まる。\n"
        "イギリスでは大学の前に1年間、働いたり海外に出たりするのが当たり前。\n"
        "韓国では深夜まで塾に通い続ける受験戦争が社会問題になっている。\n"
        "──じゃあ日本の通信制高校やオンライン大学は、世界から見てどうなのか？\n\n"
        "今日の動画は、通信制高校を検討している方、今通っている方、\n"
        "あるいは社会人でオンライン学習を考えている方に向けてお話しします。\n"
        "「将来グローバルに活躍したい」「海外大学にも興味がある」\n"
        "そんな方には特に見てほしい内容です。\n\n"
        "前回の第1回では通信制高校の選び方を解説しました。\n"
        "今回はその通信制・オンライン教育が、世界と比べてどこが進んでいて、\n"
        "どこにまだ伸びしろがあるのかを徹底比較していきます。\n"
        "最後に「今日からできる具体的なアクション」もお伝えするので、\n"
        "ぜひ最後までご覧ください。"
    )
)

# ============================================================
# SLIDE 2: 目次
# ============================================================
make_content_slide(
    "教育の選択肢シリーズ 第2回 ── 今日お話しすること",
    [
        "1. いま世界で何が起きている？── オンライン教育市場の爆発的成長",
        "",
        "2. 海外の高校生・大学生は何をしている？── 米国・英国・アジアの実態",
        "",
        "3. アメリカのオンライン高校 ── 日本と何が違うのか？",
        "",
        "4. ヨーロッパの教育 ── 「学ぶ権利」という考え方",
        "",
        "5. 日本の通信制 ── 世界と比べた「強み」と「伸びしろ」",
        "",
        "6. 学費・質保証・サポート体制 ── 日米欧まるごと比較",
        "",
        "7. 今日からできること ── MOOCs活用と具体的な次のステップ",
    ],
    notes=(
        "今日の動画では、この7つのテーマでお話ししていきます。\n"
        "前回の第1回動画では通信制高校の選び方を解説しましたが、\n"
        "今回はその通信制を「世界基準」で評価し直してみます。\n\n"
        "特に注目していただきたいのは、学費の比較です。\n"
        "同じオンラインで大学卒業するのに、アメリカと日本で5倍以上の差がある。\n"
        "この事実を知っているかどうかで、進路の考え方が変わるはずです。\n\n"
        "そして最後に、今日この動画を見終わった後に\n"
        "実際にできる具体的なアクションもお伝えします。\n"
        "ぜひ最後まで見てくださいね。"
    )
)

# ============================================================
# SLIDE 3: 市場の爆発的成長
# ============================================================
make_highlight_slide(
    "オンライン教育市場 ── いま世界で何が起きているか",
    "3,888億$",
    "2025年 世界市場規模",
    [
        "2000年から市場規模は900%に拡大",
        "2030年には5,647億ドルへ（年率7.75%成長）",
        "2029年までに約11億人がオンライン学習に参加",
        "コロナ後も米国の学習率は28%を維持",
        "　（パンデミック前はわずか10%未満）",
        "AI・VR・クラウド技術が成長をさらに加速",
    ],
    notes=(
        "まず全体像です。世界のオンライン教育市場は2025年時点で約3,888億ドル、\n"
        "日本円にするとおよそ58兆円という巨大な規模になっています。\n"
        "2000年と比べて900パーセント、つまり10倍近くに成長しました。\n\n"
        "で、この数字が皆さんにとって何を意味するかというと──\n"
        "今オンラインで学んでいるあなたの経験は、\n"
        "10年後の社会では「当たり前のスキル」になるということです。\n"
        "通信制やオンライン学習を選ぶこと自体が、実は最先端の経験なんですね。\n\n"
        "注目すべきは、コロナが終わった後も成長が止まっていないということです。\n"
        "アメリカではパンデミック前のオンライン学習率はわずか10%未満でしたが、\n"
        "コロナが収束した2022年でも28%を維持しています。\n"
        "つまり、一度オンライン学習を経験した人は「これでいい」と実感したんですね。\n\n"
        "さらにAIやVR技術の進化が、この流れをますます加速させています。\n"
        "もはやオンライン教育は緊急時の代替手段ではなく、\n"
        "教育システムの中核になりつつある。\n"
        "皆さんは、その変化のまさに渦中にいるわけです。"
    )
)

# ============================================================
# SLIDE 4: オンライン学習の利点と課題
# ============================================================
make_two_column_slide(
    "オンライン学習 ── 世界共通の「光」と「影」",
    "光：普遍的な利点",
    [
        "・いつでもどこでも学べる",
        "  → 地理的・時間的な制約がなくなる",
        "・AIで一人ひとりに最適化された学習",
        "  → アダプティブラーニングの進化",
        "・録画で繰り返し学習が可能",
        "・通学費・寮費が不要でコスト削減",
        "・年齢・外見に左右されない公平な環境",
    ],
    "影：共通の課題",
    [
        "・友達と会えない → 社会性が育ちにくい",
        "・モチベーションを自分で維持する難しさ",
        "・デジタルデバイド",
        "  → ネット環境・デバイスを持てない人がいる",
        "・実習系の科目はオンラインに限界がある",
        "・教員のオンライン指導スキルが追いついていない",
    ],
    notes=(
        "オンライン学習には世界共通の「光」と「影」があります。\n\n"
        "まず光の部分。場所や時間に縛られずに学べるのは最大の強みですよね。\n"
        "しかもAIの進化で、一人ひとりのレベルに合わせた学習が可能になっています。\n"
        "授業の録画があれば何度でも見返せますし、\n"
        "通学費や寮費がかからないのでコストも大幅に下がります。\n"
        "さらに面白いのが、オンラインでは年齢や外見で判断されないので、\n"
        "発言の内容そのもので評価されるフェアな環境が生まれるということです。\n\n"
        "一方で影の部分もあります。\n"
        "友達に会えない、孤独感がある、モチベーションが続かない。\n"
        "これはオンライン学習の最大の課題として世界中で指摘されています。\n"
        "また、ネット環境やパソコンが手に入らないという「デジタルデバイド」も深刻です。\n"
        "技術が進歩しても、その恩恵を受けられない人が取り残されてしまう。\n"
        "この「光」と「影」のバランスをどう取るかが、各国の教育政策の最大のテーマになっています。"
    ),
    left_color=GREEN_ACCENT,
    right_color=RED_ACCENT,
)

# ============================================================
# SLIDE 5: 米国・英国・アジアの学生の「今」
# ============================================================
make_three_column_slide(
    "世界の高校生・大学生は何をしている？",
    [
        ("アメリカ", [
            "成績だけでは大学に入れない",
            "課外活動の「継続性」が評価される",
            "クラブ・ボランティア・インターンを",
            "　長く深く続けることが重要",
            "大学は「なぜその活動をしたか」を",
            "　面接で深掘りしてくる",
        ], BLUE_ACCENT),
        ("イギリス", [
            "「ギャップイヤー」が当たり前",
            "大学の前後に1年間の社会経験",
            "83%が就業、56%が海外滞在",
            "企業の94%がギャップイヤー経験者を",
            "　積極採用と回答",
            "ワーホリ費用：年144万〜226万円",
        ], GREEN_ACCENT),
        ("アジア（韓国・シンガポール）", [
            "韓国：「SKY大学」がエリートの象徴",
            "高校生は深夜まで塾に通う受験戦争",
            "大学入学後も成績・資格の競争が続く",
            "シンガポール：大学進学率は3〜4割",
            "職業訓練校（ポリテク・ITE）が主要進路",
            "「効率的に学んでキャリアに直結」が鍵",
        ], ORANGE_ACCENT),
    ],
    notes=(
        "ここからが今日のハイライトです。\n"
        "世界の高校生が実際に何をしているのか、国別に見ていきましょう。\n\n"
        "まずアメリカ。これ、びっくりする方多いと思うんですが、\n"
        "アメリカではテストの点数だけでは大学に入れないんです。\n"
        "クラブ活動やボランティア、インターンシップ。\n"
        "大事なのは「たくさんやること」じゃなくて「1つを長く深く続けること」。\n"
        "大学は面接で「なぜその活動を選んだのか」を徹底的に聞いてきます。\n"
        "つまり「どんな人間か」が問われるんですね。\n"
        "皆さんはどう思いますか？日本と全然違いますよね。\n\n"
        "次にイギリス。「ギャップイヤー」という文化があって、\n"
        "大学の前後に1年間、働いたり海外に行ったりするんです。\n"
        "しかも企業の94パーセントが「ギャップイヤー経験者を積極採用したい」と回答。\n"
        "日本だと「空白期間」はマイナスですよね？\n"
        "イギリスでは逆に「主体性がある」と評価されるんです。\n\n"
        "アジアでは韓国の受験戦争が壮絶で、\n"
        "高校生が文字通り深夜まで塾に通い続けます。\n"
        "一方シンガポールは「全員が大学に行く必要はない」と割り切って、\n"
        "職業訓練校を主要な進路として確立しています。\n\n"
        "この3つの中で、皆さんが一番驚いたのはどれですか？\n"
        "ぜひコメント欄で教えてください。"
    )
)

# ============================================================
# SLIDE 6: アメリカのオンライン高校
# ============================================================
make_two_column_slide(
    "アメリカのオンライン高校 ── 大学進学を徹底サポート",
    "カリキュラムの特徴",
    [
        "・数学、科学、英語、社会科のコア科目",
        "・外国語、デジタル技術、美術史など選択科目",
        "・Honors / AP（大学レベル先取り）コース",
        "・NCAA承認 → 学生アスリートにも対応",
        "・週30時間の学習（80〜90%がPC上）",
        "・代表校：Pearson Online Academy",
        "　　　　 Connections Academy",
    ],
    "手厚いサポート体制",
    [
        "・科目別の専門教員が指導",
        "・ホームルーム教員が全科目をモニタリング",
        "・専任スクールカウンセラーが常駐",
        "・4年間の学習計画を一緒に作成",
        "・大学出願プロセスの支援",
        "・奨学金情報の提供",
        "・入試対策の無料オンラインコース",
    ],
    notes=(
        "アメリカのオンライン高校は、日本の通信制高校とはかなり違います。\n"
        "まずカリキュラム。AP、つまりAdvanced Placementという\n"
        "大学レベルの授業を高校のうちから受けられるんです。\n"
        "大学進学を目指す生徒向けに、かなり戦略的な設計になっています。\n\n"
        "でも本当にすごいのはサポート体制のほうです。\n"
        "科目ごとの専門教員、全科目をモニタリングするホームルーム教員、\n"
        "そして専任のスクールカウンセラー。この3層で、\n"
        "4年間の学習計画から大学出願、奨学金情報まで一貫してサポートしてくれます。\n"
        "「オンラインで授業が受けられる」だけじゃなくて、\n"
        "大学合格までの道筋をまるごと面倒見てくれる仕組みなんですね。\n\n"
        "学費は年間1,800ドルから2,800ドル、日本円で約27万円から42万円。\n\n"
        "で、皆さん気になりますよね。「日本から入学できるの？」って。\n"
        "実はこれらの学校の一部は、日本にいながらオンラインで入学できます。\n"
        "具体的な入学方法やステップは次回の第3回動画で詳しく解説しますので、\n"
        "ぜひチャンネル登録しておいてくださいね。\n"
        "また、このカリキュラムの考え方を日本で取り入れている\n"
        "AIE国際高等学校のような学校もあります。\n"
        "詳しくは概要欄にリンクを貼っておきます。"
    ),
    left_color=BLUE_ACCENT,
    right_color=BLUE_ACCENT,
)

# ============================================================
# SLIDE 7: ヨーロッパのオンライン教育
# ============================================================
make_two_column_slide(
    "ヨーロッパのオンライン教育 ── 「学ぶ権利」という思想",
    "高校教育",
    [
        "・「高等教育は人権である」が基本理念",
        "・英国/米国カリキュラム、IBプログラム",
        "・「ヨーロピアンアワー」で違う国の生徒と交流",
        "・CNED（仏）：欧州最大の遠隔教育機関",
        "・学費：年間€630〜€5,900",
    ],
    "大学教育",
    [
        "・Open University（英）：幅広い分野で学位",
        "・FernUniversität in Hagen（独）：",
        "  ドイツ最大の州立遠隔教育大学",
        "・ロンドン大学：キャンパスと同質の学位",
        "・EU/EEA圏の学生は授業料ほぼ無料",
        "・非EU学生でも年€5,000〜€18,000",
        "・フランスの単位がドイツでも使える",
        "  → 共通ルール「ESG」が国をまたいだ学びを保証",
    ],
    notes=(
        "次にヨーロッパです。ヨーロッパの教育の根底には、\n"
        "「教育は人権である」という考え方があります。\n"
        "だからEU圏の学生は大学の授業料がほぼ無料なんですね。\n\n"
        "オンライン高校では英国カリキュラムや国際バカロレアが選べますし、\n"
        "面白いのは「ヨーロピアンアワー」という取り組みです。\n"
        "フランスの生徒とドイツの生徒がオンラインで一緒に授業を受ける。\n"
        "日本にいながらこの環境に参加できるとしたら、すごいですよね。\n\n"
        "大学レベルでは、もっとすごいことが起きています。\n"
        "例えばフランスの大学で取った単位が、ドイツの大学でもそのまま使える。\n"
        "国をまたいで学んだり働いたりすることを前提にした仕組みなんです。\n"
        "これを可能にしているのが「ESG」という共通の質保証ルールです。\n\n"
        "また、別の国の学生とオンラインで一緒にプロジェクトをやる\n"
        "「COIL」という仕組みもあって、\n"
        "自宅にいながら国際的なチームワークを経験できます。\n\n"
        "こういう仕組みは日本にはまだありません。\n"
        "でも逆に言えば、日本のオンライン教育にはここに大きな伸びしろがあるということです。"
    ),
    left_color=GREEN_ACCENT,
    right_color=GREEN_ACCENT,
)

# ============================================================
# SLIDE 8: 日本の通信制高校
# ============================================================
make_two_column_slide(
    "日本の通信制 ── 世界と比べた「強み」と「伸びしろ」",
    "世界に誇れる強み",
    [
        "・学費の安さは世界最高レベル",
        "  公立：年間1.6万〜5万円（米国の1/10以下）",
        "  就学支援金で実質無料も可能",
        "・大学進学率が過去最高を更新中（27%）",
        "  専門学校含む合計進学率は52%超",
        "・N高：東大・海外大学への進学実績",
        "・ゼロ高：起業家育成・インターン",
    ],
    "世界と比べた「伸びしろ」",
    [
        "・サポート体制の格差が大きい",
        "  → 米国のような3層サポートは一部校のみ",
        "・国際的な質保証の仕組みがない",
        "  → 欧州のESGに相当するものが未整備",
        "・国をまたいだ単位互換の仕組みがない",
        "・キャリア支援が大学出願に比べ手薄",
        "・公立通信制では孤立リスクが課題",
    ],
    notes=(
        "さて、ここからは日本の通信制を世界と比較して評価してみましょう。\n"
        "制度の基本は前回の第1回動画で解説しましたので、\n"
        "今日は「世界と比べてどうなのか」に焦点を当てます。\n\n"
        "まず強みから。学費の安さは世界最高レベルです。\n"
        "公立なら年間1万6千円から5万円。アメリカの10分の1以下です。\n"
        "しかも就学支援金で実質無料になるケースもある。\n"
        "これは正直、すごいことなんですよ。\n\n"
        "進学実績も過去最高を更新中で、N高からは東大に行く生徒もいますし、\n"
        "ゼロ高のように起業家を育てるユニークな学校も出てきています。\n\n"
        "一方で伸びしろもはっきり見えます。\n"
        "さっきアメリカのオンライン高校には\n"
        "科目教員・ホームルーム教員・カウンセラーの3層サポートがあると言いましたよね。\n"
        "日本でこれに近いのはN高やゼロ高など一部の学校だけで、\n"
        "公立の通信制は基本的に自学自習。孤立して止まってしまうリスクがあります。\n\n"
        "また、ヨーロッパのような国際的な質保証の仕組みや、\n"
        "国をまたいだ単位互換の仕組みもまだありません。\n"
        "逆に言えば、ここが整備されれば日本の通信制は世界最強になる可能性がある。\n"
        "皆さんが通信制を選んでいるということ自体が、\n"
        "その変化の最前線にいるということなんです。"
    ),
    left_color=ORANGE_ACCENT,
    right_color=ORANGE_ACCENT,
)

# ============================================================
# SLIDE 9: ZEN大学 ── 日本の新しい挑戦
# ============================================================
make_highlight_slide(
    "ZEN大学 ── 日本初の本格オンライン大学（2025年開学）",
    "年間38万円",
    "知能情報社会学部 定員3,500名",
    [
        "6分野を横断して学ぶオーダーメイドカリキュラム",
        "　数理／情報／文化・思想／社会／経済／デジタル産業",
        "279科目から自分だけの時間割を設計",
        "大半がオンデマンド授業（自分のペースで学習）",
        "3種のメンターが手厚くサポート",
        "　クラス・コーチ／アカデミック・アドバイザー",
        "　／キャリアアドバイザー",
        "pixiv提携科目や東浩紀氏の対話型講座も",
    ],
    notes=(
        "日本のオンライン大学の新しい動きとして、\n"
        "2025年4月に開学したZEN大学を紹介しておきたいと思います。\n\n"
        "注目してほしいのは学費です。年間38万円。\n"
        "これ、次のスライドで学費比較するんですが、\n"
        "アメリカのオンライン大学と比べると5分の1以下なんですよ。\n\n"
        "279科目から自分だけの時間割を組むオーダーメイドカリキュラムで、\n"
        "大半の授業がオンデマンド。自分のペースで進められます。\n\n"
        "そして面白いのが、さっきアメリカのオンライン高校の3層サポートを紹介しましたよね。\n"
        "ZEN大学もクラス・コーチ、アカデミック・アドバイザー、\n"
        "キャリアアドバイザーという3種類のメンターがつくんです。\n"
        "日本の大学でここまで手厚いのは珍しい。\n"
        "アメリカの良いところを取り入れようとしている印象を受けます。\n\n"
        "日本のオンライン大学はまだ始まったばかりですが、\n"
        "ZEN大学がその成功モデルになれるかどうか、注目ですね。"
    )
)

# ============================================================
# SLIDE 10: 日米欧 費用比較
# ============================================================
make_three_column_slide(
    "学費の比較 ── 日米欧でこれだけ違う",
    [
        ("アメリカ", [
            "【オンライン高校】",
            "年間$1,800〜$2,800",
            "（約27万〜42万円）",
            "高額校は$20,000超もあり",
            "",
            "【オンライン大学】",
            "学士号：$40,536〜$63,185",
            "（約600万〜950万円）",
            "修士号：$19,000〜$27,000",
            "ただし通学費・寮費は不要",
        ], BLUE_ACCENT),
        ("ヨーロッパ", [
            "【オンライン高校】",
            "年間€630〜€5,900",
            "（約10万〜90万円）",
            "",
            "【オンライン大学】",
            "EU/EEA圏：無料〜低額",
            "非EU圏：年€5,000〜€18,000",
            "ドイツの州立大学は非常に安価",
            "「教育は人権」の理念が根底に",
        ], GREEN_ACCENT),
        ("日本", [
            "【通信制高校】",
            "公立：年間1.6万〜5万円",
            "私立：年間10万〜100万円",
            "就学支援金で実質無料も可能",
            "",
            "【オンライン大学】",
            "ZEN大学：年間38万円",
            "",
            "全日制高校の平均：約51万円/年",
            "通信制は大幅にコスト削減可能",
        ], ORANGE_ACCENT),
    ],
    notes=(
        "ここで学費を一気に比較してみましょう。\n"
        "まず1つ、衝撃的な数字をお見せします。\n\n"
        "アメリカのオンライン大学で学士号を取ると、600万円から950万円。\n"
        "日本のZEN大学は4年間で152万円。\n"
        "同じ「オンラインで大学卒業」なのに、5〜6倍の差があるんです。\n"
        "これ、すごくないですか？\n\n"
        "高校レベルでも差は歴然です。\n"
        "アメリカのオンライン高校が年間27万円から42万円。\n"
        "日本の公立通信制は年間1万6千円から5万円。\n"
        "桁が違いますよね。\n\n"
        "ヨーロッパは「教育は人権」という理念があるので、\n"
        "EU圏の学生は大学がほぼ無料。\n"
        "日本人でも年間5,000ユーロから学べるので、アメリカより格段に安いです。\n\n"
        "つまり日本の通信制は、費用面では世界最強クラスなんです。\n"
        "問題は「安かろう悪かろう」ではないかどうか。\n"
        "次のスライドで質保証の仕組みを見ていきましょう。\n"
        "ちなみに細かい費用の比較表は概要欄に載せておきますので、\n"
        "あとでチェックしてみてください。"
    )
)

# ============================================================
# SLIDE 11: 質保証の比較
# ============================================================
make_content_slide(
    "教育の質をどう保証している？── 認定制度の比較",
    [
        "【アメリカ】「地域認定」が信頼性の最高基準",
        "  ・Cogniaなどの認定機関が厳格に審査",
        "  ・認定なし = 単位互換できない、就職で不利になる",
        "  ・「ディプロマミル」（偽の学位販売業者）を排除する仕組み",
        "",
        "【ヨーロッパ】ESG = 欧州全体の共通品質フレームワーク",
        "  ・どの国・どの大学でも同じ基準で質を担保",
        "  ・国境を越えた単位互換・学位認定のベースになっている",
        "",
        "【日本】文部科学省の学習指導要領に基づく運営",
        "  ・通信制高校は文科省の規定で卒業要件が定められている",
        "  ・ただし、オンライン教育に特化した包括的な質保証の",
        "    フレームワークはまだ確立途上",
        "  ・国際的な認定の調和が今後の課題",
    ],
    notes=(
        "学費が安いのは良いことですが、問題は教育の質ですよね。\n"
        "各地域でどうやって質を保証しているのか見てみましょう。\n\n"
        "アメリカでは「地域認定」という仕組みが最も信頼されています。\n"
        "Cogniaなどの認定機関が学校を厳しく審査して、基準を満たしていれば認定を出す。\n"
        "認定を受けていない学校の単位は他の大学に持っていけないし、\n"
        "就職でも不利になります。\n"
        "実はアメリカには「ディプロマミル」と呼ばれる\n"
        "偽の学位を売る業者も存在するので、\n"
        "認定は消費者保護の役割も果たしているんです。\n\n"
        "ヨーロッパにはESGという共通の質保証基準があります。\n"
        "フランスの大学で取った単位がドイツの大学でも認められる。\n"
        "国をまたいで学んだり働いたりすることを前提にした仕組みで、\n"
        "これはヨーロッパの大きな強みです。\n\n"
        "日本はどうかというと、通信制高校は文科省の規定でしっかり管理されていますが、\n"
        "オンライン教育に特化した包括的な質保証フレームワークは、\n"
        "まだ確立されていないのが現状です。\n"
        "今後ZEN大学のようなオンライン専門大学が増えていくなかで、\n"
        "この仕組み作りが急務になってくると思います。"
    )
)

# ============================================================
# SLIDE 12: サポート体制の比較
# ============================================================
make_content_slide(
    "学習サポートの比較 ── 「放っておかれる」か「伴走してもらえる」か",
    [
        "【アメリカ】キャリアまで見据えた「伴走型サポート」",
        "  ・科目別教員＋ホームルーム教員＋専任カウンセラーの3層構造",
        "  ・大学出願、奨学金、入試対策まで一貫してサポート",
        "  ・課外活動やクラブで社会性も補完",
        "",
        "【ヨーロッパ】テクノロジーを活かした学習支援",
        "  ・FernUniversitätの演習システム：即時フィードバック＋難易度自動調整",
        "  ・COIL（国際共同オンライン学習）で海外の学生と協働",
        "",
        "【日本】学校によって差が大きい",
        "  ・N高：アバターチャットで交流促進＋多彩な課外授業",
        "  ・ゼロ高：専属コーチが目標設定から伴走",
        "  ・サポート校で学習面・精神面を補完するケースも",
        "  ・ただし公立通信制では自学自習が基本 → 孤立リスク",
    ],
    notes=(
        "オンライン学習で一番大事なのは、実はサポート体制です。\n"
        "放っておかれるのか、伴走してもらえるのかで、結果はまったく違ってきます。\n\n"
        "アメリカのオンライン高校は、ここが本当に手厚い。\n"
        "科目の先生、ホームルームの先生、そしてカウンセラーという3層のサポートがあって、\n"
        "大学の出願から奨学金まで一貫して面倒を見てくれます。\n"
        "社会性の面でも、オンラインのクラブ活動やNational Honor Societyのような\n"
        "課外活動で補完する仕組みが用意されています。\n\n"
        "ヨーロッパでは、テクノロジーを活かしたサポートが特徴的です。\n"
        "ドイツのFernUniversitätには演習システムがあって、\n"
        "解答すると即座にフィードバックが返ってきて、\n"
        "しかも難易度が自動で調整されるんです。\n"
        "またCOILという仕組みで、別の国の学生とオンラインで\n"
        "一緒にプロジェクトをやる機会も作られています。\n\n"
        "日本はどうかというと、正直なところ学校によって差が大きいです。\n"
        "N高やゼロ高のように手厚いサポートを提供している学校がある一方で、\n"
        "公立の通信制高校では基本的に自学自習で、\n"
        "孤立して学習が止まってしまうリスクも指摘されています。\n"
        "この格差をどう埋めていくかが、日本の通信制教育の大きな課題です。"
    )
)

# ============================================================
# SLIDE 13: MOOCsの活用
# ============================================================
make_content_slide(
    "今日からできること ── MOOCsで世界の授業を体験する",
    [
        "【Coursera】 月額$59〜 / 年額$399〜",
        "  ・Google、MIT、ハーバード等の講座が9,000コース以上",
        "  ・修了証は履歴書に書ける → 転職・社内評価に繋がった事例も",
        "",
        "【Udemy】 コースごとの買い切り（セールで1,200円台も）",
        "  ・20万コース以上、最も幅広い分野をカバー",
        "",
        "【edX】 東大がMIT・ハーバードと連携して講座を公開",
        "",
        "【オンライン英会話】 月額2,000〜13,000円",
        "",
        ">>> 今日のアクション <<<",
        "Courseraで「Google Career Certificates」と検索してみよう",
        "→ 無料で内容が見られる。これが世界基準のオンライン学習です",
    ],
    notes=(
        "ここまで各国の制度を見てきましたが、\n"
        "「で、自分は今日から何ができるの？」というのが一番大事ですよね。\n"
        "ここでは具体的なアクションをお伝えします。\n\n"
        "まず、今日この動画を見終わったら1つだけやってみてください。\n"
        "Courseraで「Google Career Certificates」と検索するんです。\n"
        "無料で講座の中身を見ることができます。\n"
        "これがアメリカの学生が実際に受けているオンライン学習のレベルです。\n"
        "「あ、自分でもできるかも」と思えたら、それが第一歩です。\n\n"
        "Courseraの修了証は実際に履歴書に書けますし、\n"
        "「Courseraの修了証がきっかけで転職できた」という事例も報告されています。\n\n"
        "もっと手軽に始めたいなら、Udemyがおすすめです。\n"
        "セールのときは1,200円台でコースが買えることもあります。\n\n"
        "edXでは東京大学がMITやハーバードと連携した講座を出していますし、\n"
        "英語力を上げたいならオンライン英会話が月額2,000円からで始められます。\n\n"
        "1つ注意点として、講義は無料でも修了証は有料というケースが多いです。\n"
        "でもGoogle認定証のような就職に直結する資格が\n"
        "自宅から取れるのは、大きな武器ですよね。\n"
        "具体的なリンクは概要欄に全部貼っておきますので、チェックしてみてください。"
    )
)

# ============================================================
# SLIDE 14: 日本への提言
# ============================================================
make_content_slide(
    "あなたが今日からできる5つのこと",
    [
        "1. Courseraで無料講座を1つ受けてみる",
        "   → 「Google Career Certificates」がおすすめ",
        "",
        "2. 自分の学校のサポート体制を確認する",
        "   → カウンセラー、キャリア相談、課外活動は使える？",
        "",
        "3. 海外オンライン高校・大学のサイトを見てみる",
        "   → 「自分でも入れるかも」と思えたら、次回の第3回動画へ",
        "",
        "4. ギャップイヤーの体験談をYouTubeで検索する",
        "   → 「gap year experience」で検索すると大量に出てくる",
        "",
        "5. この動画のコメント欄で「自分の学び方」をシェアする",
        "   → あなたの経験が、同じ悩みを持つ誰かの助けになる",
    ],
    notes=(
        "さて、ここが今日の動画で一番大事なパートです。\n"
        "世界の教育を比較してきましたが、「で、自分はどうすればいいの？」\n"
        "ということで、今日からできることを5つにまとめました。\n\n"
        "1つ目。さっきも言いましたが、Courseraで無料講座を1つ受けてみてください。\n"
        "Google Career Certificatesがおすすめです。\n"
        "世界のオンライン学習のレベルを体感できます。\n\n"
        "2つ目。自分の学校のサポート体制を確認してみてください。\n"
        "カウンセラーに相談できるのか、キャリア支援はあるのか、課外活動は何がある？\n"
        "意外と使っていないサービスがあるかもしれません。\n\n"
        "3つ目。海外のオンライン高校や大学のサイトを覗いてみてください。\n"
        "英語のサイトでも、Google翻訳で十分読めます。\n"
        "「あ、自分でも入れるかも」と思えたら、次回の第3回動画で\n"
        "具体的な入学ステップを解説しますので、そちらも見てください。\n\n"
        "4つ目。YouTubeで「gap year experience」と検索してみてください。\n"
        "イギリスやアメリカの学生がギャップイヤーの体験を\n"
        "vlogで公開しています。見るだけで視野が広がりますよ。\n\n"
        "5つ目。この動画のコメント欄で、あなたの学び方をシェアしてください。\n"
        "「通信制でこう学んでいる」「MOOCsをこう使っている」\n"
        "そういうリアルな声が、同じ悩みを持つ誰かの助けになります。\n\n"
        "日本のオンライン教育にはまだ課題があります。\n"
        "でも逆に言えば、今オンラインで学ぶことを選んでいるあなたは、\n"
        "日本の教育が変わっていく最前線にいるということです。\n"
        "あなたは課題の当事者じゃない。変化の主体なんです。"
    )
)

# ============================================================
# SLIDE 15: まとめ
# ============================================================
make_content_slide(
    "まとめ ── あなたの学びは、世界とつながっている",
    [
        "【世界の潮流】",
        "  ・オンライン教育は一時的なブームではなく、教育の「中核」へ",
        "  ・2030年には約85兆円規模に成長",
        "",
        "【日本の通信制は世界と比較して…】",
        "  ・学費の安さは世界最高レベル（強み）",
        "  ・サポート体制と質保証の仕組みに伸びしろ",
        "",
        "【あなたは「変化の最前線」にいる】",
        "  ・通信制、MOOCs、ギャップイヤー、留学…",
        "    「全日制に通う」だけが正解ではない時代",
        "  ・大事なのは「どう学ぶか」を自分で選ぶこと",
        "",
        "教育の選択肢シリーズ",
        "  第1回：通信制高校の選び方  /  第2回：本動画（日米欧比較）",
        "  第3回：海外の通信制に日本から入学する方法",
    ],
    notes=(
        "最後に、今日の内容をまとめます。\n\n"
        "オンライン教育はもう一時的なブームではなく、教育の中核になりつつあります。\n"
        "2030年には85兆円規模になると予測されている。\n"
        "その中で皆さんは、まさに最前線にいるわけです。\n\n"
        "今日見てきた中で覚えておいてほしいのは3つ。\n"
        "1つ目、日本の通信制は学費の安さでは世界最高レベル。\n"
        "2つ目、アメリカの「伴走型サポート」と\n"
        "ヨーロッパの「質保証の仕組み」は日本も取り入れるべき。\n"
        "3つ目、そして一番大事なこと──\n"
        "「全日制に通うだけが正解ではない時代」に、\n"
        "自分で学び方を選んでいるあなたは、すでに一歩先を行っています。\n\n"
        "この動画は「教育の選択肢シリーズ」の第2回です。\n"
        "第1回の「通信制高校の選び方」をまだ見ていない方は、\n"
        "概要欄にリンクがありますのでぜひチェックしてください。\n"
        "そして次回の第3回では「海外の通信制高校に日本から入学する方法」を\n"
        "具体的なステップで解説していきます。"
    )
)

# ============================================================
# SLIDE 16: CTA ＋ Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT, fill_color=DARK_BG)
add_shape(slide, 0, 0, SLIDE_WIDTH, Pt(4), fill_color=GOLD)
add_shape(slide, 0, SLIDE_HEIGHT - Pt(4), SLIDE_WIDTH, Pt(4), fill_color=GOLD)
add_text_box(slide, Inches(1), Inches(1.0), Inches(8), Inches(1.0),
             "ご視聴ありがとうございました", font_size=30, bold=True, color=GOLD_LIGHT,
             alignment=PP_ALIGN.CENTER, font_name="Yu Gothic")
add_text_box(slide, Inches(1), Inches(2.2), Inches(8), Inches(1.5),
             "チャンネル登録・高評価で応援していただけると嬉しいです\n\n"
             "コメント欄であなたの「学び方」をシェアしてください\n"
             "同じ悩みを持つ誰かの助けになります",
             font_size=15, color=CREAM,
             alignment=PP_ALIGN.CENTER, font_name="Yu Gothic")
add_text_box(slide, Inches(1), Inches(3.9), Inches(8), Inches(0.8),
             "教育の選択肢シリーズ\n"
             "第1回：通信制高校の選び方（公開中）  │  "
             "第3回：海外の通信制に日本から入学する方法（次回）",
             font_size=12, bold=True, color=GOLD_ACCENT,
             alignment=PP_ALIGN.CENTER, font_name="Yu Gothic")
add_text_box(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.4),
             "Harmonic Insight", font_size=11, color=GOLD,
             alignment=PP_ALIGN.RIGHT, font_name="Yu Gothic")
set_notes(slide, (
    "ご視聴いただきありがとうございました。\n\n"
    "今日は日本・アメリカ・ヨーロッパのオンライン教育を比較して、\n"
    "日本の通信制の「強み」と「伸びしろ」を見てきました。\n\n"
    "覚えておいてほしいのは、\n"
    "オンラインで学ぶことを選んでいるあなたは、\n"
    "日本の教育が変わっていく最前線にいるということです。\n\n"
    "今日お伝えした「5つのアクション」、\n"
    "まずは1つでいいので試してみてください。\n"
    "Courseraで「Google Career Certificates」を検索するだけなら、\n"
    "5分でできます。\n\n"
    "そしてコメント欄で、あなたの学び方をシェアしてくれたら嬉しいです。\n"
    "「通信制でこう学んでいる」「MOOCsでこれを受けた」\n"
    "そういうリアルな声が、同じ悩みを持つ誰かの背中を押します。\n\n"
    "この動画は「教育の選択肢シリーズ」の第2回です。\n"
    "第1回の「通信制高校の選び方」は概要欄にリンクがあります。\n"
    "次回の第3回では「海外の通信制高校に日本から入学する方法」を\n"
    "具体的なステップで解説していきますので、\n"
    "見逃さないようにチャンネル登録をお願いします。\n"
    "それでは、また次の動画でお会いしましょう。"
))

# ===== Save =====
output_path = "海外学習_留学_グローバルオンライン教育_プレゼンテーション.pptx"
prs.save(output_path)
print(f"プレゼンテーションを保存しました: {output_path}")
print(f"合計スライド数: {len(prs.slides)}")
