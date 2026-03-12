#!/usr/bin/env python3
"""
通信制高校選び完全ガイド YouTube動画用パワーポイント生成スクリプト
スライド = 動画の画像、プレゼンテーションノート = ナレーション台本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 定数 ──
SLIDE_WIDTH = 9144000
SLIDE_HEIGHT = 5143500
MARGIN = Emu(457200)
CONTENT_W = Emu(8229600)

# カラーパレット（既存プレゼンから踏襲）
C_DARK = RGBColor(0x33, 0x33, 0x33)
C_GOLD = RGBColor(0x8B, 0x75, 0x36)
C_GOLD_LIGHT = RGBColor(0xC9, 0xA8, 0x4C)
C_GRAY = RGBColor(0x66, 0x66, 0x66)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BG_LIGHT = RGBColor(0xF5, 0xF0, 0xE6)
C_BROWN = RGBColor(0x3C, 0x2F, 0x1E)
C_ACCENT_BLUE = RGBColor(0x2B, 0x5C, 0x8A)
C_ACCENT_GREEN = RGBColor(0x3A, 0x7D, 0x5E)
C_ACCENT_ORANGE = RGBColor(0xC4, 0x6B, 0x2B)
C_ACCENT_PURPLE = RGBColor(0x6B, 0x4C, 0x9A)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ── ヘルパー関数 ──
def add_textbox(slide, left, top, width, height, text, font_name="Calibri",
                font_size=Pt(14), bold=False, color=C_DARK, alignment=PP_ALIGN.LEFT,
                word_wrap=True, anchor=MSO_ANCHOR.TOP):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.name = font_name
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    try:
        tf.paragraphs[0].anchor = anchor
    except:
        pass
    return txbox


def add_multi_text(slide, left, top, width, height, lines, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, font_size, bold, color)"""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, (text, font_size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = "Calibri"
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return txbox


def add_rounded_rect(slide, left, top, width, height, fill_color, border=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not border:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_footer(slide, page_num):
    bar = add_rect(slide, 0, Emu(4800600), Emu(SLIDE_WIDTH), Emu(342900), C_BROWN)
    add_textbox(slide, MARGIN, Emu(4800600), Emu(5486400), Emu(342900),
                "H A R M O N I C   i n s i g h t", "Calibri", Pt(8),
                False, C_WHITE, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(8229600), Emu(4800600), Emu(731520), Emu(342900),
                str(page_num), "Calibri", Pt(8), False, C_WHITE, PP_ALIGN.CENTER)


def add_title_bar(slide, title_text, subtitle_text=None):
    add_textbox(slide, MARGIN, Emu(182880), CONTENT_W, Emu(502920),
                title_text, "Calibri", Pt(22), True, C_DARK, PP_ALIGN.LEFT)
    # gold underline
    add_rect(slide, MARGIN, Emu(685800), Emu(1097280), Emu(36576), C_GOLD_LIGHT)
    if subtitle_text:
        add_textbox(slide, MARGIN, Emu(868680), CONTENT_W, Emu(320040),
                    subtitle_text, "Calibri", Pt(13), True, C_GOLD, PP_ALIGN.LEFT)


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_card(slide, left, top, width, height, title, desc, icon_text, accent_color):
    """Add a card-style box with icon circle, title and description"""
    card = add_rounded_rect(slide, left, top, width, height, C_WHITE)
    # accent top bar
    add_rect(slide, left, top, width, Emu(36576), accent_color)
    # icon circle
    cx = left + Emu(int(width) // 2) - Emu(274320)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, top + Emu(91440), Emu(548640), Emu(548640))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()
    add_textbox(slide, cx, top + Emu(182880), Emu(548640), Emu(365760),
                icon_text, "Calibri", Pt(18), True, C_WHITE, PP_ALIGN.CENTER)
    # title
    add_textbox(slide, left + Emu(45720), top + Emu(731520), width - Emu(91440), Emu(320040),
                title, "Calibri", Pt(13), True, C_DARK, PP_ALIGN.CENTER)
    # desc
    add_textbox(slide, left + Emu(45720), top + Emu(1005840), width - Emu(91440), Emu(640080),
                desc, "Calibri", Pt(10), False, C_GRAY, PP_ALIGN.CENTER)


def new_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


# ════════════════════════════════════════════════════
# SLIDE 1: タイトルスライド
# ════════════════════════════════════════════════════
slide = new_slide()
# top & bottom gold lines
add_rect(slide, 0, 0, Emu(SLIDE_WIDTH), Emu(54864), C_GOLD_LIGHT)
add_rect(slide, 0, Emu(5088636), Emu(SLIDE_WIDTH), Emu(54864), C_GOLD_LIGHT)

# Hi icon
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(4114800), Emu(914400), Emu(914400), Emu(822960))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF0, 0xE6)
shape.line.fill.background()
add_textbox(slide, Emu(4114800), Emu(960120), Emu(914400), Emu(822960),
            "Hi", "Georgia", Pt(22), False, C_GOLD_LIGHT, PP_ALIGN.CENTER)

# decorative lines
add_rect(slide, Emu(914400), Emu(2057400), Emu(2926080), Emu(9144), C_GOLD_LIGHT)
add_rect(slide, Emu(5303520), Emu(2057400), Emu(2926080), Emu(9144), C_GOLD_LIGHT)

# Main title
add_textbox(slide, MARGIN, Emu(2194560), CONTENT_W, Emu(502920),
            "N高・ゼロ高を超える！", "Georgia", Pt(36), False, C_BROWN, PP_ALIGN.CENTER)
add_textbox(slide, MARGIN, Emu(2697480), CONTENT_W, Emu(457200),
            "通信制高校選び 完全ガイド", "Georgia", Pt(28), False, C_BROWN, PP_ALIGN.CENTER)

# subtitle line
add_rect(slide, Emu(2286000), Emu(3200400), Emu(4572000), Emu(9144), C_GOLD_LIGHT)

add_textbox(slide, MARGIN, Emu(3383280), CONTENT_W, Emu(365760),
            "〜 新たな選択肢としての通信制高校 〜", "Calibri", Pt(16), False, C_GOLD, PP_ALIGN.CENTER)

# Harmonic Insight branding
add_textbox(slide, MARGIN, Emu(3886200), CONTENT_W, Emu(274320),
            "H A R M O N I C", "Calibri", Pt(14), False, C_BROWN, PP_ALIGN.CENTER)
add_textbox(slide, MARGIN, Emu(4114800), CONTENT_W, Emu(274320),
            "i n s i g h t", "Calibri", Pt(14), False, C_GOLD_LIGHT, PP_ALIGN.CENTER)
add_textbox(slide, MARGIN, Emu(4434840), CONTENT_W, Emu(274320),
            "Harmonic Insight 2026年3月12日", "Calibri", Pt(10), False,
            RGBColor(0x8B, 0x73, 0x55), PP_ALIGN.CENTER)

add_notes(slide,
    "皆さん、こんにちは。Harmonic Insightへようこそ。\n"
    "今日は「N高・ゼロ高を超える！通信制高校選び完全ガイド」というテーマでお話しします。\n"
    "N高やゼロ高の名前は聞いたことがあるけど、実際どうなの？他にどんな選択肢があるの？\n"
    "そんな疑問に答えながら、あなたにぴったりの通信制高校を見つけるためのヒントをお伝えします。\n"
    "保護者の方も、生徒の皆さんも、ぜひ最後までご覧ください。"
)


# ════════════════════════════════════════════════════
# SLIDE 2: 今日の動画の内容（目次）
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "今日の動画の内容")
add_footer(slide, 1)

items = [
    ("1", "N高・ゼロ高で感じたワクワクを\n「私の選択軸」に変える", C_ACCENT_BLUE),
    ("2", "N高・ゼロ高を\n客観的に見つめ直す", C_GOLD),
    ("3", "タイプ別\n最適校ガイド", C_ACCENT_GREEN),
    ("4", "学費と支援制度\n完全攻略", C_ACCENT_ORANGE),
    ("5", "納得の学校選び\n実践術", C_ACCENT_PURPLE),
    ("6", "先輩たちの選択と\nその後の人生", C_ACCENT_BLUE),
    ("7", "今すぐ始める\n行動プラン", RGBColor(0x8B, 0x2E, 0x2E)),
]

start_y = Emu(960120)
card_h = Emu(457200)
gap = Emu(36576)
for i, (num, text, color) in enumerate(items):
    y = start_y + (card_h + gap) * i
    # number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, y + Emu(45720), Emu(365760), Emu(365760))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide, MARGIN, y + Emu(91440), Emu(365760), Emu(274320),
                num, "Calibri", Pt(16), True, C_WHITE, PP_ALIGN.CENTER)
    # text
    add_textbox(slide, Emu(914400), y + Emu(45720), Emu(7315200), Emu(411480),
                text.replace("\n", "  "), "Calibri", Pt(14), False, C_DARK, PP_ALIGN.LEFT)
    # line
    add_rect(slide, Emu(914400), y + card_h - Emu(9144), Emu(7772400), Emu(4572), C_BG_LIGHT)

add_notes(slide,
    "今日お話しする内容は大きく7つです。\n"
    "まず、N高やゼロ高に感じたワクワクを自分の学校選びの軸に変える方法。\n"
    "次に、N高・ゼロ高を客観的に見つめ直します。メリットだけでなくデメリットも正直にお伝えします。\n"
    "そして、あなたのタイプ別に最適な通信制高校をガイドします。\n"
    "学費や支援制度の話、学校選びの実践的なテクニック、\n"
    "実際に通信制高校で成功した先輩たちのストーリー、\n"
    "最後に、今すぐ始められる行動プランをお伝えします。\n"
    "盛りだくさんですが、一つ一つ分かりやすくお話ししていきますね。"
)


# ════════════════════════════════════════════════════
# SLIDE 3: なぜ今「通信制高校」なのか？
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "なぜ今「通信制高校」が注目されるのか？")
add_footer(slide, 2)

# Left column - 背景
add_rounded_rect(slide, MARGIN, Emu(1005840), Emu(3931920), Emu(3520440), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(1051560), Emu(3748440), Emu(365760),
            "時代の変化", "Calibri", Pt(16), True, C_GOLD, PP_ALIGN.LEFT)

bg_items = [
    "N高が通信制の常識を変えた",
    "年間20万人以上が通信制で学ぶ時代",
    "偏差値だけでは測れない力が求められる",
    "多様な学び方が社会に認められ始めた",
    "テクノロジーがオンライン学習を進化させた",
]
y = Emu(1463040)
for item in bg_items:
    add_textbox(slide, Emu(640080), y, Emu(3657600), Emu(365760),
                f"  {item}", "Calibri", Pt(12), False, C_DARK, PP_ALIGN.LEFT)
    y += Emu(365760)

# Right column - メッセージ
add_rounded_rect(slide, Emu(4572000), Emu(1005840), Emu(4114800), Emu(3520440), C_WHITE)
add_textbox(slide, Emu(4663440), Emu(1051560), Emu(3931920), Emu(365760),
            "あなたへのメッセージ", "Calibri", Pt(16), True, C_ACCENT_BLUE, PP_ALIGN.LEFT)

msg_items = [
    ("「好き」を伸ばす学校がある", C_DARK),
    ("自分のペースで学べる環境がある", C_DARK),
    ("全日制だけが正解じゃない", C_ACCENT_ORANGE),
    ("選択肢を知ることが未来を変える", C_DARK),
]
y = Emu(1463040)
for msg, clr in msg_items:
    add_textbox(slide, Emu(4754880), y, Emu(3748440), Emu(411480),
                f"  {msg}", "Calibri", Pt(13), True, clr, PP_ALIGN.LEFT)
    y += Emu(457200)

# Bottom highlight
add_rounded_rect(slide, Emu(1828800), Emu(3749040), Emu(5486400), Emu(457200), C_GOLD)
add_textbox(slide, Emu(1828800), Emu(3749040), Emu(5486400), Emu(457200),
            "N高・ゼロ高が切り拓いた道の先に、もっと多くの選択肢がある",
            "Calibri", Pt(13), True, C_WHITE, PP_ALIGN.CENTER)

add_notes(slide,
    "では、なぜ今「通信制高校」がこれほど注目されているのでしょうか。\n"
    "N高等学校の登場が通信制高校の常識を大きく変えました。\n"
    "現在、年間20万人以上の高校生が通信制で学んでいます。\n"
    "偏差値だけでは測れない力が社会で求められるようになり、\n"
    "多様な学び方が認められる時代になっています。\n"
    "大事なのは「好きを伸ばせる学校がある」「自分のペースで学べる環境がある」ということ。\n"
    "全日制だけが正解ではないんです。\n"
    "N高やゼロ高が切り拓いた道の先に、実はもっとたくさんの選択肢があります。\n"
    "今日はその選択肢を一緒に見ていきましょう。"
)


# ════════════════════════════════════════════════════
# SLIDE 4: ワクワクの正体を知る（診断導入）
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "あなたの「ワクワク」の正体は？", "4つのタイプで自分を知ろう")
add_footer(slide, 3)

types = [
    ("学習自由度\n重視タイプ", "自分のペースで\n学びたい", C_ACCENT_BLUE, "A"),
    ("専門性\n追求タイプ", "好きなことを\n深めたい", C_ACCENT_GREEN, "B"),
    ("コミュニティ\n重視タイプ", "新しい仲間と\n出会いたい", C_ACCENT_ORANGE, "C"),
    ("グローバル\n挑戦タイプ", "世界を\n広げたい", C_ACCENT_PURPLE, "D"),
]

card_w = Emu(1920240)
card_h = Emu(2743200)
start_x = Emu(457200)
gap = Emu(182880)

for i, (title, desc, color, letter) in enumerate(types):
    x = start_x + (card_w + gap) * i
    y = Emu(1188720)
    # card bg
    add_rounded_rect(slide, x, y, card_w, card_h, C_WHITE)
    # top accent
    add_rect(slide, x, y, card_w, Emu(54864), color)
    # letter circle
    cx = x + card_w // 2 - Emu(320040)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, y + Emu(182880), Emu(640080), Emu(640080))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide, cx, y + Emu(274320), Emu(640080), Emu(457200),
                letter, "Georgia", Pt(24), True, C_WHITE, PP_ALIGN.CENTER)
    # title
    add_textbox(slide, x + Emu(45720), y + Emu(914400), card_w - Emu(91440), Emu(548640),
                title, "Calibri", Pt(14), True, C_DARK, PP_ALIGN.CENTER)
    # desc
    add_textbox(slide, x + Emu(45720), y + Emu(1554480), card_w - Emu(91440), Emu(640080),
                desc, "Calibri", Pt(11), False, C_GRAY, PP_ALIGN.CENTER)

# bottom message
add_textbox(slide, MARGIN, Emu(4160520), CONTENT_W, Emu(365760),
            "→ あなたはどのタイプ？タイプに合った学校を後ほどご紹介します",
            "Calibri", Pt(12), True, C_GOLD, PP_ALIGN.CENTER)

add_notes(slide,
    "さて、前作「N高・ゼロ高って実際どうなの？」を読んで「ワクワクした」という方、\n"
    "そのワクワクの正体を一緒に探りましょう。\n"
    "大きく4つのタイプに分けられます。\n"
    "Aタイプは「自分のペースで学びたい」学習自由度重視タイプ。\n"
    "Bタイプは「好きなことを深めたい」専門性追求タイプ。\n"
    "Cタイプは「新しい仲間と出会いたい」コミュニティ重視タイプ。\n"
    "Dタイプは「世界を広げたい」グローバル挑戦タイプ。\n"
    "皆さんはどのタイプに近いですか？\n"
    "この後、タイプ別におすすめの学校を紹介していきますので、楽しみにしてください。"
)


# ════════════════════════════════════════════════════
# SLIDE 5: 自己診断の詳細
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "4つの診断であなたの最適タイプを特定", "自分を知ることが学校選びの第一歩")
add_footer(slide, 4)

diagnostics = [
    ("診断1", "学習スタイル診断", "オンライン完結派？\n週数日通学派？\n個別指導派？", C_ACCENT_BLUE),
    ("診断2", "生活リズム診断", "朝型？夜型？\n部活や習い事と\n両立したい？", C_ACCENT_GREEN),
    ("診断3", "目標志向診断", "大学進学？専門職？\n起業？やりたいことを\n探したい？", C_ACCENT_ORANGE),
    ("診断4", "サポートニーズ診断", "手厚いケア希望？\n自律学習型？\nメンター必要？", C_ACCENT_PURPLE),
]

card_w = Emu(1920240)
gap = Emu(182880)
for i, (num, title, questions, color) in enumerate(diagnostics):
    x = Emu(457200) + (card_w + gap) * i
    y = Emu(1097280)
    add_rounded_rect(slide, x, y, card_w, Emu(3200400), C_WHITE)
    add_rect(slide, x, y, card_w, Emu(54864), color)
    # number badge
    add_rounded_rect(slide, x + Emu(137160), y + Emu(137160), Emu(822960), Emu(320040), color)
    add_textbox(slide, x + Emu(137160), y + Emu(137160), Emu(822960), Emu(320040),
                num, "Calibri", Pt(11), True, C_WHITE, PP_ALIGN.CENTER)
    # title
    add_textbox(slide, x + Emu(45720), y + Emu(548640), card_w - Emu(91440), Emu(457200),
                title, "Calibri", Pt(14), True, C_DARK, PP_ALIGN.CENTER)
    # divider
    add_rect(slide, x + Emu(274320), y + Emu(1005840), card_w - Emu(548640), Emu(9144), C_BG_LIGHT)
    # questions
    add_textbox(slide, x + Emu(91440), y + Emu(1097280), card_w - Emu(182880), Emu(1645920),
                questions, "Calibri", Pt(11), False, C_GRAY, PP_ALIGN.CENTER)

add_textbox(slide, MARGIN, Emu(4389120), CONTENT_W, Emu(274320),
            "→ 診断結果を組み合わせて、第3章のタイプ別ガイドで最適校を見つけましょう",
            "Calibri", Pt(11), True, C_GOLD, PP_ALIGN.CENTER)

add_notes(slide,
    "では、具体的に4つの診断で自分のタイプを特定しましょう。\n"
    "診断1は学習スタイル診断。オンラインで完結したいのか、週に数日は通学したいのか、\n"
    "個別指導を受けたいのか。自分に合った学習の形を考えてみてください。\n"
    "診断2は生活リズム診断。朝型か夜型か、部活や習い事との両立を考えているかがポイントです。\n"
    "診断3は目標志向診断。大学進学を目指すのか、専門職に就きたいのか、\n"
    "起業に興味があるのか、それともまだやりたいことを探している段階なのか。\n"
    "診断4はサポートニーズ診断。手厚いケアが欲しいのか、自分で進められるタイプなのか。\n"
    "これらの診断結果を組み合わせることで、あなたに最適な学校タイプが見えてきます。\n"
    "この後のタイプ別ガイドで、具体的な学校を紹介していきますね。"
)


# ════════════════════════════════════════════════════
# SLIDE 6: N高・S高グループの客観評価
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "N高等学校・S高等学校グループ", "その「革新性」と「現実」")
add_footer(slide, 5)

# Merit box
add_rounded_rect(slide, MARGIN, Emu(1005840), Emu(3931920), Emu(1645920), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Emu(548640), Emu(1051560), Emu(3748440), Emu(365760),
            "メリット（強み）", "Calibri", Pt(14), True, C_ACCENT_GREEN, PP_ALIGN.LEFT)
merits = [
    "先進的なオンライン学習システム",
    "豊富な課外活動・部活動",
    "プログラミング・起業など実践カリキュラム",
    "全国どこからでも入学可能",
]
y = Emu(1417320)
for m in merits:
    add_textbox(slide, Emu(640080), y, Emu(3657600), Emu(274320),
                f"  {m}", "Calibri", Pt(11), False, C_DARK)
    y += Emu(274320)

# Demerit box
add_rounded_rect(slide, Emu(4663440), Emu(1005840), Emu(3931920), Emu(1645920), RGBColor(0xFF, 0xEB, 0xEE))
add_textbox(slide, Emu(4754880), Emu(1051560), Emu(3748440), Emu(365760),
            "デメリット（注意点）", "Calibri", Pt(14), True, RGBColor(0xC6, 0x28, 0x28), PP_ALIGN.LEFT)
demerits = [
    "自己管理能力が求められる",
    "対面サポートが限定的（ネットコース）",
    "学費はコースにより大きく異なる",
    "生徒数が多く個別対応に限界がある場合も",
]
y = Emu(1417320)
for d in demerits:
    add_textbox(slide, Emu(4846320), y, Emu(3657600), Emu(274320),
                f"  {d}", "Calibri", Pt(11), False, C_DARK)
    y += Emu(274320)

# Data section
add_rounded_rect(slide, MARGIN, Emu(2834640), CONTENT_W, Emu(1645920), C_WHITE)
add_textbox(slide, Emu(548640), Emu(2880360), CONTENT_W, Emu(320040),
            "数字で見るN高・S高", "Calibri", Pt(14), True, C_GOLD, PP_ALIGN.LEFT)

data_items = [
    ("ネットコース学費", "年間約25〜38万円\n（就学支援金適用前）"),
    ("通学コース学費", "年間約70〜120万円"),
    ("在籍生徒数", "約2万7千人\n（2024年時点）"),
    ("大学進学実績", "東大・京大・早慶\nなど実績あり"),
]

dx = Emu(2057400)
for i, (label, val) in enumerate(data_items):
    x = Emu(548640) + dx * i
    add_textbox(slide, x, Emu(3200400), Emu(1828800), Emu(274320),
                label, "Calibri", Pt(10), True, C_GOLD, PP_ALIGN.CENTER)
    add_textbox(slide, x, Emu(3474720), Emu(1828800), Emu(548640),
                val, "Calibri", Pt(10), False, C_DARK, PP_ALIGN.CENTER)

add_notes(slide,
    "まず、N高等学校・S高等学校グループを客観的に見ていきましょう。\n"
    "メリットとしては、先進的なオンライン学習システム、豊富な課外活動や部活動、\n"
    "プログラミングや起業体験など実践的なカリキュラムが挙げられます。\n"
    "全国どこからでも入学できるのも大きな魅力です。\n"
    "一方、デメリットや注意点もあります。\n"
    "自己管理能力が求められること、ネットコースでは対面サポートが限定的なこと、\n"
    "生徒数が多いため個別対応に限界がある場合もあります。\n"
    "学費はネットコースで年間約25万から38万円、通学コースだと70万から120万円程度です。\n"
    "在籍生徒数は約2万7千人。大学進学実績も東大や早慶への合格者がいます。\n"
    "N高は素晴らしい学校ですが、万人にとってベストとは限りません。"
)


# ════════════════════════════════════════════════════
# SLIDE 6: ゼロ高等学院の客観評価
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "ゼロ高等学院", "堀江氏の「教育論」と体験型学習の真相")
add_footer(slide, 6)

# 特徴
add_rounded_rect(slide, MARGIN, Emu(1005840), Emu(5120640), Emu(1554480), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(1051560), Emu(4937760), Emu(320040),
            "ゼロ高の特徴", "Calibri", Pt(14), True, C_GOLD, PP_ALIGN.LEFT)

features = [
    "「偏差値0」= 座学にとらわれない独自の教育理念",
    "イベント企画・起業体験・地方創生プロジェクトなど実践中心",
    "堀江貴文氏が主宰、行動力と挑戦マインドを重視",
    "少人数制で密度の濃い体験が可能",
]
y = Emu(1371600)
for f in features:
    add_textbox(slide, Emu(640080), y, Emu(4846320), Emu(274320),
                f"  {f}", "Calibri", Pt(11), False, C_DARK)
    y += Emu(274320)

# 向き不向き
add_rounded_rect(slide, MARGIN, Emu(2697480), Emu(3931920), Emu(1737360), RGBColor(0xE8, 0xF5, 0xE9))
add_textbox(slide, Emu(548640), Emu(2743200), Emu(3748440), Emu(274320),
            "向いている人", "Calibri", Pt(13), True, C_ACCENT_GREEN, PP_ALIGN.LEFT)
suited = ["自分でやりたいことが明確にある", "行動力・チャレンジ精神がある", "起業やフリーランスに興味がある"]
y = Emu(3017520)
for s in suited:
    add_textbox(slide, Emu(640080), y, Emu(3657600), Emu(274320),
                f"  {s}", "Calibri", Pt(10), False, C_DARK)
    y += Emu(274320)

add_rounded_rect(slide, Emu(4572000), Emu(2697480), Emu(4114800), Emu(1737360), RGBColor(0xFF, 0xEB, 0xEE))
add_textbox(slide, Emu(4663440), Emu(2743200), Emu(3931920), Emu(274320),
            "向いていない人", "Calibri", Pt(13), True, RGBColor(0xC6, 0x28, 0x28), PP_ALIGN.LEFT)
not_suited = ["手厚い学習サポートを求める方", "大学受験を最優先にしたい方", "安定した通学スタイルを望む方"]
y = Emu(3017520)
for ns in not_suited:
    add_textbox(slide, Emu(4754880), y, Emu(3840480), Emu(274320),
                f"  {ns}", "Calibri", Pt(10), False, C_DARK)
    y += Emu(274320)

add_notes(slide,
    "次にゼロ高等学院です。\n"
    "ゼロ高は堀江貴文さんが主宰する学校で、「偏差値0」を掲げています。\n"
    "これは座学にとらわれず、実体験から学ぶという教育理念を表しています。\n"
    "イベント企画、起業体験、地方創生プロジェクトなど、実践的な活動が中心です。\n"
    "少人数制なので、密度の濃い体験ができるのが魅力です。\n"
    "向いているのは、やりたいことが明確にあって行動力がある人。\n"
    "起業やフリーランスに興味がある人にもぴったりです。\n"
    "一方、手厚い学習サポートを求める方や、大学受験を最優先にしたい方には、\n"
    "他の選択肢のほうが合うかもしれません。\n"
    "ゼロ高も素晴らしい学校ですが、自分に合うかどうかが大切です。"
)


# ════════════════════════════════════════════════════
# SLIDE 7: 専門性追求タイプ向け学校
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "【Bタイプ】「好き」を仕事に！専門特化校", "専門性追求タイプ向け")
add_footer(slide, 7)

schools = [
    ("ヒューマンキャンパス\n高等学校", "声優・eスポーツ・美容\nマンガなど40分野", "40以上の\n専門分野", C_ACCENT_GREEN),
    ("AOIKE高等学校", "パティシエ・調理師の\n夢を育む専門校", "製菓・調理\n特化", C_ACCENT_ORANGE),
    ("IT・プログラミング\n特化校", "最先端のIT技術を\n高校から学べる", "テック系\nスキル", C_ACCENT_BLUE),
]

card_w = Emu(2651760)
gap = Emu(182880)
for i, (name, desc, tag, color) in enumerate(schools):
    x = MARGIN + (card_w + gap) * i
    y = Emu(1097280)
    add_rounded_rect(slide, x, y, card_w, Emu(2560320), C_WHITE)
    add_rect(slide, x, y, card_w, Emu(54864), color)
    # tag
    tag_shape = add_rounded_rect(slide, x + Emu(91440), y + Emu(137160), Emu(1097280), Emu(365760), color)
    add_textbox(slide, x + Emu(91440), y + Emu(137160), Emu(1097280), Emu(365760),
                tag, "Calibri", Pt(9), True, C_WHITE, PP_ALIGN.CENTER)
    # school name
    add_textbox(slide, x + Emu(91440), y + Emu(594360), card_w - Emu(182880), Emu(548640),
                name, "Calibri", Pt(14), True, C_DARK, PP_ALIGN.LEFT)
    # description
    add_textbox(slide, x + Emu(91440), y + Emu(1188720), card_w - Emu(182880), Emu(822960),
                desc, "Calibri", Pt(11), False, C_GRAY, PP_ALIGN.LEFT)

# bottom message
add_rounded_rect(slide, MARGIN, Emu(3840480), CONTENT_W, Emu(548640), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(3886200), Emu(8046720), Emu(457200),
            "ポイント：在学中から実践的なスキルが身につく！卒業後の就職・デビューサポートも充実",
            "Calibri", Pt(12), True, C_GOLD, PP_ALIGN.LEFT)

add_notes(slide,
    "ここからはタイプ別のおすすめ学校を紹介していきます。\n"
    "まずはBタイプ、「好きなことを深めたい」専門性追求タイプの方へ。\n"
    "ヒューマンキャンパス高等学校は、声優、eスポーツ、美容、マンガなど40以上の専門分野を学べます。\n"
    "高校の学習と並行して、プロの講師から実践的なスキルを身につけられるのが特徴です。\n"
    "AOIKE高等学校は、パティシエや調理師を目指す方に特化した学校です。\n"
    "製菓・調理の専門知識を高校時代から学べます。\n"
    "IT・プログラミング特化校では、最先端の技術を若いうちから習得できます。\n"
    "これらの学校の魅力は、在学中から実践的なスキルが身につくこと。\n"
    "卒業後の就職やデビューのサポートも充実しています。"
)


# ════════════════════════════════════════════════════
# SLIDE 8: グローバル・挑戦タイプ向け学校
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "【Dタイプ】世界へ羽ばたく！海外大学進学の道", "グローバル・挑戦志向タイプ向け")
add_footer(slide, 8)

# NIC
add_rounded_rect(slide, MARGIN, Emu(1005840), Emu(3931920), Emu(2834640), C_WHITE)
add_rect(slide, MARGIN, Emu(1005840), Emu(3931920), Emu(54864), C_ACCENT_PURPLE)
add_textbox(slide, Emu(548640), Emu(1097280), Emu(3748440), Emu(365760),
            "NIC International College in Japan", "Calibri", Pt(15), True, C_DARK, PP_ALIGN.LEFT)
add_textbox(slide, Emu(548640), Emu(1463040), Emu(3748440), Emu(274320),
            "37年の実績が証明する「転換教育 」", "Calibri", Pt(12), True, C_ACCENT_PURPLE, PP_ALIGN.LEFT)
nic_points = [
    "独自の「転換教育」で英語力と思考力を徹底強化",
    "欧米トップ大学への進学実績が豊富",
    "日本にいながら海外大学準備が完結",
    "帰国後のキャリアサポートも万全",
]
y = Emu(1828800)
for p in nic_points:
    add_textbox(slide, Emu(640080), y, Emu(3657600), Emu(274320),
                f"  {p}", "Calibri", Pt(10), False, C_DARK)
    y += Emu(274320)

# AIE
add_rounded_rect(slide, Emu(4572000), Emu(1005840), Emu(4114800), Emu(2834640), C_WHITE)
add_rect(slide, Emu(4572000), Emu(1005840), Emu(4114800), Emu(54864), C_ACCENT_BLUE)
add_textbox(slide, Emu(4663440), Emu(1097280), Emu(3931920), Emu(365760),
            "AIE国際高等学校", "Calibri", Pt(15), True, C_DARK, PP_ALIGN.LEFT)
add_textbox(slide, Emu(4663440), Emu(1463040), Emu(3931920), Emu(274320),
            "国際バカロレア（IB）教育", "Calibri", Pt(12), True, C_ACCENT_BLUE, PP_ALIGN.LEFT)
aie_points = [
    "国際バカロレア認定で世界基準の教育",
    "バイリンガル環境で英語力を自然に習得",
    "海外大学への出願準備を手厚くサポート",
    "少人数制のきめ細かい指導",
]
y = Emu(1828800)
for p in aie_points:
    add_textbox(slide, Emu(4754880), y, Emu(3840480), Emu(274320),
                f"  {p}", "Calibri", Pt(10), False, C_DARK)
    y += Emu(274320)

# bottom
add_rounded_rect(slide, MARGIN, Emu(4023360), CONTENT_W, Emu(457200), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(4069080), Emu(8046720), Emu(365760),
            "海外大学進学は「遠い夢」ではない。通信制高校から世界へ羽ばたく道がある",
            "Calibri", Pt(12), True, C_GOLD, PP_ALIGN.LEFT)

add_notes(slide,
    "Dタイプ、グローバル・挑戦志向の方におすすめの学校です。\n"
    "NIC International College in Japanは37年の実績を誇る学校で、\n"
    "独自の「転換教育」で英語力と思考力を徹底的に強化します。\n"
    "欧米のトップ大学への進学実績が豊富で、\n"
    "日本にいながら海外大学への準備が完結するのが大きな魅力です。\n"
    "AIE国際高等学校は、国際バカロレア認定校として世界基準の教育を提供しています。\n"
    "バイリンガル環境で英語力を自然に身につけられ、\n"
    "少人数制できめ細かい指導が受けられます。\n"
    "海外大学進学は決して「遠い夢」ではありません。\n"
    "通信制高校から世界へ羽ばたく道は、確かに存在しているんです。"
)


# ════════════════════════════════════════════════════
# SLIDE 9: 学習自由度重視＆コミュニティ重視タイプ
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "【A・Cタイプ】自分のペースで ＆ 仲間と成長", "学習自由度重視 × コミュニティ重視")
add_footer(slide, 9)

# Left: A type
add_rounded_rect(slide, MARGIN, Emu(1005840), Emu(4023360), Emu(3520440), C_WHITE)
add_rect(slide, MARGIN, Emu(1005840), Emu(4023360), Emu(54864), C_ACCENT_BLUE)
add_textbox(slide, Emu(548640), Emu(1097280), Emu(3840480), Emu(320040),
            "Aタイプ：学習自由度重視", "Calibri", Pt(14), True, C_ACCENT_BLUE, PP_ALIGN.LEFT)

a_schools = [
    ("ルネサンス高等学校グループ", "スマホで学べる新スタイル。スクーリングは年4日程度"),
    ("広域通信制高校", "完全オンライン対応。自宅学習中心で時間を有効活用"),
]
y = Emu(1463040)
for name, desc in a_schools:
    add_textbox(slide, Emu(548640), y, Emu(3840480), Emu(274320),
                name, "Calibri", Pt(12), True, C_DARK, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(640080), y + Emu(274320), Emu(3748440), Emu(365760),
                desc, "Calibri", Pt(10), False, C_GRAY, PP_ALIGN.LEFT)
    y += Emu(731520)

add_textbox(slide, Emu(548640), y + Emu(182880), Emu(3840480), Emu(640080),
            "→ 部活・習い事・仕事と両立したい方に\n→ 自律的に学習を進められる方に最適",
            "Calibri", Pt(10), True, C_ACCENT_BLUE, PP_ALIGN.LEFT)

# Right: C type
add_rounded_rect(slide, Emu(4663440), Emu(1005840), Emu(4023360), Emu(3520440), C_WHITE)
add_rect(slide, Emu(4663440), Emu(1005840), Emu(4023360), Emu(54864), C_ACCENT_ORANGE)
add_textbox(slide, Emu(4754880), Emu(1097280), Emu(3840480), Emu(320040),
            "Cタイプ：コミュニティ重視", "Calibri", Pt(14), True, C_ACCENT_ORANGE, PP_ALIGN.LEFT)

c_schools = [
    ("クラーク記念国際高等学校", "全国展開・多様なコース。週5日通学も可能"),
    ("第一学院高等学校", "一人ひとりに寄り添う教育と豊富なキャンパスライフ"),
    ("おおぞら高等学院", "手厚いメンタルケアと進路サポート"),
    ("公立通信制高校", "学費を抑えたい方に。通学圏内で探せる"),
]
y = Emu(1463040)
for name, desc in c_schools:
    add_textbox(slide, Emu(4754880), y, Emu(3840480), Emu(274320),
                name, "Calibri", Pt(12), True, C_DARK, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(4846320), y + Emu(274320), Emu(3748440), Emu(274320),
                desc, "Calibri", Pt(9), False, C_GRAY, PP_ALIGN.LEFT)
    y += Emu(594360)

add_notes(slide,
    "続いて、AタイプとCタイプの方におすすめの学校です。\n"
    "Aタイプ、学習自由度を重視する方には、\n"
    "ルネサンス高等学校グループがおすすめです。スマホで学べる新しいスタイルで、\n"
    "スクーリングは年4日程度。部活や習い事、仕事との両立がしやすい学校です。\n"
    "広域通信制高校も、完全オンライン対応で自宅学習中心の方にぴったりです。\n"
    "Cタイプ、仲間との交流を大切にしたい方には、\n"
    "クラーク記念国際高等学校は全国にキャンパスがあり、週5日通学も可能です。\n"
    "第一学院高等学校は一人ひとりに寄り添った教育が特徴。\n"
    "おおぞら高等学院は手厚いメンタルケアが魅力です。\n"
    "学費を抑えたい方には公立の通信制高校という選択肢もあります。"
)


# ════════════════════════════════════════════════════
# SLIDE 11: 学校比較データシート
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "学校比較データシート", "主要通信制高校を一覧で比較")
add_footer(slide, 10)

# comparison table
add_rounded_rect(slide, MARGIN, Emu(960120), CONTENT_W, Emu(3657600), C_WHITE)

# Header row
comp_headers = ["学校名", "タイプ", "学費目安（年）", "特徴", "向いている人"]
comp_widths = [Emu(1645920), Emu(1005840), Emu(1280160), Emu(2194560), Emu(2103120)]
hx = MARGIN
for i, (h, w) in enumerate(zip(comp_headers, comp_widths)):
    add_rect(slide, hx, Emu(960120), w, Emu(320040), C_GOLD)
    add_textbox(slide, hx, Emu(960120), w, Emu(320040),
                h, "Calibri", Pt(9), True, C_WHITE, PP_ALIGN.CENTER)
    hx += w

# Data rows
comp_data = [
    ("N高・S高", "オンライン/通学", "25〜120万", "先進的IT教育", "自己管理力がある人"),
    ("ゼロ高", "体験型", "約50万〜", "起業・実践中心", "行動力がある人"),
    ("ヒューマンキャンパス", "専門特化", "約40〜80万", "40以上の専門分野", "好きを極めたい人"),
    ("クラーク記念国際", "通学型", "約60〜100万", "全国キャンパス", "仲間と学びたい人"),
    ("ルネサンス", "オンライン", "約25〜40万", "スマホ学習", "自分のペースの人"),
    ("NIC International", "グローバル", "約150万〜", "海外大学進学", "海外を目指す人"),
    ("公立通信制", "通信型", "約3〜5万", "低コスト", "費用を抑えたい人"),
]

for ri, (name, typ, cost, feat, suited) in enumerate(comp_data):
    bg = C_BG_LIGHT if ri % 2 == 0 else C_WHITE
    ry = Emu(1280160) + Emu(320040) * ri
    vals = [name, typ, cost, feat, suited]
    rx = MARGIN
    for ci, (val, w) in enumerate(zip(vals, comp_widths)):
        add_rect(slide, rx, ry, w, Emu(320040), bg)
        add_textbox(slide, rx, ry, w, Emu(320040),
                    val, "Calibri", Pt(8), ci == 0, C_DARK, PP_ALIGN.CENTER)
        rx += w

add_textbox(slide, MARGIN, Emu(4663440), CONTENT_W, Emu(274320),
            "※ 学費は就学支援金適用前の目安です。詳細は各学校の公式サイトをご確認ください",
            "Calibri", Pt(9), False, C_GRAY, PP_ALIGN.CENTER)

add_notes(slide,
    "ここで、主要な通信制高校を一覧で比較してみましょう。\n"
    "N高・S高はオンラインと通学の両方に対応し、学費は年間25万から120万円程度。\n"
    "ゼロ高は体験型の学びが中心で、行動力のある人に向いています。\n"
    "ヒューマンキャンパスは40以上の専門分野を学べるのが最大の特徴。\n"
    "クラーク記念国際は全国にキャンパスがあり、通学して仲間と学びたい方に最適です。\n"
    "ルネサンスはスマホで学べるスタイルで、自分のペースを大切にしたい方向け。\n"
    "NIC Internationalは海外大学進学に特化した学校です。\n"
    "公立通信制は年間3万から5万円と圧倒的に費用を抑えられます。\n"
    "この表を参考に、自分に合いそうな学校を絞り込んでみてください。\n"
    "学費は就学支援金適用前の目安なので、実際にはもっと安くなる可能性があります。"
)


# ════════════════════════════════════════════════════
# SLIDE 12: 学費と支援制度
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "「お金の不安」を超える！学費と支援制度", "知らないと損する支援制度を完全攻略")
add_footer(slide, 11)

# 学費比較テーブル風
add_rounded_rect(slide, MARGIN, Emu(1005840), CONTENT_W, Emu(1554480), C_WHITE)
add_textbox(slide, Emu(548640), Emu(1005840), CONTENT_W, Emu(320040),
            "通信制高校の学費目安（年間）", "Calibri", Pt(14), True, C_GOLD, PP_ALIGN.LEFT)

# Header row
headers = ["学校タイプ", "授業料", "その他費用", "年間合計目安"]
col_w = Emu(2057400)
for i, h in enumerate(headers):
    x = MARGIN + col_w * i
    add_rect(slide, x, Emu(1326564), col_w, Emu(274320), C_GOLD)
    add_textbox(slide, x, Emu(1326564), col_w, Emu(274320),
                h, "Calibri", Pt(10), True, C_WHITE, PP_ALIGN.CENTER)

# Data rows
rows = [
    ("公立通信制", "約1〜3万円", "教材費等", "約3〜5万円"),
    ("私立（ネット型）", "約25〜40万円", "施設費等", "約30〜50万円"),
    ("私立（通学型）", "約50〜80万円", "施設+活動費", "約70〜120万円"),
]
for ri, row in enumerate(rows):
    bg = C_BG_LIGHT if ri % 2 == 0 else C_WHITE
    y = Emu(1600884) + Emu(274320) * ri
    for ci, val in enumerate(row):
        x = MARGIN + col_w * ci
        add_rect(slide, x, y, col_w, Emu(274320), bg)
        add_textbox(slide, x, y, col_w, Emu(274320),
                    val, "Calibri", Pt(10), False, C_DARK, PP_ALIGN.CENTER)

# 支援制度
add_rounded_rect(slide, MARGIN, Emu(2743200), CONTENT_W, Emu(1737360), C_WHITE)
add_textbox(slide, Emu(548640), Emu(2743200), CONTENT_W, Emu(320040),
            "活用すべき支援制度", "Calibri", Pt(14), True, C_ACCENT_GREEN, PP_ALIGN.LEFT)

supports = [
    ("高等学校等就学支援金", "年収約910万円未満の世帯が対象。私立は最大39.6万円/年"),
    ("都道府県独自の助成金", "お住まいの自治体で追加支援あり。要確認！"),
    ("学校独自の奨学金・特待生制度", "成績優秀者や特技のある生徒向けの減免制度"),
    ("学び直し支援金", "再入学者向けの国の支援制度"),
]
y = Emu(3063240)
for name, desc in supports:
    add_textbox(slide, Emu(640080), y, Emu(2651760), Emu(274320),
                name, "Calibri", Pt(10), True, C_DARK, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(3383280), y, Emu(5212080), Emu(274320),
                desc, "Calibri", Pt(9), False, C_GRAY, PP_ALIGN.LEFT)
    y += Emu(320040)

add_notes(slide,
    "次に、多くの方が気になる学費の話です。\n"
    "通信制高校の学費は、学校のタイプによって大きく異なります。\n"
    "公立の通信制なら年間3万から5万円程度と非常にリーズナブルです。\n"
    "私立のネット型は年間30万から50万円、通学型は70万から120万円が目安です。\n"
    "ただし、学費を大幅に軽減できる支援制度がたくさんあります。\n"
    "まず、高等学校等就学支援金。年収約910万円未満の世帯が対象で、\n"
    "私立の場合は最大で年間39万6千円が支給されます。\n"
    "さらに、都道府県独自の助成金もあります。お住まいの自治体のHPで必ず確認してください。\n"
    "学校独自の奨学金や特待生制度、学び直し支援金などもあります。\n"
    "知らないと損する制度ばかりなので、しっかりチェックしましょう。"
)


# ════════════════════════════════════════════════════
# SLIDE 13: 支援制度申請ステップガイド
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "就学支援金 申請ステップガイド", "申請漏れゼロ！確実に支援を受けるために")
add_footer(slide, 12)

# 年収別シミュレーション
add_rounded_rect(slide, MARGIN, Emu(960120), Emu(4023360), Emu(1920240), C_WHITE)
add_textbox(slide, Emu(548640), Emu(960120), Emu(3840480), Emu(320040),
            "年収別 支援金シミュレーション", "Calibri", Pt(13), True, C_ACCENT_BLUE, PP_ALIGN.LEFT)

sim_data = [
    ("年収270万円未満", "最大39.6万円/年 支給", "私立でも実質負担が大幅減"),
    ("年収270〜590万円", "最大39.6万円/年 支給", "多くの家庭がこの区分に該当"),
    ("年収590〜910万円", "最大11.88万円/年 支給", "一部支給で負担を軽減"),
    ("年収910万円以上", "支給対象外", "学校独自の奨学金を活用"),
]
sy = Emu(1280160)
for income, amount, note in sim_data:
    add_textbox(slide, Emu(548640), sy, Emu(1554480), Emu(274320),
                income, "Calibri", Pt(9), True, C_DARK, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(2103120), sy, Emu(1280160), Emu(274320),
                amount, "Calibri", Pt(9), True, C_ACCENT_GREEN, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(3383280), sy, Emu(1188720), Emu(274320),
                note, "Calibri", Pt(8), False, C_GRAY, PP_ALIGN.LEFT)
    sy += Emu(320040)

# 申請ステップ
add_rounded_rect(slide, Emu(4663440), Emu(960120), Emu(4023360), Emu(1920240), C_WHITE)
add_textbox(slide, Emu(4754880), Emu(960120), Emu(3840480), Emu(320040),
            "申請の3ステップ", "Calibri", Pt(13), True, C_ACCENT_ORANGE, PP_ALIGN.LEFT)

app_steps = [
    ("STEP 1", "入学時に学校から申請書類を受け取る"),
    ("STEP 2", "マイナンバーカード等で\n保護者の所得情報を提出"),
    ("STEP 3", "審査後、学校の授業料から\n支援金が差し引かれる"),
]
asy = Emu(1280160)
for step, desc in app_steps:
    add_rounded_rect(slide, Emu(4754880), asy, Emu(822960), Emu(274320), C_ACCENT_ORANGE)
    add_textbox(slide, Emu(4754880), asy, Emu(822960), Emu(274320),
                step, "Calibri", Pt(9), True, C_WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, Emu(5669280), asy, Emu(2926080), Emu(457200),
                desc, "Calibri", Pt(9), False, C_DARK, PP_ALIGN.LEFT)
    asy += Emu(502920)

# 注意事項
add_rounded_rect(slide, MARGIN, Emu(3063240), CONTENT_W, Emu(1371600), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Emu(548640), Emu(3063240), CONTENT_W, Emu(320040),
            "申請時の注意点", "Calibri", Pt(13), True, C_ACCENT_ORANGE, PP_ALIGN.LEFT)

warnings = [
    "申請期限を必ず確認（入学後すぐに手続きが必要な場合も）",
    "都道府県の上乗せ助成金は別途申請が必要",
    "転入・編入の場合は前の学校での在籍期間に注意",
    "7月頃に収入状況届出（継続届）の提出が必要",
]
wy = Emu(3383280)
for w in warnings:
    add_textbox(slide, Emu(640080), wy, Emu(8046720), Emu(274320),
                f"  {w}", "Calibri", Pt(10), False, C_DARK, PP_ALIGN.LEFT)
    wy += Emu(274320)

add_notes(slide,
    "就学支援金の具体的な申請方法をお伝えします。\n"
    "まず年収別のシミュレーションです。\n"
    "年収270万円未満の世帯は、最大で年間39万6千円が支給されます。\n"
    "年収590万円までの世帯も同額の支給を受けられます。多くの家庭がこの区分に該当します。\n"
    "年収910万円までの世帯は一部支給で、最大11万8800円です。\n"
    "申請は3ステップで簡単です。\n"
    "ステップ1、入学時に学校から申請書類を受け取ります。\n"
    "ステップ2、マイナンバーカード等で保護者の所得情報を提出します。\n"
    "ステップ3、審査後、授業料から支援金が差し引かれます。\n"
    "注意点として、申請期限の確認、都道府県の上乗せ助成金の別途申請、\n"
    "7月頃の継続届の提出を忘れないようにしましょう。"
)


# ════════════════════════════════════════════════════
# SLIDE 14: 学校選びの実践術
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "後悔しない！学校選び実践術", "ネット情報を鵜呑みにしない賢い選び方")
add_footer(slide, 13)

# 3 step cards
steps = [
    ("STEP 1", "情報収集", "学校公式サイトを読み解く\nパンフレットを取り寄せる\nSNS・口コミは参考程度に", C_ACCENT_BLUE),
    ("STEP 2", "体験・比較", "オープンキャンパスに参加\n個別相談で本音を聞く\n3校以上を比較検討する", C_ACCENT_GREEN),
    ("STEP 3", "家族で対話", "親子で本音を共有する\n費用・通学・将来を話し合う\n「自分で決めた」実感を大切に", C_ACCENT_ORANGE),
]

card_w = Emu(2651760)
gap = Emu(182880)
for i, (step, title, desc, color) in enumerate(steps):
    x = MARGIN + (card_w + gap) * i
    y = Emu(1005840)
    add_rounded_rect(slide, x, y, card_w, Emu(3200400), C_WHITE)
    # step label
    add_rounded_rect(slide, x + Emu(91440), y + Emu(91440), Emu(914400), Emu(320040), color)
    add_textbox(slide, x + Emu(91440), y + Emu(91440), Emu(914400), Emu(320040),
                step, "Calibri", Pt(11), True, C_WHITE, PP_ALIGN.CENTER)
    # title
    add_textbox(slide, x + Emu(91440), y + Emu(502920), card_w - Emu(182880), Emu(411480),
                title, "Calibri", Pt(18), True, C_DARK, PP_ALIGN.CENTER)
    # divider
    add_rect(slide, x + Emu(457200), y + Emu(914400), card_w - Emu(914400), Emu(9144), C_BG_LIGHT)
    # desc
    lines = desc.split("\n")
    dy = y + Emu(1051560)
    for line in lines:
        add_textbox(slide, x + Emu(137160), dy, card_w - Emu(274320), Emu(365760),
                    f"  {line}", "Calibri", Pt(11), False, C_DARK, PP_ALIGN.LEFT)
        dy += Emu(365760)

# bottom warning
add_rounded_rect(slide, MARGIN, Emu(4297680), CONTENT_W, Emu(365760), RGBColor(0xFF, 0xF3, 0xE0))
add_textbox(slide, Emu(548640), Emu(4297680), Emu(8046720), Emu(365760),
            "注意：「〇〇ランキング」だけで学校を決めないこと。自分の目で確かめよう",
            "Calibri", Pt(11), True, C_ACCENT_ORANGE, PP_ALIGN.LEFT)

add_notes(slide,
    "学校選びで後悔しないための実践的なステップをお伝えします。\n"
    "ステップ1は情報収集。学校の公式サイトを丁寧に読み解き、パンフレットを取り寄せましょう。\n"
    "SNSや口コミサイトの情報は参考程度にとどめてください。\n"
    "ステップ2は体験と比較。必ずオープンキャンパスに参加して、雰囲気を肌で感じてください。\n"
    "個別相談で先生に本音の質問をすることも大切です。\n"
    "最低3校以上は比較検討することをおすすめします。\n"
    "ステップ3は家族での対話。費用のこと、通学のこと、将来のこと、\n"
    "親子で本音を共有し、一緒に考えることが大切です。\n"
    "最終的に「自分で決めた」という実感を持てることが、入学後のモチベーションにつながります。\n"
    "ネットの「ランキング」だけで決めず、必ず自分の目で確かめましょう。"
)


# ════════════════════════════════════════════════════
# SLIDE 15: 親子コミュニケーション術
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "親子で「未来の進路」を語ろう", "通信制高校を家族で考えるコミュニケーション術")
add_footer(slide, 14)

# Left: 保護者向け
add_rounded_rect(slide, MARGIN, Emu(1005840), Emu(4023360), Emu(3200400), C_WHITE)
add_rect(slide, MARGIN, Emu(1005840), Emu(4023360), Emu(54864), C_ACCENT_BLUE)
add_textbox(slide, Emu(548640), Emu(1097280), Emu(3840480), Emu(320040),
            "保護者の方へ", "Calibri", Pt(14), True, C_ACCENT_BLUE, PP_ALIGN.LEFT)

parent_tips = [
    ("子どもの話を最後まで聞く", "否定せず、まずは気持ちを受け止める"),
    ("「なぜ通信制？」ではなく\n「何を学びたい？」と聞く", "動機を理解することが大切"),
    ("一緒に学校を調べる", "子どもと同じ情報を共有する"),
    ("最終決定は本人に委ねる", "自己決定感が入学後の意欲に直結"),
]
py = Emu(1463040)
for tip, note in parent_tips:
    add_textbox(slide, Emu(548640), py, Emu(2468880), Emu(457200),
                f"  {tip}", "Calibri", Pt(10), True, C_DARK, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(3017520), py, Emu(1463040), Emu(457200),
                note, "Calibri", Pt(8), False, C_GRAY, PP_ALIGN.LEFT)
    py += Emu(548640)

# Right: 生徒向け
add_rounded_rect(slide, Emu(4663440), Emu(1005840), Emu(4023360), Emu(3200400), C_WHITE)
add_rect(slide, Emu(4663440), Emu(1005840), Emu(4023360), Emu(54864), C_ACCENT_GREEN)
add_textbox(slide, Emu(4754880), Emu(1097280), Emu(3840480), Emu(320040),
            "生徒の皆さんへ", "Calibri", Pt(14), True, C_ACCENT_GREEN, PP_ALIGN.LEFT)

student_tips = [
    ("自分の気持ちを言葉にする", "「なんとなく」ではなく具体的に"),
    ("調べた情報を親に共有する", "一方的でなく、一緒に考える姿勢"),
    ("不安なことも正直に伝える", "隠すより話した方が解決が早い"),
    ("「こうしたい」を明確にする", "目標があると親も安心する"),
]
sy = Emu(1463040)
for tip, note in student_tips:
    add_textbox(slide, Emu(4754880), sy, Emu(2468880), Emu(457200),
                f"  {tip}", "Calibri", Pt(10), True, C_DARK, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(7223760), sy, Emu(1463040), Emu(457200),
                note, "Calibri", Pt(8), False, C_GRAY, PP_ALIGN.LEFT)
    sy += Emu(548640)

# bottom
add_rounded_rect(slide, MARGIN, Emu(4297680), CONTENT_W, Emu(365760), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(4297680), Emu(8046720), Emu(365760),
            "家族で納得して決めた進路は、入学後の大きな支えになります",
            "Calibri", Pt(12), True, C_GOLD, PP_ALIGN.CENTER)

add_notes(slide,
    "進路を決める上で、親子のコミュニケーションはとても大切です。\n"
    "保護者の方へのアドバイスです。\n"
    "まず、お子さんの話を最後まで聞いてあげてください。否定せず、気持ちを受け止めましょう。\n"
    "「なぜ通信制なの？」ではなく「何を学びたいの？」と聞いてみてください。\n"
    "一緒に学校を調べることで、同じ情報を共有できます。\n"
    "そして最終決定はお子さん本人に委ねてください。自分で決めたという実感が入学後の意欲に直結します。\n"
    "生徒の皆さんへ。\n"
    "自分の気持ちを言葉にして、具体的に伝えましょう。\n"
    "調べた情報を親に共有して、一緒に考える姿勢を見せてください。\n"
    "不安なことも正直に伝えた方が、解決が早くなります。\n"
    "家族で納得して決めた進路は、入学後の大きな支えになりますよ。"
)


# ════════════════════════════════════════════════════
# SLIDE 16: 先輩たちのストーリー
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "先輩たちの選択とその後の人生", "通信制高校から広がる多様な未来")
add_footer(slide, 15)

stories = [
    ("不登校→再起", "通信制で心身を立て直し\n新たな興味を発見して活躍", "心のケアが\n充実", C_ACCENT_GREEN),
    ("好き→プロへ", "eスポーツ・動画編集など\n在学中にスキルを磨きプロに", "実践力が\n武器に", C_ACCENT_BLUE),
    ("通信制→難関大", "自分のペースで学習し\n難関国立大学に現役合格", "戦略的\n進学", C_ACCENT_PURPLE),
    ("起業・フリーランス", "高校在学中にビジネスを\n立ち上げて成功", "行動力が\n未来を拓く", C_ACCENT_ORANGE),
]

card_w = Emu(1966440)
gap = Emu(137160)
for i, (title, desc, tag, color) in enumerate(stories):
    x = MARGIN + (card_w + gap) * i
    y = Emu(1005840)
    add_rounded_rect(slide, x, y, card_w, Emu(2743200), C_WHITE)
    add_rect(slide, x, y, card_w, Emu(54864), color)
    # circle icon
    cx = x + card_w // 2 - Emu(274320)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, y + Emu(137160), Emu(548640), Emu(548640))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide, cx, y + Emu(182880), Emu(548640), Emu(457200),
                tag, "Calibri", Pt(8), True, C_WHITE, PP_ALIGN.CENTER)
    # title
    add_textbox(slide, x + Emu(45720), y + Emu(777240), card_w - Emu(91440), Emu(365760),
                title, "Calibri", Pt(13), True, C_DARK, PP_ALIGN.CENTER)
    # desc
    add_textbox(slide, x + Emu(45720), y + Emu(1188720), card_w - Emu(91440), Emu(914400),
                desc, "Calibri", Pt(10), False, C_GRAY, PP_ALIGN.CENTER)

# bottom message
add_rounded_rect(slide, MARGIN, Emu(3886200), CONTENT_W, Emu(548640), C_BG_LIGHT)
add_textbox(slide, Emu(548640), Emu(3931920), Emu(8046720), Emu(411480),
            "「やり直し」や「回り道」もポジティブな選択肢。あなたの未来は、あなたが創る。",
            "Calibri", Pt(13), True, C_GOLD, PP_ALIGN.CENTER)

add_notes(slide,
    "ここで、実際に通信制高校を選んだ先輩たちのストーリーをご紹介します。\n"
    "まず、不登校から再起した先輩。通信制高校で心身を立て直し、\n"
    "個別サポートの中で新たな興味を発見して、今は社会で活躍しています。\n"
    "次に、好きなことをプロレベルまで極めた先輩。\n"
    "eスポーツや動画編集のスキルを在学中に磨き、フリーランスとして活躍しています。\n"
    "通信制から難関大学に合格した先輩もいます。\n"
    "自分のペースで学習できる環境を活かして、戦略的に受験勉強に取り組みました。\n"
    "起業やフリーランスの道を選んだ先輩は、高校在学中にビジネスを立ち上げて成功しています。\n"
    "「やり直し」や「回り道」もポジティブな選択肢です。\n"
    "あなたの未来は、あなた自身が創るものです。"
)


# ════════════════════════════════════════════════════
# SLIDE 13: 今すぐ始める行動プラン
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "今すぐ始める！行動プラン", "最初の一歩を踏み出そう")
add_footer(slide, 16)

actions = [
    ("今日", "気になる学校を3つリストアップ\n公式サイトとパンフレットをチェック", C_ACCENT_GREEN),
    ("今週中", "オープンキャンパスや説明会に申し込む\n支援制度を調べる（自治体HPも確認）", C_ACCENT_BLUE),
    ("今月中", "実際にオープンキャンパスに参加\n個別相談で直接質問する", C_ACCENT_PURPLE),
    ("来月まで", "親子で話し合い、比較表を作成\n志望校を絞り込む", C_ACCENT_ORANGE),
]

card_w = Emu(1966440)
gap = Emu(137160)
for i, (when, what, color) in enumerate(actions):
    x = MARGIN + (card_w + gap) * i
    y = Emu(1005840)
    add_rounded_rect(slide, x, y, card_w, Emu(2743200), C_WHITE)
    # when label
    add_rounded_rect(slide, x + Emu(182880), y + Emu(137160), card_w - Emu(365760), Emu(457200), color)
    add_textbox(slide, x + Emu(182880), y + Emu(182880), card_w - Emu(365760), Emu(365760),
                when, "Calibri", Pt(18), True, C_WHITE, PP_ALIGN.CENTER)
    # arrow
    add_textbox(slide, x, y + Emu(640080), card_w, Emu(274320),
                "▼", "Calibri", Pt(14), False, color, PP_ALIGN.CENTER)
    # what
    add_textbox(slide, x + Emu(91440), y + Emu(960120), card_w - Emu(182880), Emu(1554480),
                what, "Calibri", Pt(11), False, C_DARK, PP_ALIGN.LEFT)

# bottom
add_rounded_rect(slide, MARGIN, Emu(3886200), CONTENT_W, Emu(548640), C_GOLD)
add_textbox(slide, Emu(548640), Emu(3931920), Emu(8046720), Emu(411480),
            "大切なのは「完璧な選択」ではなく「自分で選ぶ」こと。今日がその第一歩！",
            "Calibri", Pt(13), True, C_WHITE, PP_ALIGN.CENTER)

add_notes(slide,
    "最後に、今すぐ始められる行動プランをお伝えします。\n"
    "まず今日やること。気になる学校を3つリストアップして、公式サイトとパンフレットをチェックしましょう。\n"
    "今週中には、オープンキャンパスや説明会に申し込みましょう。\n"
    "支援制度も調べてみてください。お住まいの自治体のホームページも忘れずにチェック。\n"
    "今月中には実際にオープンキャンパスに足を運んで、雰囲気を肌で感じてください。\n"
    "個別相談で先生に直接質問することも大切です。\n"
    "来月までには親子で話し合って、比較表を作り、志望校を絞り込んでいきましょう。\n"
    "大切なのは「完璧な選択」をすることではなく、「自分で選ぶ」ということです。\n"
    "今日がその第一歩。ぜひ行動に移してみてください。"
)


# ════════════════════════════════════════════════════
# SLIDE 18: 入学準備チェックリスト
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "入学準備チェックリスト", "迷わず進むためのロードマップ")
add_footer(slide, 17)

# Left: 入学前
add_rounded_rect(slide, MARGIN, Emu(960120), Emu(2651760), Emu(3657600), C_WHITE)
add_rect(slide, MARGIN, Emu(960120), Emu(2651760), Emu(54864), C_ACCENT_GREEN)
add_textbox(slide, Emu(548640), Emu(1005840), Emu(2468880), Emu(320040),
            "入学前", "Calibri", Pt(14), True, C_ACCENT_GREEN, PP_ALIGN.LEFT)

before_items = [
    "願書・必要書類の準備",
    "面接対策（志望動機を整理）",
    "転入・編入の場合は単位確認",
    "就学支援金の申請準備",
    "学習環境の整備（PC・ネット）",
    "スクーリング日程の確認",
]
by = Emu(1371600)
for item in before_items:
    add_textbox(slide, Emu(548640), by, Emu(2468880), Emu(274320),
                f"  {item}", "Calibri", Pt(10), False, C_DARK, PP_ALIGN.LEFT)
    by += Emu(320040)

# Middle: 入学後（最初の1ヶ月）
add_rounded_rect(slide, Emu(3291840), Emu(960120), Emu(2651760), Emu(3657600), C_WHITE)
add_rect(slide, Emu(3291840), Emu(960120), Emu(2651760), Emu(54864), C_ACCENT_BLUE)
add_textbox(slide, Emu(3383280), Emu(1005840), Emu(2468880), Emu(320040),
            "入学後 最初の1ヶ月", "Calibri", Pt(14), True, C_ACCENT_BLUE, PP_ALIGN.LEFT)

after_items = [
    "学習計画を立てる",
    "レポート提出の仕組みを理解",
    "担任・メンターと面談",
    "オンライン学習ツールに慣れる",
    "同級生との交流機会に参加",
    "困ったら早めに相談する",
]
ay = Emu(1371600)
for item in after_items:
    add_textbox(slide, Emu(3383280), ay, Emu(2468880), Emu(274320),
                f"  {item}", "Calibri", Pt(10), False, C_DARK, PP_ALIGN.LEFT)
    ay += Emu(320040)

# Right: 面接対策
add_rounded_rect(slide, Emu(6126480), Emu(960120), Emu(2560320), Emu(3657600), C_WHITE)
add_rect(slide, Emu(6126480), Emu(960120), Emu(2560320), Emu(54864), C_ACCENT_PURPLE)
add_textbox(slide, Emu(6217920), Emu(1005840), Emu(2377440), Emu(320040),
            "面接でよくある質問", "Calibri", Pt(14), True, C_ACCENT_PURPLE, PP_ALIGN.LEFT)

interview_qs = [
    "Q. なぜ通信制高校を\n    選びましたか？",
    "Q. 入学後にやりたい\n    ことは何ですか？",
    "Q. 将来の目標を\n    教えてください",
    "Q. 自分の長所は\n    何だと思いますか？",
]
iy = Emu(1371600)
for q in interview_qs:
    add_rounded_rect(slide, Emu(6217920), iy, Emu(2377440), Emu(502920), C_BG_LIGHT)
    add_textbox(slide, Emu(6309360), iy + Emu(45720), Emu(2194560), Emu(457200),
                q, "Calibri", Pt(9), False, C_DARK, PP_ALIGN.LEFT)
    iy += Emu(548640)

add_notes(slide,
    "入学準備のチェックリストをまとめました。\n"
    "入学前にやるべきこと。願書や必要書類の準備、面接対策として志望動機を整理しておきましょう。\n"
    "転入・編入の場合は、前の学校での単位がどれだけ引き継げるか確認してください。\n"
    "就学支援金の申請準備も忘れずに。学習に必要なパソコンやネット環境も整えましょう。\n"
    "入学後の最初の1ヶ月は、学習計画を立てることが最も大切です。\n"
    "レポート提出の仕組みを理解して、担任やメンターとの面談を早めに行いましょう。\n"
    "困ったことがあれば、早めに相談することがポイントです。\n"
    "面接では「なぜ通信制を選んだか」「入学後にやりたいこと」「将来の目標」が\n"
    "よく聞かれます。事前に自分の言葉で話せるように準備しておきましょう。"
)


# ════════════════════════════════════════════════════
# SLIDE 19: よくある質問（FAQ）
# ════════════════════════════════════════════════════
slide = new_slide()
add_title_bar(slide, "よくある質問（FAQ）", "視聴者の皆さんからの疑問に回答")
add_footer(slide, 18)

faqs = [
    ("Q. 通信制高校の卒業資格は全日制と同じですか？",
     "A. はい、同じです。「高校卒業」の資格は全日制と全く同じ扱いです。\n"
     "    履歴書にも「○○高等学校 卒業」と記載できます。"),
    ("Q. 大学受験に不利になりませんか？",
     "A. 不利にはなりません。推薦入試やAO入試では、\n"
     "    通信制ならではの経験がアピールポイントになることもあります。"),
    ("Q. 友達はできますか？",
     "A. できます！スクーリング、部活動、オンラインコミュニティなど、\n"
     "    交流の機会は多くの学校で用意されています。"),
    ("Q. 途中で全日制に戻ることはできますか？",
     "A. 制度上は可能ですが、カリキュラムの違いがあるため、\n"
     "    事前に転入先の学校に相談することをおすすめします。"),
]

fy = Emu(960120)
for i, (q, a) in enumerate(faqs):
    colors = [C_ACCENT_BLUE, C_ACCENT_GREEN, C_ACCENT_ORANGE, C_ACCENT_PURPLE]
    add_rounded_rect(slide, MARGIN, fy, CONTENT_W, Emu(822960), C_WHITE)
    add_rect(slide, MARGIN, fy, Emu(36576), Emu(822960), colors[i])
    add_textbox(slide, Emu(594360), fy + Emu(45720), Emu(8046720), Emu(274320),
                q, "Calibri", Pt(11), True, C_DARK, PP_ALIGN.LEFT)
    add_textbox(slide, Emu(594360), fy + Emu(320040), Emu(8046720), Emu(457200),
                a, "Calibri", Pt(10), False, C_GRAY, PP_ALIGN.LEFT)
    fy += Emu(868680)

add_notes(slide,
    "よくある質問にお答えします。\n"
    "まず、通信制高校の卒業資格は全日制と同じですか？という質問。\n"
    "はい、全く同じです。高校卒業の資格に違いはありません。\n"
    "次に、大学受験に不利になりませんか？\n"
    "不利にはなりません。むしろ推薦入試やAO入試では、\n"
    "通信制ならではのユニークな経験がアピールポイントになることもあります。\n"
    "友達はできますか？という質問も多いですが、できます。\n"
    "スクーリングや部活動、オンラインコミュニティなど、交流の機会はたくさんあります。\n"
    "途中で全日制に戻ることは制度上可能ですが、カリキュラムの違いがあるので、\n"
    "事前に転入先の学校に相談することをおすすめします。"
)


# ════════════════════════════════════════════════════
# SLIDE 20: まとめ＆エンディング
# ════════════════════════════════════════════════════
slide = new_slide()
# top & bottom gold lines
add_rect(slide, 0, 0, Emu(SLIDE_WIDTH), Emu(54864), C_GOLD_LIGHT)
add_rect(slide, 0, Emu(5088636), Emu(SLIDE_WIDTH), Emu(54864), C_GOLD_LIGHT)

add_textbox(slide, MARGIN, Emu(457200), CONTENT_W, Emu(502920),
            "まとめ", "Calibri", Pt(28), True, C_BROWN, PP_ALIGN.CENTER)

summary = [
    "N高・ゼロ高だけが通信制高校ではない",
    "あなたのタイプに合った学校が必ずある",
    "学費の不安は支援制度で解消できる",
    "自分の目で見て、自分で決めることが大切",
    "「回り道」も立派な選択肢",
]
y = Emu(1097280)
for i, s in enumerate(summary):
    num = str(i + 1)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(2286000), y, Emu(365760), Emu(365760))
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_GOLD
    circle.line.fill.background()
    add_textbox(slide, Emu(2286000), y + Emu(45720), Emu(365760), Emu(274320),
                num, "Calibri", Pt(14), True, C_WHITE, PP_ALIGN.CENTER)
    add_textbox(slide, Emu(2834640), y, Emu(5486400), Emu(365760),
                s, "Calibri", Pt(15), False, C_DARK, PP_ALIGN.LEFT)
    y += Emu(457200)

# decorative line
add_rect(slide, Emu(2286000), y + Emu(137160), Emu(4572000), Emu(9144), C_GOLD_LIGHT)

# CTA
add_textbox(slide, MARGIN, y + Emu(365760), CONTENT_W, Emu(411480),
            "チャンネル登録・高評価お願いします！", "Calibri", Pt(16), True, C_GOLD, PP_ALIGN.CENTER)
add_textbox(slide, MARGIN, y + Emu(777240), CONTENT_W, Emu(320040),
            "コメント欄であなたの気になる通信制高校を教えてください", "Calibri", Pt(12), False, C_GRAY, PP_ALIGN.CENTER)

# Harmonic Insight branding
add_textbox(slide, MARGIN, Emu(4206240), CONTENT_W, Emu(274320),
            "H A R M O N I C   i n s i g h t", "Calibri", Pt(12), False, C_BROWN, PP_ALIGN.CENTER)

add_notes(slide,
    "今日のまとめです。\n"
    "1つ目、N高やゼロ高だけが通信制高校ではありません。もっとたくさんの選択肢があります。\n"
    "2つ目、あなたのタイプに合った学校が必ずあります。自分を知ることが学校選びの第一歩です。\n"
    "3つ目、学費の不安は支援制度で解消できます。知らないと損する制度がたくさんあるので、しっかり調べましょう。\n"
    "4つ目、自分の目で見て、自分で決めることが大切です。ネットの情報だけで決めないでください。\n"
    "5つ目、「回り道」も立派な選択肢です。人と違う道を選ぶ勇気を持ってください。\n"
    "今日の動画が皆さんの学校選びの参考になれば嬉しいです。\n"
    "チャンネル登録と高評価をお願いします。\n"
    "コメント欄で、あなたの気になる通信制高校があれば教えてください。\n"
    "それでは、また次の動画でお会いしましょう。ありがとうございました。"
)


# ── 保存 ──
output_path = "N高ゼロ高を超える_通信制高校選び完全ガイド_YouTube.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
